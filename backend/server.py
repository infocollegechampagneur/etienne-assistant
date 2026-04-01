"""
╔══════════════════════════════════════════════════════════════════════════════╗
║                    ÉTIENNE - Backend API Server                               ║
║                    Assistant IA pour le personnel scolaire québécois          ║
╚══════════════════════════════════════════════════════════════════════════════╝

📋 INDEX DES SECTIONS (Ctrl+F pour naviguer):
─────────────────────────────────────────────
[SECTION 1]  IMPORTS & CONFIGURATION         - Ligne ~30
[SECTION 2]  QUOTA MANAGER                   - Ligne ~55  (GeminiQuotaManager)
[SECTION 3]  GRAPHIQUES MATHÉMATIQUES        - Ligne ~150 (matplotlib)
[SECTION 4]  DATABASE & APP SETUP            - Ligne ~520
[SECTION 5]  MODÈLES PYDANTIC (Requêtes)     - Ligne ~575
[SECTION 6]  FONCTIONS UTILITAIRES           - Ligne ~680 (trust_score, extract_text, etc.)
[SECTION 7]  DÉTECTION IA & CONTENU          - Ligne ~850
[SECTION 8]  RÉPONSE IA (GEMINI)             - Ligne ~1190 (get_ai_response + PROTOCOLE MEQ CORRECTION)
[SECTION 9]  ROUTES BASIQUES                 - Ligne ~2350 (health, chat, quota)
[SECTION 10] GÉNÉRATION DOCUMENTS            - Ligne ~2500 (PDF, DOCX, PPTX, XLSX)
[SECTION 11] UPLOAD & ANALYSE FICHIERS       - Ligne ~2900
[SECTION 12] MODÈLES PYDANTIC (Admin)        - Ligne ~3100
[SECTION 13] ROUTES DÉTECTION IA             - Ligne ~3520
[SECTION 14] ROUTES ADMIN (Logs, Reports)    - Ligne ~3680
[SECTION 15] GÉNÉRATION IMAGES               - Ligne ~3960 (Hugging Face)
[SECTION 16] AUTHENTIFICATION                - Ligne ~4060 (signup, login)
[SECTION 17] GESTION LICENCES                - Ligne ~4220
[SECTION 18] MOTS BLOQUÉS                    - Ligne ~4430
[SECTION 19] ADMIN MULTI-NIVEAUX             - Ligne ~4560
[SECTION 20] CONVERSATIONS CLOUD             - Ligne ~4850
[SECTION 21] RÉINITIALISATION MOT DE PASSE   - Ligne ~5040
[SECTION 22] CHANGEMENT EMAIL                - Ligne ~5300
[SECTION 23] CHANGEMENT MOT DE PASSE         - Ligne ~5480
[SECTION 24] STATISTIQUES ADMIN              - Ligne ~5640

🔑 CREDENTIALS TEST:
- Super Admin: informatique@champagneur.qc.ca / !0910Hi8ki8+

⚙️ VARIABLES ENVIRONNEMENT REQUISES:
- MONGO_URL, DB_NAME, JWT_SECRET_KEY, GOOGLE_API_KEY
- SMTP_HOST, SMTP_PORT, SMTP_USERNAME, SMTP_PASSWORD, SMTP_FROM_EMAIL (pour emails)
- HUGGINGFACE_API_KEY (optionnel, pour images artistiques)
"""

from fastapi import FastAPI, APIRouter, HTTPException, Request, Depends
from fastapi.security import HTTPBearer, HTTPAuthorizationCredentials
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import sys
import logging
from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Optional
import uuid
from datetime import datetime, timezone, timedelta
import asyncio
import re
from collections import defaultdict
from threading import Lock
from fastapi.responses import StreamingResponse, FileResponse
from fastapi import UploadFile, File, Form
from io import BytesIO
from ai_detector_advanced import AdvancedAIDetector
import google.generativeai as genai
from docx import Document
from docx.shared import Inches as DocxInches, RGBColor, Pt as DocxPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from pptx import Presentation
from pptx.util import Inches, Pt
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
import PyPDF2
import docx2txt
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
import secrets
import time

# ══════════════════════════════════════════════════════════════════════════════
# [SECTION 2] SYSTÈME DE GESTION DU QUOTA GEMINI
# ══════════════════════════════════════════════════════════════════════════════
class GeminiQuotaManager:
    """Gestionnaire de quota pour l'API Gemini avec retry automatique"""
    
    def __init__(self):
        self.max_requests_per_minute = 15  # Limite conservatrice (API free: 15-20/min)
        self.requests = []  # Liste des timestamps des requêtes
        self.lock = Lock()
        self.retry_after = None  # Timestamp après lequel on peut réessayer
        self.last_error_message = None
    
    def get_status(self):
        """Obtenir le statut actuel du quota"""
        with self.lock:
            now = time.time()
            # Nettoyer les requêtes de plus d'une minute
            self.requests = [t for t in self.requests if now - t < 60]
            
            used = len(self.requests)
            remaining = max(0, self.max_requests_per_minute - used)
            
            # Calculer le temps jusqu'au reset
            if self.requests:
                oldest_request = min(self.requests)
                reset_in_seconds = max(0, 60 - (now - oldest_request))
            else:
                reset_in_seconds = 0
            
            # Si on est en période de retry forcé
            if self.retry_after and now < self.retry_after:
                wait_seconds = self.retry_after - now
                return {
                    "used": used,
                    "max": self.max_requests_per_minute,
                    "remaining": 0,
                    "reset_in_seconds": int(wait_seconds),
                    "can_request": False,
                    "status": "quota_exceeded",
                    "message": f"Quota dépassé. Réessai dans {int(wait_seconds)}s"
                }
            
            return {
                "used": used,
                "max": self.max_requests_per_minute,
                "remaining": remaining,
                "reset_in_seconds": int(reset_in_seconds) if used > 0 else 0,
                "can_request": remaining > 0,
                "status": "ok" if remaining > 0 else "limit_reached",
                "message": f"{remaining} requêtes disponibles" if remaining > 0 else "Limite atteinte, patientez..."
            }
    
    def record_request(self):
        """Enregistrer une requête"""
        with self.lock:
            self.requests.append(time.time())
    
    def record_quota_error(self, retry_delay_seconds=60):
        """Enregistrer une erreur de quota (429)"""
        with self.lock:
            self.retry_after = time.time() + retry_delay_seconds
            self.last_error_message = f"API quota exceeded. Retry in {retry_delay_seconds}s"
    
    def can_make_request(self):
        """Vérifier si on peut faire une requête"""
        status = self.get_status()
        return status["can_request"]
    
    def wait_time_seconds(self):
        """Obtenir le temps d'attente avant la prochaine requête possible"""
        status = self.get_status()
        if status["can_request"]:
            return 0
        return status["reset_in_seconds"]

# Instance globale du gestionnaire de quota
gemini_quota = GeminiQuotaManager()

import pdfplumber
from PIL import Image
import pandas as pd
import tempfile
import json
import jwt

# Clé secrète JWT (doit correspondre à celle utilisée pour créer les tokens)
JWT_SECRET = os.environ.get("JWT_SECRET_KEY", "etienne_security_jwt_secret_key_2025")

# Import des fonctionnalites de securite avancees
from security_advanced import (
    get_current_admin, check_rate_limit, get_client_ip,
    send_critical_alert, send_daily_report,
    generate_csv_report, generate_pdf_report,
    hash_password, verify_password, create_access_token
)

# === GÉNÉRATION DE GRAPHIQUES MATHÉMATIQUES (CODE INLINE) ===
# Solution de contournement: code inline au lieu du module utils
try:
    import matplotlib
    matplotlib.use('Agg')  # Backend sans GUI pour serveur
    import matplotlib.pyplot as plt
    import numpy as np
    import base64
    from io import BytesIO
    import re
    
    GRAPH_GENERATOR_AVAILABLE = True
    logging.info("✅ Graphiques mathématiques disponibles (code inline)")
    
    def detect_graph_request_inline(text: str) -> bool:
        """Détecte si l'utilisateur demande un graphique mathématique (FR + EN)"""
        # Mots-clés SPÉCIFIQUES aux graphiques - pas de termes génériques
        graph_keywords_fr = [
            'graphique', 'courbe', 'plot'
        ]
        
        # Mots-clés anglais pour graphiques  
        graph_keywords_en = [
            'plot', 'graph', 'chart'
        ]
        
        # Termes mathématiques SPÉCIFIQUES (pas de termes scolaires génériques)
        math_terms = [
            'f(x)', 'fonction mathématique', 'mathematical function',
            'x²', 'x³', 'x**', 'sin(', 'cos(', 'tan(', 'exp(', 'log(',
            'polynôme', 'polynomial', 'parabole', 'parabola',
            'trigonométrique', 'trigonometric'
        ]
        
        text_lower = text.lower()
        
        # Si le message contient "corrige" ou "correction" ou "texte de l'élève",
        # ce n'est PAS une demande de graphique
        correction_indicators = [
            'corrige', 'correction', "texte de l'élève", 'texte de l\'élève',
            'protocole meq', 'barème', 'critère', 'pondération',
            'nombre de mots', 'erreurs c4', 'erreurs c5'
        ]
        if any(ind in text_lower for ind in correction_indicators):
            return False
        
        # Phrases spécifiques (FR + EN)
        specific_phrases = [
            'dessine un graphique', 'dessine moi un graphique', 'dessine le graphique',
            'trace un graphique', 'trace moi un graphique', 'trace le graphique',
            'trace la courbe', 'trace moi la courbe',
            'graphique pour', 'graphique de', 'courbe de', 'représentation graphique',
            'draw a graph', 'draw the graph', 'plot a graph', 'plot the graph',
            'show me a graph', 'show the graph', 'graph of', 'chart of',
            'plot f(x)', 'graph f(x)', 'draw f(x)'
        ]
        
        # Vérifier phrases spécifiques en premier
        if any(phrase in text_lower for phrase in specific_phrases):
            return True
        
        # Vérifier combinaison mot-clé + terme mathématique
        has_keyword = any(kw in text_lower for kw in graph_keywords_fr + graph_keywords_en)
        has_math = any(term in text_lower for term in math_terms)
        
        return has_keyword and has_math
    
    def extract_math_expression_inline(text: str) -> str:
        """Extrait l'expression mathématique du texte"""
        # Chercher f(x) = ...
        pattern = r'f\(x\)\s*=\s*([^\n,\.]+)'
        match = re.search(pattern, text, re.IGNORECASE)
        if match:
            expr = match.group(1).strip()
            # Nettoyer l'expression
            expr = expr.replace('^', '**')  # Puissances
            expr = expr.replace('²', '**2')
            expr = expr.replace('³', '**3')
            return expr
        
        # Chercher directement une expression
        patterns = [
            r'x\*\*\d+',  # x**2, x**3
            r'x²|x³',      # Unicode
            r'\d*x[\+\-\*/]',  # 2x+3
            r'sin\(x\)|cos\(x\)|tan\(x\)',  # Trigonométrie
        ]
        for pattern in patterns:
            match = re.search(pattern, text, re.IGNORECASE)
            if match:
                return match.group(0)
        
        return None
    
    def generate_graph_inline(expression: str) -> dict:
        """Génère un graphique mathématique"""
        try:
            # Préparer l'expression pour numpy
            expr_safe = expression.replace('x', 'x_values')
            expr_safe = expr_safe.replace('sin', 'np.sin')
            expr_safe = expr_safe.replace('cos', 'np.cos')
            expr_safe = expr_safe.replace('tan', 'np.tan')
            expr_safe = expr_safe.replace('exp', 'np.exp')
            expr_safe = expr_safe.replace('log', 'np.log')
            expr_safe = expr_safe.replace('sqrt', 'np.sqrt')
            
            # Générer les valeurs
            x_values = np.linspace(-10, 10, 400)
            y_values = eval(expr_safe)
            
            # Créer le graphique
            plt.figure(figsize=(10, 6))
            plt.plot(x_values, y_values, 'b-', linewidth=2, label=f'f(x) = {expression}')
            plt.grid(True, alpha=0.3)
            plt.axhline(y=0, color='k', linewidth=0.5)
            plt.axvline(x=0, color='k', linewidth=0.5)
            plt.xlabel('x', fontsize=12)
            plt.ylabel('f(x)', fontsize=12)
            plt.title(f'Graphique de f(x) = {expression}', fontsize=14, fontweight='bold')
            plt.legend()
            
            # Convertir en base64
            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=100, bbox_inches='tight')
            plt.close()
            
            buffer.seek(0)
            image_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            return {
                'expression': expression,
                'graph_base64': image_base64,
                'markdown': f"![Graphique](data:image/png;base64,{image_base64})"
            }
            
        except Exception as e:
            logging.error(f"Erreur génération graphique: {e}")
            return None
    
    def process_graph_request_inline(message: str) -> dict:
        """Traite une demande de graphique"""
        if not detect_graph_request_inline(message):
            return None
        
        expression = extract_math_expression_inline(message)
        
        # Si pas d'expression trouvée, utiliser une fonction par défaut pour démonstration
        if not expression:
            # Pour les demandes génériques, proposer un exemple
            logging.info("Demande de graphique sans expression spécifique, utilisation d'un exemple")
            expression = "x**2"  # Parabole simple comme exemple
        
        result = generate_graph_inline(expression)
        if result:
            logging.info(f"✅ Graphique généré pour: {expression}")
        
        return result
    
    # ==================== GÉNÉRATEUR D'ANGLES GÉOMÉTRIQUES ====================
    
    def generate_angle_diagram(angle_type: str) -> dict:
        """Génère un diagramme d'angle carré pour utilisation dans Word"""
        try:
            from matplotlib.patches import Wedge
            
            # Image carrée de taille appropriée pour Word
            fig, ax = plt.subplots(1, 1, figsize=(5, 5))
            ax.set_aspect('equal')
            ax.axis('off')
            
            # Couleurs pastel
            color_yellow = '#FFE566'
            color_blue = '#99CCFF'
            color_green = '#99FF99'
            color_pink = '#FFB3B3'
            
            def draw_wedge(center, radius, t1, t2, color):
                w = Wedge(center, radius, t1, t2, facecolor=color, edgecolor='black', lw=1.5, alpha=0.7)
                ax.add_patch(w)
            
            if angle_type == 'adjacents':
                ax.set_xlim(-0.5, 4)
                ax.set_ylim(-0.5, 4)
                O = (1.5, 1.5)
                ax.plot([O[0], 3.8], [O[1], 1.5], 'k-', lw=2)
                ax.plot([O[0], 3.2], [O[1], 3.2], 'k-', lw=2)
                ax.plot([O[0], 1.5], [O[1], 3.8], 'k-', lw=2)
                draw_wedge(O, 0.5, 0, 45, color_green)
                draw_wedge(O, 0.5, 45, 90, color_blue)
                ax.plot(O[0], O[1], 'ko', markersize=8)
                ax.text(O[0]-0.2, O[1]-0.35, 'O', fontsize=14, fontweight='bold')
                ax.text(3.8, 1.3, 'x', fontsize=14, fontweight='bold')
                ax.text(3.3, 3.3, 'y', fontsize=14, fontweight='bold')
                ax.text(1.3, 3.9, 'z', fontsize=14, fontweight='bold')
                ax.set_title('Angles adjacents', fontsize=14, fontweight='bold', pad=10)
                
            elif angle_type == 'complementaires':
                ax.set_xlim(-0.5, 4)
                ax.set_ylim(-0.5, 4)
                O = (1.2, 1.2)
                ax.plot([O[0], 3.8], [O[1], 1.2], 'k-', lw=2)
                ax.plot([O[0], 1.2], [O[1], 3.8], 'k-', lw=2)
                ax.plot([O[0], 3], [O[1], 2.5], 'k-', lw=2)
                # Angle droit
                ax.plot([O[0], O[0]+0.25], [O[1], O[1]], 'k-', lw=1.5)
                ax.plot([O[0]+0.25, O[0]+0.25], [O[1], O[1]+0.25], 'k-', lw=1.5)
                draw_wedge(O, 0.5, 0, 32, color_yellow)
                draw_wedge(O, 0.5, 32, 90, color_blue)
                ax.text(2.0, 1.35, '32°', fontsize=12, fontweight='bold')
                ax.text(1.35, 2.0, '58°', fontsize=12, fontweight='bold')
                ax.plot(O[0], O[1], 'ko', markersize=8)
                ax.text(O[0]-0.25, O[1]-0.35, 'A', fontsize=14, fontweight='bold')
                ax.set_title('Angles complémentaires (90°)', fontsize=14, fontweight='bold', pad=10)
                
            elif angle_type == 'supplementaires':
                ax.set_xlim(-2, 4)
                ax.set_ylim(-0.5, 4)
                O = (1, 1.5)
                ax.plot([-1.5, 3.5], [1.5, 1.5], 'k-', lw=2)
                ax.plot([O[0], 2.8], [O[1], 3.3], 'k-', lw=2)
                draw_wedge(O, 0.5, 0, 50, color_green)
                draw_wedge(O, 0.5, 50, 180, color_pink)
                ax.text(1.7, 1.65, '50°', fontsize=12, fontweight='bold')
                ax.text(0.2, 1.9, '130°', fontsize=12, fontweight='bold')
                ax.plot(O[0], O[1], 'ko', markersize=8)
                ax.text(O[0]-0.15, O[1]-0.35, 'O', fontsize=14, fontweight='bold')
                ax.text(3.5, 1.3, 'x', fontsize=14, fontweight='bold')
                ax.text(-1.6, 1.3, "x'", fontsize=14, fontweight='bold')
                ax.set_title('Angles supplémentaires (180°)', fontsize=14, fontweight='bold', pad=10)
                
            elif angle_type == 'opposes_sommet':
                ax.set_xlim(-2.5, 2.5)
                ax.set_ylim(-2.5, 2.5)
                O = (0, 0)
                ax.plot([-2, 2], [-1.2, 1.2], 'k-', lw=2)
                ax.plot([-2, 2], [1, -1], 'k-', lw=2)
                draw_wedge(O, 0.45, -31, 27, color_yellow)
                draw_wedge(O, 0.45, 149, 207, color_yellow)
                draw_wedge(O, 0.45, 27, 149, color_blue)
                draw_wedge(O, 0.45, 207, 329, color_blue)
                ax.plot(O[0], O[1], 'ko', markersize=8)
                ax.text(0.5, 0.15, 'α', fontsize=14, color='#B8860B', fontweight='bold')
                ax.text(-0.7, -0.25, 'α', fontsize=14, color='#B8860B', fontweight='bold')
                ax.text(0.15, 0.6, 'β', fontsize=14, color='#4169E1', fontweight='bold')
                ax.text(-0.35, -0.7, 'β', fontsize=14, color='#4169E1', fontweight='bold')
                ax.set_title('Angles opposés par le sommet', fontsize=14, fontweight='bold', pad=10)
                
            elif angle_type == 'alternes_internes':
                ax.set_xlim(-2.5, 3.5)
                ax.set_ylim(-1.5, 4)
                ax.plot([-2, 3], [2.5, 2.5], '--', color='gray', lw=2)
                ax.plot([-2, 3], [0.8, 0.8], '--', color='gray', lw=2)
                ax.plot([-0.5, 2.2], [-0.5, 3.5], 'k-', lw=2)
                A = (1.1, 2.5)
                B = (0.47, 0.8)
                draw_wedge(A, 0.4, 180, 243, color_yellow)
                draw_wedge(B, 0.4, 0, 63, color_yellow)
                ax.plot(A[0], A[1], 'ko', markersize=6)
                ax.plot(B[0], B[1], 'ko', markersize=6)
                ax.text(A[0]-0.15, A[1]+0.15, 'A', fontsize=12, fontweight='bold')
                ax.text(B[0]-0.15, B[1]-0.35, 'B', fontsize=12, fontweight='bold')
                ax.text(-2.2, 2.65, 'd₁', fontsize=12, fontweight='bold')
                ax.text(-2.2, 0.95, 'd₂', fontsize=12, fontweight='bold')
                ax.text(2.3, 3.5, 't', fontsize=12, fontweight='bold')
                ax.set_title('Angles alternes-internes', fontsize=14, fontweight='bold', pad=10)
                
            elif angle_type == 'alternes_externes':
                ax.set_xlim(-2.5, 3.5)
                ax.set_ylim(-1.5, 4)
                ax.plot([-2, 3], [2.5, 2.5], '--', color='gray', lw=2)
                ax.plot([-2, 3], [0.8, 0.8], '--', color='gray', lw=2)
                ax.plot([-0.5, 2.2], [-0.5, 3.5], 'k-', lw=2)
                A = (1.1, 2.5)
                B = (0.47, 0.8)
                draw_wedge(A, 0.4, 0, 63, color_pink)
                draw_wedge(B, 0.4, 180, 243, color_pink)
                ax.plot(A[0], A[1], 'ko', markersize=6)
                ax.plot(B[0], B[1], 'ko', markersize=6)
                ax.text(A[0]+0.2, A[1]+0.15, 'A', fontsize=12, fontweight='bold')
                ax.text(B[0]-0.35, B[1]-0.15, 'B', fontsize=12, fontweight='bold')
                ax.text(-2.2, 2.65, 'd₁', fontsize=12, fontweight='bold')
                ax.text(-2.2, 0.95, 'd₂', fontsize=12, fontweight='bold')
                ax.set_title('Angles alternes-externes', fontsize=14, fontweight='bold', pad=10)
                
            elif angle_type == 'correspondants':
                ax.set_xlim(-2.5, 3.5)
                ax.set_ylim(-1.5, 4)
                ax.plot([-2, 3], [2.5, 2.5], '--', color='gray', lw=2)
                ax.plot([-2, 3], [0.8, 0.8], '--', color='gray', lw=2)
                ax.plot([-0.5, 2.2], [-0.5, 3.5], 'k-', lw=2)
                A = (1.1, 2.5)
                B = (0.47, 0.8)
                draw_wedge(A, 0.4, 0, 63, color_blue)
                draw_wedge(B, 0.4, 0, 63, color_blue)
                ax.plot(A[0], A[1], 'ko', markersize=6)
                ax.plot(B[0], B[1], 'ko', markersize=6)
                ax.text(A[0]+0.2, A[1]+0.15, 'A', fontsize=12, fontweight='bold')
                ax.text(B[0]+0.2, B[1]-0.35, 'B', fontsize=12, fontweight='bold')
                ax.text(-2.2, 2.65, 'd₁', fontsize=12, fontweight='bold')
                ax.text(-2.2, 0.95, 'd₂', fontsize=12, fontweight='bold')
                ax.set_title('Angles correspondants', fontsize=14, fontweight='bold', pad=10)
            
            else:
                ax.set_xlim(-1, 1)
                ax.set_ylim(-1, 1)
                ax.text(0, 0, f'{angle_type}', fontsize=14, ha='center')
            
            # Convertir en base64
            buffer = BytesIO()
            plt.savefig(buffer, format='png', dpi=120, bbox_inches='tight', facecolor='white', edgecolor='none')
            plt.close()
            
            buffer.seek(0)
            image_base64 = base64.b64encode(buffer.getvalue()).decode()
            
            return {
                'angle_type': angle_type,
                'image_base64': image_base64
            }
            
        except Exception as e:
            logging.error(f"Erreur génération diagramme d'angles: {e}")
            return None
    
    def generate_multiple_angles() -> list:
        """Génère plusieurs images d'angles individuelles"""
        angle_types = ['adjacents', 'complementaires', 'supplementaires', 'opposes_sommet', 'alternes_internes', 'correspondants']
        results = []
        for angle_type in angle_types:
            result = generate_angle_diagram(angle_type)
            if result:
                results.append(result)
        return results
    
    def detect_angle_request(message: str) -> str:
        """Détecte si le message demande un diagramme d'angles et retourne le type"""
        message_lower = message.lower()
        
        # Vérifier d'abord si c'est une demande liée aux angles
        angle_keywords = ['angle', 'angles', 'géométrie', 'geometrie', 'droites parallèles', 
                         'droites paralleles', 'sécante', 'secante', 'sommet']
        
        is_angle_request = any(kw in message_lower for kw in angle_keywords)
        
        if not is_angle_request:
            return None
        
        # Patterns pour détecter les types d'angles spécifiques
        if any(word in message_lower for word in ['adjacent', 'adjacents']):
            return 'adjacents'
        elif any(word in message_lower for word in ['complémentaire', 'complementaire', 'complémentaires', 'complementaires']):
            return 'complementaires'
        elif any(word in message_lower for word in ['supplémentaire', 'supplementaire', 'supplémentaires', 'supplementaires']):
            return 'supplementaires'
        elif any(word in message_lower for word in ['opposé par le sommet', 'opposés par le sommet', 
                                                     'oppose par le sommet', 'opposes par le sommet',
                                                     'opposé au sommet', 'opposés au sommet']):
            return 'opposes_sommet'
        elif any(word in message_lower for word in ['alterne-interne', 'alternes-internes', 
                                                     'alterne interne', 'alternes internes',
                                                     'alterne intérieur', 'alternes intérieurs']):
            return 'alternes_internes'
        elif any(word in message_lower for word in ['alterne-externe', 'alternes-externes', 
                                                     'alterne externe', 'alternes externes',
                                                     'alterne extérieur', 'alternes extérieurs']):
            return 'alternes_externes'
        elif any(word in message_lower for word in ['correspondant', 'correspondants']):
            return 'correspondants'
        
        # Si demande générique d'angles avec image/schéma → générer TOUS les angles
        if any(word in message_lower for word in ['image', 'diagramme', 'schéma', 'schema', 
                                                   'dessine', 'montre', 'génère', 'genere',
                                                   'crée', 'cree', 'affiche', 'voir', 'tous',
                                                   'all', 'différents', 'types']):
            return 'all'
        
        # Demande d'angles sans précision → générer tous
        return 'all'
    
except ImportError as e:
    GRAPH_GENERATOR_AVAILABLE = False
    logging.warning(f"⚠️ Graphiques mathématiques désactivés: {e}")
    
    def process_graph_request_inline(message: str) -> dict:
        return None

# === FIN CODE GRAPHIQUES INLINE ===


ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

# MongoDB connection
mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

# Configuration de Google Gemini (100% gratuit)
genai.configure(api_key=os.environ.get('GOOGLE_API_KEY', ''))

# Initialiser le detecteur IA avance avec auto-apprentissage
advanced_detector = AdvancedAIDetector()

# Charger les poids depuis la base de donnees au demarrage
async def load_detector_weights():
    """Charge les poids du detecteur depuis MongoDB"""
    try:
        weights_doc = await db.detector_weights.find_one({"version": "latest"})
        if weights_doc and "weights" in weights_doc:
            advanced_detector.set_weights(weights_doc["weights"])
            logging.info("Poids du detecteur charges depuis la base de donnees")
        else:
            logging.info("Utilisation des poids par defaut du detecteur")
    except Exception as e:
        logging.error(f"Erreur chargement poids: {e}")

# Sauvegarder les poids dans la base de donnees
async def save_detector_weights():
    """Sauvegarde les poids du detecteur dans MongoDB"""
    try:
        weights = advanced_detector.get_weights()
        await db.detector_weights.update_one(
            {"version": "latest"},
            {
                "$set": {
                    "weights": weights,
                    "updated_at": datetime.now(timezone.utc)
                }
            },
            upsert=True
        )
        logging.info("Poids du detecteur sauvegardes")
    except Exception as e:
        logging.error(f"Erreur sauvegarde poids: {e}")

# Create the main app without a prefix
app = FastAPI()

# Create a router with the /api prefix
api_router = APIRouter(prefix="/api")


# Define Models
class ChatMessage(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: str
    message: str
    response: str
    message_type: str  # "je_veux", "je_recherche", "sources_fiables", "activites"
    trust_score: Optional[float] = None
    sources: Optional[List[str]] = None
    image_base64: Optional[str] = None  # Pour les images generees
    images: Optional[List[str]] = None  # Pour plusieurs images
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class ChatRequest(BaseModel):
    message: str
    message_type: str
    session_id: Optional[str] = None
    conversation_id: Optional[str] = None  # Pour lier à une conversation
    conversation_history: Optional[List[dict]] = None  # Historique pour la mémoire

# ==================== MODELES POUR CONVERSATIONS CLOUD ====================

class ConversationCreate(BaseModel):
    title: str
    first_message: Optional[str] = None

class ConversationMessage(BaseModel):
    role: str  # "user" ou "assistant"
    content: str
    timestamp: Optional[str] = None
    image_base64: Optional[str] = None

class ConversationUpdate(BaseModel):
    title: Optional[str] = None
    is_pinned: Optional[bool] = None
    is_favorite: Optional[bool] = None
    tags: Optional[List[str]] = None

# ==================== MODELES POUR RESET MOT DE PASSE ====================

class PasswordResetRequest(BaseModel):
    email: str

class PasswordResetConfirm(BaseModel):
    token: str
    new_password: str

class SearchResult(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    query: str
    sources: List[dict]  # [{"url": str, "title": str, "trust_score": float, "content_preview": str}]
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

class DocumentRequest(BaseModel):
    content: str
    title: str = "Document Etienne"
    format: str = "pdf"  # pdf, docx, pptx, xlsx
    filename: Optional[str] = None

class TextAnalysisRequest(BaseModel):
    text: str
    analysis_type: str = "complete"  # "ai_detection", "plagiarism", "complete"

class FileAnalysisRequest(BaseModel):
    question: str
    extracted_text: str
    filename: str
    message_type: str

class SecurityIncidentLog(BaseModel):
    id: str = Field(default_factory=lambda: str(uuid.uuid4()))
    session_id: Optional[str] = None
    user_message: str
    detected_category: str  # "violence", "illegal_activities", "hacking", "drugs", etc.
    keywords_matched: List[str]
    severity: str  # "low", "medium", "high", "critical"
    ip_address: Optional[str] = None
    user_agent: Optional[str] = None
    blocked: bool = True
    timestamp: datetime = Field(default_factory=lambda: datetime.now(timezone.utc))

# Systeme de confiance des sources
TRUSTED_DOMAINS = {
    ".gouv.qc.ca": 0.98,
    ".gouv.ca": 0.95,
    ".edu": 0.90,
    "quebec.ca": 0.97,
    "education.gouv.qc.ca": 0.98,
    "mees.gouv.qc.ca": 0.98,  # Ministere de l Education du Quebec
    "banq.qc.ca": 0.88,  # Bibliotheque nationale du Quebec
    "uqam.ca": 0.85,
    "umontreal.ca": 0.85,
    "ulaval.ca": 0.85,
    "mcgill.ca": 0.85,
    "cegep": 0.75,
    "universit": 0.75,
    ".ca": 0.70,
    ".org": 0.60,
    ".com": 0.40,
    "wikipedia": 0.65
}

def calculate_trust_score(url: str, content: str = "") -> float:
    """Calcule le score de confiance de une source base sur le domaine et le contenu"""
    base_score = 0.5
    
    # Verification du domaine
    for domain, score in TRUSTED_DOMAINS.items():
        if domain in url.lower():
            base_score = max(base_score, score)
    
    # Analyse basique du contenu si fourni
    if content:
        quality_indicators = [
            "bibliographie", "references", "source", "etude", "recherche", 
            "academique", "officiel", "ministere", "universite", "peer-review"
        ]
        quality_count = sum(1 for indicator in quality_indicators if indicator in content.lower())
        content_bonus = min(0.15, quality_count * 0.03)
        base_score = min(0.98, base_score + content_bonus)
    
    return round(base_score, 2)

async def extract_text_from_file(file: UploadFile) -> str:
    """Extrait le texte de un fichier uploade - VERSION OPTIMISÉE"""
    try:
        # Creer un fichier temporaire
        with tempfile.NamedTemporaryFile(delete=False, suffix=f"_{file.filename}") as temp_file:
            content = await file.read()
            temp_file.write(content)
            temp_file_path = temp_file.name
        
        extracted_text = ""
        file_extension = file.filename.lower().split('.')[-1] if '.' in file.filename else ''
        
        try:
            if file_extension == 'pdf':
                # Extraction PDF ULTRA-OPTIMISÉE avec traitement asynchrone
                def extract_pdf_sync(path):
                    """Extraction PDF synchrone dans un thread séparé"""
                    text_parts = []
                    try:
                        # Essayer PyPDF2 en premier (plus rapide que pdfplumber)
                        with open(path, 'rb') as pdf_file:
                            pdf_reader = PyPDF2.PdfReader(pdf_file)
                            total_pages = len(pdf_reader.pages)
                            
                            # Limiter à 50 pages pour les gros PDFs (optimisation timeout)
                            # Pour PDFs > 50 pages, traiter uniquement les premières pages
                            max_pages = min(total_pages, 50)
                            
                            logging.info(f"PDF: {total_pages} pages, extraction de {max_pages} pages")
                            
                            for i in range(max_pages):
                                try:
                                    page_text = pdf_reader.pages[i].extract_text()
                                    if page_text and page_text.strip():
                                        text_parts.append(page_text)
                                    
                                    # Pause tous les 10 pages pour éviter blocage
                                    if (i + 1) % 10 == 0:
                                        logging.info(f"Traité {i + 1}/{max_pages} pages")
                                        
                                except Exception as page_error:
                                    logging.warning(f"Erreur page {i+1}: {page_error}")
                                    continue
                            
                            if total_pages > max_pages:
                                text_parts.append(f"\n\n[📄 Note: {max_pages} premières pages extraites sur {total_pages} pages totales. Pour analyser plus de pages, divisez le document ou utilisez un PDF plus court.]\n")
                            
                            return "\n".join(text_parts)
                    
                    except Exception as pypdf_error:
                        logging.warning(f"PyPDF2 échoué: {pypdf_error}, essai pdfplumber...")
                        # Fallback avec pdfplumber (plus lent mais meilleur pour PDFs complexes)
                        try:
                            text_parts = []
                            with pdfplumber.open(path) as pdf:
                                total_pages = len(pdf.pages)
                                max_pages = min(total_pages, 30)  # Encore plus limité pour pdfplumber
                                
                                logging.info(f"pdfplumber: {total_pages} pages, extraction de {max_pages} pages")
                                
                                for i in range(max_pages):
                                    try:
                                        page_text = pdf.pages[i].extract_text()
                                        if page_text:
                                            text_parts.append(page_text)
                                    except:
                                        continue
                                
                                if total_pages > max_pages:
                                    text_parts.append(f"\n\n[📄 Note: {max_pages} premières pages extraites sur {total_pages} pages totales.]\n")
                                
                                return "\n".join(text_parts)
                        
                        except Exception as plumber_error:
                            logging.error(f"pdfplumber aussi échoué: {plumber_error}")
                            raise HTTPException(
                                status_code=500,
                                detail="Impossible d'extraire le texte de ce PDF. Le fichier est peut-être corrompu, protégé, ou contient uniquement des images."
                            )
                
                # Exécuter l'extraction PDF dans un thread séparé (non-bloquant)
                extracted_text = await asyncio.to_thread(extract_pdf_sync, temp_file_path)
            
            elif file_extension in ['docx', 'doc']:
                # Extraction Word
                extracted_text = docx2txt.process(temp_file_path)
            
            elif file_extension == 'txt':
                # Fichier texte simple
                with open(temp_file_path, 'r', encoding='utf-8') as f:
                    extracted_text = f.read()
            
            elif file_extension in ['xlsx', 'xls']:
                # Extraction Excel
                try:
                    df = pd.read_excel(temp_file_path, sheet_name=None)  # Toutes les feuilles
                    for sheet_name, sheet_data in df.items():
                        extracted_text += f"\n=== Feuille: {sheet_name} ===\n"
                        extracted_text += sheet_data.to_string(index=False, na_rep='') + "\n"
                except Exception:
                    # Fallback pour anciens formats Excel
                    df = pd.read_excel(temp_file_path, engine='openpyxl')
                    extracted_text = df.to_string(index=False, na_rep='')
            
            elif file_extension == 'csv':
                # Fichier CSV
                df = pd.read_csv(temp_file_path)
                extracted_text = df.to_string(index=False, na_rep='')
            
            elif file_extension == 'pptx':
                # PowerPoint (extraction basique)
                from pptx import Presentation
                prs = Presentation(temp_file_path)
                for slide_num, slide in enumerate(prs.slides, 1):
                    extracted_text += f"\n=== Slide {slide_num} ===\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            
            elif file_extension in ['png', 'jpg', 'jpeg', 'bmp', 'tiff', 'webp']:
                # Extraction de texte d'image via Gemini Vision (multimodal)
                try:
                    img = Image.open(temp_file_path)
                    model = genai.GenerativeModel('gemini-2.5-flash')
                    response = model.generate_content(
                        [
                            "Extrais TOUT le texte visible dans cette image. "
                            "Retourne UNIQUEMENT le texte brut, sans commentaire ni explication. "
                            "Respecte la mise en page originale (paragraphes, retours à la ligne).",
                            img
                        ]
                    )
                    extracted_text = response.text.strip() if response.text else ""
                    if not extracted_text:
                        raise Exception("Aucun texte détecté dans l'image")
                except Exception as ocr_error:
                    logging.error(f"Erreur OCR Gemini: {ocr_error}")
                    raise HTTPException(
                        status_code=400,
                        detail="Impossible de lire le texte de cette image. Assurez-vous que l'image contient du texte lisible."
                    )
            
            else:
                raise HTTPException(
                    status_code=400, 
                    detail=f"Format de fichier non supporte: {file_extension}. "
                          f"Formats acceptes: PDF, DOCX, TXT, XLSX, CSV, PPTX, PNG, JPG, JPEG"
                )
        
        finally:
            # Nettoyer le fichier temporaire
            if os.path.exists(temp_file_path):
                os.unlink(temp_file_path)
        
        if not extracted_text.strip():
            raise HTTPException(
                status_code=400,
                detail="Impossible de extraire le texte de ce fichier. Verifiez que le fichier n'est pas protege ou corrompu."
            )
        
        return extracted_text.strip()
    
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur extraction texte: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Erreur lors de l extraction du texte: {str(e)}"
        )

# ============================================================================
# DETECTION IA GRATUITE - HuggingFace + Heuristique Avancee
# Precision cible: 75-85% sans coût
# ============================================================================

async def detect_ai_content_with_huggingface(text: str) -> dict:
    """
    Detection IA gratuite utilisant HuggingFace Inference API.
    Utilise un modele RoBERTa fine-tune specifiquement pour detecter le contenu IA.
    """
    try:
        import aiohttp
        
        # Limiter la longueur du texte
        max_text_length = 2000
        analyzed_text = text[:max_text_length] if len(text) > max_text_length else text
        
        # Modele gratuit specialise pour detection IA sur HuggingFace
        # Ce modele est entraîne specifiquement pour detecter GPT/Claude/etc.
        API_URL = "https://api-inference.huggingface.co/models/Hello-SimpleAI/chatgpt-detector-roberta"
        
        async with aiohttp.ClientSession() as session:
            async with session.post(
                API_URL,
                headers={"Content-Type": "application/json"},
                json={"inputs": analyzed_text},
                timeout=aiohttp.ClientTimeout(total=10)
            ) as response:
                if response.status == 200:
                    result = await response.json()
                    
                    # Le modele retourne [{"label": "Real/Fake", "score": 0.XX}]
                    if isinstance(result, list) and len(result) > 0:
                        # Trouver le score "Fake" (AI-generated)
                        ai_score = 0.5
                        for item in result:
                            if item.get("label") == "Fake":
                                ai_score = item.get("score", 0.5)
                                break
                        
                        # Determiner les patterns bases sur le score
                        detected_patterns = []
                        if ai_score > 0.7:
                            detected_patterns = ["high_ai_confidence", "model_detected_fake", "consistent_patterns"]
                        elif ai_score > 0.5:
                            detected_patterns = ["moderate_ai_indicators", "some_patterns_detected"]
                        else:
                            detected_patterns = ["likely_human_written"]
                        
                        return {
                            "ai_probability": round(ai_score, 2),
                            "is_likely_ai": ai_score > 0.5,
                            "confidence": "High" if ai_score > 0.7 or ai_score < 0.3 else "Medium",
                            "detected_patterns": detected_patterns,
                            "reasoning": f"Modele RoBERTa detecte {int(ai_score*100)}% de probabilite IA",
                            "method": "huggingface_roberta"
                        }
                
                # Si l API HuggingFace echoue, utiliser le fallback
                logging.warning(f"HuggingFace API status: {response.status}")
                return detect_ai_content_fallback(text)
                
    except asyncio.TimeoutError:
        logging.warning("HuggingFace API timeout - using fallback")
        return detect_ai_content_fallback(text)
    except Exception as e:
        logging.error(f"Erreur HuggingFace API: {e}")
        return detect_ai_content_fallback(text)


async def detect_ai_content_with_llm(text: str) -> dict:
    """
    Point de entree principal pour la detection IA.
    Essaie de abord HuggingFace (gratuit), puis fallback vers heuristique avancee.
    """
    # Essayer de abord HuggingFace
    result = await detect_ai_content_with_huggingface(text)
    
    # Si la methode est deja le fallback, le retourner directement
    if result.get("method") == "heuristic_advanced":
        return result
    
    # Si HuggingFace a fonctionne mais avec faible confiance, combiner avec heuristique
    if result.get("confidence") == "Medium":
        heuristic_result = detect_ai_content_fallback(text)
        
        # Moyenne ponderee : 70% HuggingFace, 30% heuristique
        combined_probability = (result["ai_probability"] * 0.7) + (heuristic_result["ai_probability"] * 0.3)
        
        return {
            "ai_probability": round(combined_probability, 2),
            "is_likely_ai": combined_probability > 0.5,
            "confidence": "High" if combined_probability > 0.7 or combined_probability < 0.3 else "Medium",
            "detected_patterns": result["detected_patterns"][:2] + heuristic_result["detected_patterns"][:1],
            "reasoning": f"Analyse combinee HuggingFace + heuristique: {int(combined_probability*100)}%",
            "method": "hybrid_huggingface_heuristic"
        }
    
    return result


def detect_ai_content_fallback(text: str) -> dict:
    """
    Detection IA ultra-avancee avec 15+ analyses et auto-apprentissage
    Precision: 80-85% (ameliore au fil du temps)
    """
    return advanced_detector.analyze_text(text)

# ============================================================================
# FIN DE LA NOUVELLE DETECTION IA
# ============================================================================


def check_plagiarism(text: str) -> dict:
    """Verificateur de plagiat basique"""
    try:
        # Phrases communes qui peuvent indiquer du plagiat
        common_academic_phrases = [
            "according to the study", "research shows that", "studies have shown",
            "it has been proven that", "experts agree that", "the data suggests",
            "furthermore", "in addition", "however", "therefore", "consequently"
        ]
        
        # Verification de phrases trop parfaites/academiques
        text_lower = text.lower()
        academic_score = 0
        found_phrases = []
        
        for phrase in common_academic_phrases:
            if phrase in text_lower:
                academic_score += 0.1
                found_phrases.append(phrase)
        
        # Verification de la diversite du vocabulaire
        words = text_lower.split()
        unique_words = set(words)
        vocabulary_diversity = len(unique_words) / len(words) if words else 0
        
        # Score de risque de plagiat
        plagiarism_risk = min(academic_score, 0.9)
        if vocabulary_diversity < 0.6:  # Faible diversite = risque
            plagiarism_risk += 0.1
        
        plagiarism_risk = min(plagiarism_risk, 0.99)
        
        return {
            "plagiarism_risk": round(plagiarism_risk, 2),
            "is_suspicious": plagiarism_risk > 0.4,
            "vocabulary_diversity": round(vocabulary_diversity, 2),
            "risk_level": "High" if plagiarism_risk > 0.6 else "Medium" if plagiarism_risk > 0.3 else "Low",
            "found_phrases": found_phrases[:3],
            "recommendation": "Verifiez l originalite avec des sources academiques" if plagiarism_risk > 0.4 else "Contenu semble original"
        }
        
    except Exception as e:
        return {
            "plagiarism_risk": 0.0,
            "is_suspicious": False,
            "error": str(e)
        }

# Detection de langue pour reponses adaptees
def detect_language(text: str) -> str:
    """Detecte la langue du message - détection améliorée pour enseignants"""
    text_lower = text.lower().strip()
    
    # Si le texte est très court (moins de 10 caractères), défaut français
    if len(text_lower) < 10:
        return "fr"
    
    # ÉTAPE 1: Phrases complètes en anglais (très fiables)
    english_phrases = [
        "can you help me", "could you help", "i need help", "i would like",
        "please help me", "i am looking for", "i want to", "how can i",
        "what is the", "what are the", "how do i", "how to",
        "create a lesson", "make a quiz", "write an", "generate a",
        "i'm teaching", "my students", "for my class", "in my classroom",
        "thank you", "thanks for", "hello", "hi there", "good morning",
        "i have a question", "can you explain", "could you create",
        "i need a", "please create", "please make", "please write"
    ]
    
    # ÉTAPE 2: Phrases complètes en français (très fiables)
    french_phrases = [
        "peux-tu m'aider", "pouvez-vous m'aider", "j'ai besoin", "j'aimerais",
        "s'il vous plaît", "s'il te plaît", "je cherche", "je veux",
        "comment puis-je", "qu'est-ce que", "comment faire", "comment est-ce",
        "crée un cours", "fais un quiz", "écris un", "génère un",
        "j'enseigne", "mes élèves", "pour ma classe", "dans ma classe",
        "merci", "merci pour", "bonjour", "salut", "bonsoir", "bon matin",
        "j'ai une question", "peux-tu expliquer", "pouvez-vous créer",
        "j'ai besoin d'un", "crée-moi", "fais-moi", "écris-moi",
        "pour mes élèves", "pour les élèves", "au secondaire", "en secondaire",
        "cours de français", "cours de math", "cours d'anglais", "cours de science",
        "je suis enseignant", "je suis enseignante", "je suis prof"
    ]
    
    # Vérifier phrases anglaises
    for phrase in english_phrases:
        if phrase in text_lower:
            return "en"
    
    # Vérifier phrases françaises
    for phrase in french_phrases:
        if phrase in text_lower:
            return "fr"
    
    # ÉTAPE 3: Compter les mots distinctifs avec contexte
    # Mots UNIQUEMENT anglais (jamais utilisés en français)
    english_only = [
        " the ", " is ", " are ", " am ", " was ", " were ", " been ", " being ",
        " have ", " has ", " had ", " do ", " does ", " did ", " will ", " would ",
        " could ", " should ", " can ", " may ", " might ", " must ",
        " this ", " that ", " these ", " those ", " there ", " here ",
        " what ", " which ", " who ", " whom ", " whose ", " where ", " when ",
        " why ", " how ", " if ", " then ", " than ", " because ", " although ",
        " however ", " therefore ", " moreover ", " furthermore ",
        " my ", " your ", " his ", " her ", " its ", " our ", " their ",
        " myself ", " yourself ", " himself ", " herself ", " itself ",
        " and ", " or ", " but ", " not ", " no ", " yes ", " very ",
        " about ", " after ", " before ", " between ", " during ", " without "
    ]
    
    # Mots UNIQUEMENT français (jamais utilisés en anglais)
    french_only = [
        " le ", " la ", " les ", " un ", " une ", " des ", " du ", " de la ",
        " est ", " sont ", " suis ", " sommes ", " êtes ", " était ", " étaient ",
        " avoir ", " avons ", " avez ", " ont ", " eu ", " été ",
        " faire ", " fais ", " fait ", " faisons ", " faites ", " font ",
        " ce ", " cette ", " ces ", " cet ", " celui ", " celle ", " ceux ",
        " qui ", " que ", " quoi ", " dont ", " où ", " lequel ", " laquelle ",
        " mon ", " ma ", " mes ", " ton ", " ta ", " tes ", " son ", " sa ", " ses ",
        " notre ", " nos ", " votre ", " vos ", " leur ", " leurs ",
        " et ", " ou ", " mais ", " donc ", " car ", " ni ", " or ",
        " ne ", " pas ", " plus ", " moins ", " très ", " bien ", " mal ",
        " dans ", " sur ", " sous ", " avec ", " sans ", " pour ", " par ", " chez ",
        " je ", " tu ", " il ", " elle ", " nous ", " vous ", " ils ", " elles ",
        " moi ", " toi ", " lui ", " eux ", " y ", " en ",
        " aussi ", " encore ", " toujours ", " jamais ", " souvent ", " parfois "
    ]
    
    text_spaced = f" {text_lower} "
    
    english_count = sum(1 for word in english_only if word in text_spaced)
    french_count = sum(1 for word in french_only if word in text_spaced)
    
    # ÉTAPE 4: Décision finale
    # Il faut une différence significative pour basculer en anglais
    if english_count >= 5 and english_count > french_count * 2:
        return "en"
    elif french_count >= 3:
        return "fr"
    elif english_count > french_count + 2:
        return "en"
    else:
        return "fr"  # Par défaut français pour app québécoise


# ============================================
# DETECTION DE CONTENU ILLEGAL ET DANGEREUX
# ============================================

ILLEGAL_CONTENT_PATTERNS = {
    "violence": {
        "keywords": [
            "tuer", "assassiner", "meurtre", "comment tuer", "faire du mal", "blesser quelqu'un",
            "torturer", "violence physique", "attaque", "agresser", "frapper quelqu'un",
            "murder", "kill someone", "how to kill", "hurt someone", "attack someone"
        ],
        "severity": "critical"
    },
    "explosives": {
        "keywords": [
            "fabriquer une bombe", "faire exploser", "explosif", "bombe artisanale",
            "faire une bombe", "creer une bombe", "dynamite", "detonateur",
            "make a bomb", "build a bomb", "explosive device", "homemade bomb"
        ],
        "severity": "critical"
    },
    "hacking": {
        "keywords": [
            "pirater", "hacker", "voler des donnees", "acces illegal", "briser un mot de passe",
            "cracker un systeme", "ddos", "intrusion", "exploit", "backdoor",
            "hack into", "break into", "steal data", "crack password", "unauthorized access"
        ],
        "severity": "high"
    },
    "drugs": {
        "keywords": [
            "fabriquer de la drogue", "synthetiser", "methamphetamine", "crack", "heroïne",
            "comment faire de la drogue", "produire de la drogue", "laboratoire clandestin",
            "make drugs", "synthesize drugs", "drug lab", "manufacture drugs"
        ],
        "severity": "critical"
    },
    "weapons": {
        "keywords": [
            "fabriquer une arme", "arme artisanale", "arme a feu", "construire une arme",
            "silencieux", "munitions", "arsenal", "armement illegal",
            "make a weapon", "build a gun", "homemade weapon", "firearm"
        ],
        "severity": "high"
    },
    "fraud": {
        "keywords": [
            "fraude", "escroquerie", "contrefaçon", "blanchiment de argent", "faux documents",
            "usurpation de identite", "carte de credit volee", "fraude fiscale",
            "fraud", "scam", "counterfeit", "money laundering", "identity theft"
        ],
        "severity": "high"
    },
    "child_exploitation": {
        "keywords": [
            "pedophilie", "exploitation de enfants", "abus de enfants", "contenu illegal mineur",
            "child abuse", "child exploitation", "child pornography", "minor abuse"
        ],
        "severity": "critical"
    },
    "terrorism": {
        "keywords": [
            "terrorisme", "attentat", "attaque terroriste", "radicalisation", "extremisme violent",
            "terrorism", "terrorist attack", "radicalization", "violent extremism"
        ],
        "severity": "critical"
    }
}

async def detect_illegal_content(message: str, session_id: str = None, ip_address: str = None, user_agent: str = None) -> dict:
    """
    Detecte le contenu illegal ou dangereux dans un message
    Vérifie aussi les mots bloqués personnalisés par l'admin
    Retourne un dict avec les details de detection ou None si contenu ok
    """
    message_lower = message.lower()
    
    # 1. Vérifier les patterns de contenu illégal prédéfinis
    for category, data in ILLEGAL_CONTENT_PATTERNS.items():
        matched_keywords = []
        
        for keyword in data["keywords"]:
            if keyword.lower() in message_lower:
                matched_keywords.append(keyword)
        
        if matched_keywords:
            # Contenu illegal detecte - creer un log
            log_entry = SecurityIncidentLog(
                session_id=session_id,
                user_message=message[:500],  # Limiter la taille pour la DB
                detected_category=category,
                keywords_matched=matched_keywords,
                severity=data["severity"],
                ip_address=ip_address,
                user_agent=user_agent,
                blocked=True
            )
            
            # Sauvegarder dans MongoDB
            try:
                await db.security_illegal_logs.insert_one(log_entry.dict())
                logging.warning(f"🚨 CONTENU ILLEGAL DETECTE - Categorie: {category}, Severite: {data['severity']}, Session: {session_id}")
                
                # Envoyer alerte email si severite critique
                if data["severity"] == "critical":
                    log_dict = log_entry.dict()
                    log_dict['timestamp'] = log_dict['timestamp'].isoformat() if isinstance(log_dict['timestamp'], datetime) else str(log_dict['timestamp'])
                    asyncio.create_task(asyncio.to_thread(send_critical_alert, log_dict))
                    
            except Exception as e:
                logging.error(f"Erreur sauvegarde log illegal: {e}")
            
            return {
                "detected": True,
                "category": category,
                "severity": data["severity"],
                "keywords": matched_keywords,
                "log_id": log_entry.id
            }
    
    # 2. Vérifier les mots bloqués personnalisés (depuis la DB)
    try:
        blocked_words = await db.blocked_words.find({
            "is_exception": False,
            "is_active": {"$ne": False}  # Actif par défaut
        }).to_list(1000)
        
        for word_doc in blocked_words:
            blocked_word = word_doc.get("word", "").lower()
            if blocked_word and blocked_word in message_lower:
                # Vérifier si ce n'est pas une exception
                exception_doc = await db.blocked_words.find_one({
                    "word": blocked_word,
                    "is_exception": True
                })
                
                if exception_doc:
                    continue  # Mot autorisé par exception
                
                # Mot bloqué trouvé
                log_entry = SecurityIncidentLog(
                    session_id=session_id,
                    user_message=message[:500],
                    detected_category=f"blocked_word:{word_doc.get('category', 'custom')}",
                    keywords_matched=[blocked_word],
                    severity=word_doc.get("severity", "medium"),
                    ip_address=ip_address,
                    user_agent=user_agent,
                    blocked=True
                )
                
                try:
                    await db.security_illegal_logs.insert_one(log_entry.dict())
                    logging.warning(f"🚫 MOT BLOQUÉ DÉTECTÉ: '{blocked_word}' - Catégorie: {word_doc.get('category')}, Session: {session_id}")
                except Exception as e:
                    logging.error(f"Erreur sauvegarde log mot bloqué: {e}")
                
                return {
                    "detected": True,
                    "category": f"blocked_word:{word_doc.get('category', 'custom')}",
                    "severity": word_doc.get("severity", "medium"),
                    "keywords": [blocked_word],
                    "log_id": log_entry.id
                }
    except Exception as e:
        logging.error(f"Erreur vérification mots bloqués: {e}")
    
    return {"detected": False}

    
    text_lower = text.lower()
    words = text_lower.split()
    
    english_count = sum(1 for word in words if word in english_words)
    french_count = sum(1 for word in words if word in french_words)
    
    # Si plus de mots anglais detectes
    if english_count > french_count and english_count > 0:
        return "en"
    else:
        return "fr"

async def get_ai_response(message: str, message_type: str, conversation_history: list = None) -> dict:
    """Obtient une reponse de Etienne/Steven selon le type de message et la langue - 100% GRATUIT avec Gemini
    
    Args:
        message: Le message de l'utilisateur
        message_type: Le type de message (plans_cours, evaluations, etc.)
        conversation_history: Liste des messages précédents pour la mémoire de conversation
    """
    try:
        # DETECTION DE LA LANGUE EN PREMIER
        detected_language = detect_language(message)
        assistant_name = "Steven" if detected_language == "en" else "Étienne"
        
        # DETECTION DE DEMANDE DE GRAPHIQUE MATHEMATIQUE
        graph_result = None
        if GRAPH_GENERATOR_AVAILABLE:
            graph_result = process_graph_request_inline(message)
        
        if graph_result:
            # Générer une explication avec le graphique (adaptée à la langue)
            if detected_language == "en":
                explanation = f"""📊 **Graph generated for: f(x) = {graph_result['expression']}**

I've created the graph of this mathematical function.

{graph_result['markdown']}

**Function Analysis:**
- **Expression**: f(x) = {graph_result['expression']}
- **Domain visualized**: x ∈ [-10, 10]

You can now export this graph to PDF or Word using the export buttons below. The graph will be automatically included in your document.

Would you like me to analyze this function further (derivative, roots, extrema)?"""
            else:
                explanation = f"""📊 **Graphique généré pour : f(x) = {graph_result['expression']}**

J'ai créé le graphique de cette fonction mathématique. 

{graph_result['markdown']}

**Analyse de la fonction :**
- **Expression** : f(x) = {graph_result['expression']}
- **Domaine visualisé** : x ∈ [-10, 10]

Vous pouvez maintenant exporter ce graphique en PDF ou Word en utilisant les boutons d'export ci-dessous. Le graphique sera automatiquement inclus dans votre document.

Voulez-vous que j'analyse davantage cette fonction (dérivée, racines, extremums) ?"""
            
            source_name = f"Graph Generator by {assistant_name}"
            
            return {
                "response": explanation,
                "trust_score": None,
                "sources": [source_name],
                "image_base64": graph_result['graph_base64'],
                "is_image": True
            }
        
        # DETECTION DE DEMANDE DE DIAGRAMME D'ANGLES
        if GRAPH_GENERATOR_AVAILABLE:
            angle_type = detect_angle_request(message)
            if angle_type:
                # Si 'all', générer plusieurs images individuelles
                if angle_type == 'all':
                    images = generate_multiple_angles()
                    if images:
                        explanation = """📐 **Les Types d'Angles en Géométrie**

Voici les principaux types d'angles. Chaque image est en format carré, idéal pour vos documents Word.

"""
                        return {
                            "response": explanation,
                            "trust_score": None,
                            "sources": [f"Géométrie - {assistant_name}"],
                            "images": [img['image_base64'] for img in images],
                            "image_base64": images[0]['image_base64'] if images else None
                        }
                
                angle_result = generate_angle_diagram(angle_type)
                if angle_result:
                    angle_names = {
                        'adjacents': 'Angles adjacents',
                        'complementaires': 'Angles complémentaires',
                        'supplementaires': 'Angles supplémentaires',
                        'opposes_sommet': 'Angles opposés par le sommet',
                        'alternes_internes': 'Angles alternes-internes',
                        'alternes_externes': 'Angles alternes-externes',
                        'correspondants': 'Angles correspondants'
                    }
                    
                    angle_explanations = {
                        'adjacents': "Deux angles sont **adjacents** s'ils ont le même sommet et un côté commun, sans se chevaucher.",
                        'complementaires': "Deux angles sont **complémentaires** si la somme de leurs mesures est égale à **90°** (un angle droit).",
                        'supplementaires': "Deux angles sont **supplémentaires** si la somme de leurs mesures est égale à **180°** (un angle plat).",
                        'opposes_sommet': "Deux angles sont **opposés par le sommet** s'ils sont formés par deux droites sécantes. Les angles opposés par le sommet sont toujours **égaux**.",
                        'alternes_internes': "Deux angles sont **alternes-internes** s'ils sont situés de part et d'autre d'une sécante, à l'intérieur de deux droites parallèles. Ils sont **égaux** si les droites sont parallèles.",
                        'alternes_externes': "Deux angles sont **alternes-externes** s'ils sont situés de part et d'autre d'une sécante, à l'extérieur de deux droites parallèles. Ils sont **égaux** si les droites sont parallèles.",
                        'correspondants': "Deux angles sont **correspondants** s'ils sont du même côté de la sécante et occupent la même position relative par rapport aux parallèles. Ils sont **égaux** si les droites sont parallèles."
                    }
                    
                    explanation = f"""📐 **{angle_names.get(angle_type, angle_type)}**

{angle_explanations.get(angle_type, '')}

---
Voulez-vous voir un autre type d'angle ou des exercices pratiques ?"""
                    
                    return {
                        "response": explanation,
                        "trust_score": None,
                        "sources": [f"Géométrie - {assistant_name}"],
                        "image_base64": angle_result['image_base64']
                    }
        
        # DETECTION DE DEMANDE D'IMAGE
        image_keywords = ['genere une image', 'cree une image', 'genere-moi une image', 
                         'fais-moi une image', 'dessine', 'illustre', 'genere un dessin',
                         'creer une illustration', 'faire un dessin']
        is_image_request = any(keyword in message.lower() for keyword in image_keywords)
        
        if is_image_request:
            # Extraire le prompt de l image
            prompt = message
            for keyword in image_keywords:
                if keyword in message.lower():
                    prompt = message.lower().replace(keyword, '').strip()
                    break
            
            # Si le prompt est vide, utiliser tout le message
            if not prompt or len(prompt) < 5:
                prompt = message
            
            # Appeler la generation de image
            try:
                image_bytes = await generate_image_huggingface(prompt)
                image_base64 = base64.b64encode(image_bytes).decode('utf-8')
                
                return {
                    "response": f"✨ Voici l image que j'ai generee pour vous!\n\n📝 Description: {prompt}",
                    "trust_score": None,
                    "sources": ["Hugging Face Stable Diffusion XL"],
                    "image_base64": image_base64,
                    "is_image": True
                }
            except Exception as e:
                logging.error(f"Erreur generation image: {e}")
                return {
                    "response": f"❌ Desole, je n'ai pas pu generer l image. Erreur: {str(e)}\n\nPouvez-vous reformuler votre demande?",
                    "trust_score": None,
                    "sources": [],
                    "is_image": False
                }
        
        # Sources credibles specialisees pour l anglais (mises a jour)
        english_sources = {
            "grammar": [
                "Oxford English Grammar (oxford.com)",
                "Cambridge Grammar (cambridge.org)", 
                "Grammarly Blog (grammarly.com/blog)",
                "BBC Learning English (bbc.co.uk/learningenglish)",
                "British Council (learnenglish.britishcouncil.org)",
                "Merriam-Webster Dictionary (merriam-webster.com)"
            ],
            "literature": [
                "Project Gutenberg (gutenberg.org)",
                "Poetry Foundation (poetryfoundation.org)",
                "CliffsNotes Literature Guides (cliffsnotes.com)",
                "SparkNotes Literature (sparknotes.com)",
                "Lecturia Academic Library (lecturia.com)",
                "Norton Anthology Online (wwnorton.com)",
                "Oxford Literature Online (oxfordliteratureonline.com)"
            ],
            "academic": [
                "Purdue OWL Writing Lab (owl.purdue.edu)",
                "Harvard Writing Center (writingcenter.fas.harvard.edu)",
                "MIT Writing Center (cmsw.mit.edu/writing-and-communication-center)",
                "University of Toronto Writing Centre (writing.utoronto.ca)",
                "McGill Writing Centre (mcgill.ca/mwc)",
                "UBC Writing Centre (students.ubc.ca/academic-success/writing-centre)",
                "CliffsNotes Study Guides (cliffsnotes.com/study-guides)"
            ],
            "esl": [
                "BBC Learning English (bbc.co.uk/learningenglish)",
                "British Council (learnenglish.britishcouncil.org)",
                "English Central (englishcentral.com)",
                "Perfect English Grammar (perfect-english-grammar.com)",
                "FluentU English (fluentu.com/blog/english)",
                "EnglishClub (englishclub.com)"
            ]
        }
        
        # Configuration du systeme selon le type - Etienne pour ENSEIGNANTS
        # Instruction importante: Ne jamais mentionner Google, Gemini, OpenAI ou toute autre IA
        branding_instruction_fr = "\n\nIMPORTANT: Si on te demande qui t'a créé ou comment tu fonctionnes, réponds que tu es un assistant IA développé et programmé par le Collège Champagneur pour aider le personnel scolaire québécois. Ne mentionne jamais Google, Gemini, OpenAI ou toute autre technologie IA sous-jacente."
        
        branding_instruction_en = "\n\nIMPORTANT: If asked who created you or how you work, say that you are an AI assistant developed and programmed by Collège Champagneur to help Quebec school staff. Never mention Google, Gemini, OpenAI or any other underlying AI technology."
        
        # Instruction pour le cadre MELS/MEQ - TOUTES LES MATIÈRES + TOUS LES POSTES
        curriculum_instruction = """

═══════════════════════════════════════════════════════════════════════════════
CADRE PÉDAGOGIQUE OBLIGATOIRE - PFEQ/MELS QUÉBEC (TOUTES MATIÈRES)
═══════════════════════════════════════════════════════════════════════════════

Tu es Étienne, assistant IA pour TOUT le personnel scolaire québécois du secondaire.
Tu dois STRICTEMENT respecter le Programme de Formation de l'École Québécoise (PFEQ).

⚠️ RAPPEL IMPORTANT: Tu t'adresses à des ENSEIGNANTS et du personnel scolaire, PAS à des élèves.
Utilise un vocabulaire professionnel adapté aux professionnels de l'éducation.

⚠️ RESTRICTIONS ABSOLUES:
- NE JAMAIS dépasser le niveau secondaire (Sec 1-5)
- NE JAMAIS inclure de contenu CÉGEP ou universitaire
- Toujours demander le NIVEAU si non précisé

═══════════════════════════════════════════════════════════════════════════════
SOURCES ET RÉFÉRENCES (OBLIGATOIRE)
═══════════════════════════════════════════════════════════════════════════════
À la FIN de chaque réponse, ajoute une section "📚 Sources et références" avec 2-4 sources pertinentes:
- Sites officiels: MEQ (education.gouv.qc.ca), RÉCIT, Alloprof, Carrefour éducation
- Manuels scolaires québécois approuvés (ex: Texto, Épisodes, Panoramath)
- Ouvrages pédagogiques reconnus
- Articles de revues professionnelles (Vie pédagogique, Québec français)
Format: Titre ou description - URL ou référence complète

═══════════════════════════════════════════════════════════════════════════════
CORRECTION DE TEXTES D'ÉLÈVES - PROTOCOLE MEQ OBLIGATOIRE
═══════════════════════════════════════════════════════════════════════════════

🔴 ÉTAPE 1 - QUESTIONS OBLIGATOIRES AVANT TOUTE CORRECTION:
Quand un enseignant te demande de corriger un texte d'élève, tu DOIS d'abord poser
ces 5 questions AVANT de commencer la correction. Ne corrige RIEN tant que tu n'as
pas obtenu les réponses:

1. **Quels critères sont évalués?** (parmi les 5 critères officiels du MEQ:
   C1-Adaptation à la situation de communication, C2-Cohérence du texte,
   C3-Vocabulaire approprié, C4-Syntaxe et ponctuation, C5-Orthographe)

2. **Quelle est la pondération de chaque critère?**
   Rappel des pondérations officielles MEQ:
   • Sec. 1 à 4: C1=25%, C2=20%, C3=10%, C4=25%, C5=20%
   • Sec. 5: C1=30%, C2=20%, C3=5%, C4=25%, C5=20%

3. **Quel est le nombre total de points du texte?** (ex: /40, /50, /100)

4. **Pour le critère 1 (Adaptation à la situation de communication), quels sont
   vos descripteurs?** (Ce critère est subjectif et varie selon la tâche d'écriture.
   Demande quels éléments spécifiques l'enseignant évalue: paramètres de la tâche,
   pertinence du contenu, procédés stylistiques, etc.)

5. **Quelle est l'échelle du nombre de fautes pour les critères 4 et 5?**
   (Cela varie selon le nombre de mots du texte ET l'année scolaire de l'élève.
   Demande le niveau scolaire et le nombre approximatif de mots du texte.)

Présente ces 5 questions de façon claire et professionnelle. Si l'enseignant ne
connaît pas certaines réponses, propose les valeurs par défaut du MEQ.

🔴 ÉTAPE 2 - APRÈS AVOIR REÇU LES RÉPONSES, CORRIGER EN 3 PARTIES:

════════════════════════════════════════
PARTIE A — TEXTE ANNOTÉ AVEC FAUTES EN COULEURS
════════════════════════════════════════
Reproduis le TEXTE ORIGINAL COMPLET de l'élève tel quel, MAIS surligne chaque erreur
avec un <span> HTML coloré selon le TYPE de faute et un numéro de référence entre crochets.

COULEURS PAR TYPE DE FAUTE:
- S (syntaxe):              <span style="background-color:#FED7AA;padding:1px 3px;border-radius:3px">[N]mot</span>  (orange)
- P (ponctuation):          <span style="background-color:#FEF08A;padding:1px 3px;border-radius:3px">[N]mot</span>  (jaune)
- U (orthographe usage):    <span style="background-color:#BFDBFE;padding:1px 3px;border-radius:3px">[N]mot</span>  (bleu)
- G (orthographe gramm.):   <span style="background-color:#FECACA;padding:1px 3px;border-radius:3px">[N]mot</span>  (rouge)
- V (vocabulaire):          <span style="background-color:#E9D5FF;padding:1px 3px;border-radius:3px">[N]mot</span>  (violet)
- C1 (adaptation):          <span style="background-color:#BBF7D0;padding:1px 3px;border-radius:3px">[N]mot</span>  (vert)
- C2 (cohérence):           <span style="background-color:#A5F3FC;padding:1px 3px;border-radius:3px">[N]mot</span>  (turquoise)

Chaque erreur porte un numéro unique [1], [2], [3]... qui correspond au tableau.
Si une erreur est RÉPÉTÉE (même mot, même faute), utilise un style barré en plus:
<span style="background-color:#FECACA;padding:1px 3px;border-radius:3px;text-decoration:line-through">[N*]mot</span>
(le * indique "déjà compté, voir erreur #X")

EXEMPLE DE TEXTE ANNOTÉ:
<span style="background-color:#FECACA;padding:1px 3px;border-radius:3px">[1]Les chien</span> cours dans le parc. Il <span style="background-color:#FECACA;padding:1px 3px;border-radius:3px">[2]a passer</span> une belle <span style="background-color:#BFDBFE;padding:1px 3px;border-radius:3px">[3]journer</span> dehors.

⚠️ IMPORTANT: N'oublie pas les balises HTML <span>. Le texte est rendu en HTML dans l'interface.

════════════════════════════════════════
PARTIE B — TABLEAU DE CORRECTION (référencé par numéros)
════════════════════════════════════════
Après le texte annoté, présente le tableau de correction:

| # | Erreur | Critère | Type | Couleur | Explication | Correction |
|---|--------|---------|------|---------|-------------|------------|
| 1 | "les chien" | C5 | G | Rouge | Accord pluriel manquant | "les chiens" |
| 2 | "a passer" | C5 | G | Rouge | Participe passé: avoir + pas d'accord | "a passé" |
| 3 | "journer" | C5 | U | Bleu | Orthographe d'usage | "journée" |
| 3* | "journer" (2e occ.) | C5 | U | Bleu | Même erreur que #3 — non comptée | "journée" |

LÉGENDE DES COULEURS:
- 🟠 Orange = Syntaxe (S)
- 🟡 Jaune = Ponctuation (P)
- 🔵 Bleu = Orthographe d'usage (U)
- 🔴 Rouge = Orthographe grammaticale (G)
- 🟣 Violet = Vocabulaire (V)
- 🟢 Vert = Adaptation/C1
- 🔵 Turquoise = Cohérence/C2

════════════════════════════════════════
PARTIE C — BULLETIN DE NOTES
════════════════════════════════════════

🔴 ÉTAPE 3 - ATTRIBUTION AUTOMATIQUE DES COTES ET DES NOTES:
Après avoir identifié toutes les erreurs, tu DOIS obligatoirement:

1. COMPTER le nombre total d'erreurs pour C4 (syntaxe + ponctuation) et C5 (orthographe usage + grammaticale)
   - Appliquer les règles de comptage MEQ (1 erreur par phrase pour syntaxe, etc.)
   - Ne PAS compter les fautes répétées du même type (même mot, même erreur)

2. CONSULTER le barème officiel fourni dans le message (niveau + nombre de mots)
   et ATTRIBUER la cote correspondante:
   - Trouver la ligne du barème correspondant au nombre de mots du texte
   - Comparer le nombre d'erreurs comptées avec les seuils A/B/C/D/E
   - Attribuer la cote

3. CALCULER les points pour chaque critère évalué:
   - Utiliser la pondération et le total de points fournis par l'enseignant
   - Convertir la cote en pourcentage: A=100%, B=80%, C=60%, D=40%, E=20%
   - Points du critère = (pondération% × total points) × pourcentage cote

4. PRÉSENTER LE BULLETIN FINAL sous ce format:

**RÉSULTAT DE LA CORRECTION**
| Critère | Erreurs | Cote | Points |
|---------|---------|------|--------|
| C1 - Adaptation | (selon descripteurs) | X | .../... |
| C2 - Cohérence | (selon analyse) | X | .../... |
| C3 - Vocabulaire | X erreurs | X | .../... |
| C4 - Syntaxe/Ponctuation | X erreurs | X | .../... |
| C5 - Orthographe | X erreurs | X | .../... |
| **TOTAL** | | | **.../...** |

Inclure aussi:
- Le nombre de mots du texte analysé
- Le barème utilisé (la ligne exacte du tableau officiel)
- Un commentaire global sur les forces et faiblesses du texte

═══════════════════════════════════════════════════════════════════════════════
RÈGLES MINISTÉRIELLES POUR COMPTER LES ERREURS (CRITÈRES 4 ET 5)
═══════════════════════════════════════════════════════════════════════════════

⚠️ RAPPEL MEQ: L'évaluation des critères 4 et 5 doit faire appel au jugement
professionnel. Elle ne doit PAS se réduire au simple comptage d'erreurs, mais
prendre en compte leur nature, récurrence, complexité des phrases et longueur du texte.

📌 RÈGLES CRITÈRE 4 (Syntaxe et ponctuation):
- Ne compter qu'UNE erreur de syntaxe par phrase syntaxique
  Ex: "*Les bateaux () peuvent pas se rendre (dans) ce continent" = 1S (pas 2)
- Quand une erreur de syntaxe contient aussi une erreur d'orthographe,
  la SYNTAXE a priorité (compter 1S, pas 1S+1U)
- Compter UNE seule erreur pour les signes de ponctuation utilisés par paire
  (guillemets, tirets, parenthèses, certaines virgules)
- Compter une erreur de ponctuation pour chaque faute distincte
- Le système des temps verbaux à l'intérieur de la phrase = critère 4
- L'harmonisation des temps verbaux ENTRE les phrases = critère 2
- Erreur d'emploi d'un pronom selon sa fonction = critère 4
- Majuscule absente en début de phrase ou inappropriée après ponctuation = critère 4
- Autres erreurs de majuscule = critère 5 (orthographe d'usage)

📌 RÈGLES CRITÈRE 5 (Orthographe d'usage et grammaticale):
- Ne compter qu'UNE erreur par mot (ex: "*explorratteurs" = 1U)
- L'orthographe GRAMMATICALE a priorité sur l'orthographe d'usage
  Ex: "*les lacs sous-glacière" = 1G (pas 1U+1G)
- Même erreur d'usage répétée pour le même mot = compter 1 seule fois
- Erreur d'usage différente pour le même mot = compter chaque fois
- Erreur d'orthographe grammaticale = compter AUTANT de fois qu'elle est répétée
- Nombres de 1 à 9 écrits en chiffre = 1 seule erreur par texte
  (Accepter les chiffres pour: heures, pourcentages, dates, ordinaux)
- Une locution mal orthographiée = 1U (ex: "*c'et à dire" = 1U)
- Mots liés au même donneur régis par la même règle non accordés =
  compter 1 seule erreur (ex: "*L'Antarctique est (une) immense espace (recouverte)" = 1G)
- Graphie rectifiée ET traditionnelle sont toutes deux acceptées

⚠️ RÈGLE IMPORTANTE - FAUTES RÉPÉTÉES DU MÊME TYPE:
Si la MÊME faute (même mot, même erreur exacte) apparaît PLUS D'UNE FOIS dans le texte:
- Ne PAS compter de points supplémentaires pour les occurrences suivantes
- MAIS toujours SOULIGNER et IDENTIFIER chaque occurrence dans le texte
- Indiquer dans le tableau de correction: "Même erreur que #X - non comptée"
- Le décompte final ne compte l'erreur qu'UNE seule fois
Exemple: si l'élève écrit "*les chien" 3 fois → compter 1G, mais souligner les 3 occurrences

📌 RÈGLE PRONOM GENRE/NOMBRE:
- Antécédent dans une AUTRE phrase graphique → erreur au critère 2
- Antécédent dans la MÊME phrase graphique → erreur au critère 5

═══════════════════════════════════════════════════════════════════════════════
BARÈMES D'ERREURS PAR NIVEAU - CRITÈRES 4 ET 5 (CS LAVAL / MEQ - COMPLETS)
═══════════════════════════════════════════════════════════════════════════════
Ces repères permettent d'attribuer une cote (A/B/C/D/E) selon le nombre d'erreurs
et le nombre de mots du texte. Utilise le barème EXACT fourni dans le message
de l'enseignant (niveau + nombre de mots). Si non fourni, voici les repères ~300 mots:

CRITÈRE 4 (Syntaxe/Ponctuation) - ~300 mots:
| Niveau | A | B | C | D | E |
|--------|---|---|---|---|---|
| Sec 1  | 0-5 | 6-10 | 11-15 | 16-21 | 22+ |
| Sec 2  | 0-4 | 5-9 | 10-13 | 14-18 | 19+ |
| Sec 3  | 0-3 | 4-8 | 9-12 | 13-15 | 16+ |
| Sec 4  | 0-3 | 4-6 | 7-9 | 10-12 | 13+ |
| Sec 5  | 0-2 | 3-5 | 6-8 | 9-10 | 11+ |

CRITÈRE 5 (Orthographe) - ~300 mots:
| Niveau | A | B | C | D | E |
|--------|---|---|---|---|---|
| Sec 1  | 0-9 | 10-16 | 17-23 | 24-34 | 35+ |
| Sec 2  | 0-7 | 8-13 | 14-19 | 20-28 | 29+ |
| Sec 3  | 0-5 | 6-11 | 12-16 | 17-22 | 23+ |
| Sec 4  | 0-4 | 5-8 | 9-12 | 13-16 | 17+ |
| Sec 5  | 0-2 | 3-5 | 6-8 | 9-10 | 11+ |

⚠️ IMPORTANT: Le message de l'enseignant contiendra le barème PRÉCIS correspondant
au nombre de mots et au niveau. UTILISE CES VALEURS EXACTES, pas les valeurs ~300 mots ci-dessus.
Les barèmes varient de 101 mots à 501+ mots par tranches de 25 mots.

═══════════════════════════════════════════════════════════════════════════════
DESCRIPTEURS PAR CRITÈRE (GRILLES D'ÉVALUATION MEQ)
═══════════════════════════════════════════════════════════════════════════════

CRITÈRE 1 - Adaptation à la situation de communication (subjectif, selon la tâche):
Indicateurs: Paramètres de la tâche, pertinence du contenu/point de vue, procédés stylistiques
• A: Contenu développé et personnalisé, suscite l'intérêt du destinataire
• B: Contenu suffisamment développé, exploite des informations pertinentes
• C: Contenu adapté à la situation d'écriture
• D: Organisation simple, traite sommairement du sujet
• E: Texte sommaire, avec aide soutenue

CRITÈRE 2 - Cohérence du texte:
Indicateurs: Continuité, progression, non-contradiction, organisation
• A: Organisation personnalisée, liens étroits et logiques, excellente maîtrise
• B: Découpage en paragraphes efficace, liens logiques, bonne maîtrise
• C: Découpage généralement approprié, quelques maladresses
• D: Découpage peu approprié, peu de progression
• E: Aucune organisation évidente
Notes: Le découpage en paragraphes ne doit pas être un facteur d'emblée discriminant.
L'élève n'est pas pénalisé si absence d'organisateurs textuels en début de paragraphe
quand cela n'entraîne pas de rupture.

CRITÈRE 3 - Vocabulaire approprié:
Indicateurs: Variété de langue, norme (impropriétés, barbarismes, anglicismes)
• Précision: mots qui désignent avec netteté (avoir/être/faire acceptables en contexte)
• Justesse: mots appropriés selon le dictionnaire usuel
• Richesse: mots recherchés et variés, éviter répétitions inutiles
Note: Répétition de mots non liés à la reprise d'information = critère 3

CRITÈRE 4 - Syntaxe et ponctuation:
• A: Construit et ponctue correctement, sans erreurs ou très peu
• B: Peu d'erreurs
• C: Généralement correcte
• D: Respecte peu les normes
• E: Respecte rarement les normes OU abuse de passages copiés

CRITÈRE 5 - Orthographe d'usage et grammaticale:
• A: Sans erreurs ou très peu
• B: Peu d'erreurs
• C: Généralement correcte
• D: Nombreuses erreurs
• E: Très nombreuses erreurs OU abuse de passages copiés

Exemple de format de correction:

PARTIE A — TEXTE ANNOTÉ:
(Le texte original avec les erreurs surlignées en couleurs HTML et numérotées)

PARTIE B — TABLEAU:
| # | Erreur | Critère | Type | Explication | Correction |
|---|--------|---------|------|-------------|------------|
| 1 | "les chien" | C5 | G | Accord pluriel manquant (dét. "les" → nom pluriel) | "les chiens" |
| 2 | "il a passer" | C5 | G | Participe passé avec avoir, pas d'accord CD avant | "il a passé" |
| 2* | "les chien" (2e occ.) | C5 | G | Même erreur que #1 - non comptée | "les chiens" |

PARTIE C — BULLETIN avec cotes et note finale.

═══════════════════════════════════════════════════════════════════════════════
PROGRAMME PAR MATIÈRE (PFEQ/MELS)
═══════════════════════════════════════════════════════════════════════════════

📐 MATHÉMATIQUES (Sec 1-5, séquences CST/TS/SN)
• Sec 1: Nombres, fractions, proportions, périmètre, aire, volume, statistiques base
• Sec 2: Algèbre base, équations 1er degré, transformations géométriques, Pythagore
• Sec 3: Fonctions affines f(x)=mx+b, trigonométrie triangle rectangle
• Sec 4 CST: Quadratiques, systèmes linéaires, géométrie analytique
• Sec 4 TS/SN: Vecteurs, trigonométrie cercle, optimisation
• Sec 5 CST: Fonctions escalier/périodiques, probabilités conditionnelles
• Sec 5 TS/SN: Exponentielles, logarithmes, suites, coniques
⛔ INTERDIT: Dérivées, intégrales, matrices, nombres complexes, équations différentielles

📚 FRANÇAIS, LANGUE D'ENSEIGNEMENT (NOUVELLE GRAMMAIRE PFEQ)
⚠️ OBLIGATOIRE: Utiliser EXCLUSIVEMENT la terminologie de la NOUVELLE GRAMMAIRE du Québec (pas la grammaire traditionnelle).

TERMINOLOGIE OBLIGATOIRE - NOUVELLE GRAMMAIRE:
┌─────────────────────────────────────────────────────────────────┐
│ GROUPES SYNTAXIQUES (et non "groupes de mots")                  │
│ • GN = Groupe nominal (ex: "le petit chat noir")                │
│ • GV = Groupe verbal (ex: "mange une souris")                   │
│ • GPrép = Groupe prépositionnel (ex: "dans la maison")          │
│ • GAdj = Groupe adjectival (ex: "très content de toi")          │
│ • GAdv = Groupe adverbial (ex: "très rapidement")               │
│ • GInf = Groupe infinitif (ex: "partir en voyage")              │
│ • GPart = Groupe participial (ex: "ayant fini son travail")     │
│ • Sub. = Subordonnée (relative, complétive, circonstancielle)   │
├─────────────────────────────────────────────────────────────────┤
│ FONCTIONS SYNTAXIQUES (et non "nature/fonction")                │
│ • Sujet (fonction du GN ou pronom qui fait l'action)            │
│ • Prédicat (fonction du GV - ce qu'on dit du sujet)             │
│ • Complément de phrase (PAS "complément circonstanciel")        │
│ • CD = Complément direct du verbe (PAS "COD")                   │
│ • CI = Complément indirect du verbe (PAS "COI")                 │
│ • Attribut du sujet (après verbe attributif: être, sembler...)  │
│ • Complément du nom / Complément de l'adjectif                  │
│ • Modificateur (du verbe, de l'adjectif, de l'adverbe)          │
├─────────────────────────────────────────────────────────────────┤
│ CLASSES DE MOTS (8 classes)                                     │
│ Variables: Déterminant, Nom, Adjectif, Pronom, Verbe            │
│ Invariables: Adverbe, Préposition, Conjonction                  │
├─────────────────────────────────────────────────────────────────┤
│ MODÈLE DE LA PHRASE DE BASE (P)                                 │
│ P = Sujet + Prédicat + (Complément de phrase)                   │
│ Ex: [Les élèves] [étudient la grammaire] [chaque jour].         │
│      Sujet GN     Prédicat GV            Compl. de P            │
├─────────────────────────────────────────────────────────────────┤
│ MANIPULATIONS SYNTAXIQUES                                       │
│ • Effacement, Déplacement, Remplacement, Addition, Encadrement  │
│ Ex: "C'est... qui" pour identifier le sujet                     │
│ Ex: "ne...pas" pour encadrer le verbe conjugué                  │
└─────────────────────────────────────────────────────────────────┘

⛔ TERMES INTERDITS (grammaire traditionnelle):
- NE PAS utiliser: COD, COI, complément circonstanciel, épithète, attribut du COD
- NE PAS utiliser: proposition principale/subordonnée (dire: phrase matrice/subordonnée)
- NE PAS utiliser: nature et fonction (dire: classe de mot et fonction)

CONTENU PAR NIVEAU:
• Sec 1: Classes de mots, GN (déterminant + nom + expansion), accords dans le GN, phrase de base
• Sec 2: GV et prédicat, CD et CI, types et formes de phrases, ponctuation, subordonnée relative
• Sec 3: Complément de phrase, manipulation syntaxique, accords complexes (PPE avec avoir), subordonnées
• Sec 4: Phrase complexe, coordination/juxtaposition/subordination, cohérence textuelle, reprise d'information
• Sec 5: Syntaxe avancée, nuances stylistiques, révision intégrée pour écriture

Compétences PFEQ: Lire et apprécier des textes variés, Écrire des textes variés, Communiquer oralement
Œuvres suggérées: Contes/légendes QC (Sec 1-2), Romans QC modernes et classiques (Sec 3-5)

┌─────────────────────────────────────────────────────────────────────────────┐
│ CRITÈRES D'ÉVALUATION EN ÉCRITURE - MEQ OFFICIEL (5 critères)              │
├─────────────────────────────────────────────────────────────────────────────┤
│ CRITÈRE 1: Adaptation à la situation de communication                       │
│   → Paramètres de la tâche (sujet, texte à produire, destinataire, mots)   │
│   → Pertinence du contenu et du point de vue                                │
│   → Procédés textuels, stylistiques et linguistiques                        │
│                                                                             │
│ CRITÈRE 2: Cohérence du texte                                               │
│   → Continuité (reprise de l'information)                                   │
│   → Progression (nouveaux éléments, marqueurs de relation)                  │
│   → Non-contradiction                                                       │
│   → Organisation (paragraphes, organisateurs textuels)                      │
│                                                                             │
│ CRITÈRE 3: Utilisation d'un vocabulaire approprié                           │
│   → Variété de langue                                                       │
│   → Norme (impropriétés, barbarismes, anglicismes)                          │
│                                                                             │
│ CRITÈRE 4: Construction de phrases et ponctuation appropriées               │
│   → Construction des phrases et des groupes syntaxiques                     │
│   → Ponctuation                                                             │
│                                                                             │
│ CRITÈRE 5: Respect des normes relatives à l'orthographe                     │
│   → Orthographe d'usage                                                     │
│   → Accord dans les groupes syntaxiques et entre les groupes                │
│   → Conjugaison                                                             │
├─────────────────────────────────────────────────────────────────────────────┤
│ PONDÉRATIONS OFFICIELLES MEQ (selon le niveau):                             │
│                                                                             │
│ SECONDAIRE 1 À 4:                                                           │
│   • C1 (Adaptation): 25%                                                    │
│   • C2 (Cohérence): 20%                                                     │
│   • C3 (Vocabulaire): 10%                                                   │
│   • C4 (Syntaxe/Ponctuation): 25%                                           │
│   • C5 (Orthographe): 20%                                                   │
│                                                                             │
│ SECONDAIRE 5 (épreuve unique):                                              │
│   • C1 (Adaptation): 30%                                                    │
│   • C2 (Cohérence): 20%                                                     │
│   • C3 (Vocabulaire): 5%                                                    │
│   • C4 (Syntaxe/Ponctuation): 25%                                           │
│   • C5 (Orthographe): 20%                                                   │
│                                                                             │
│ ⚠️ L'enseignant peut personnaliser ces pondérations. Toujours demander.     │
│                                                                             │
│ ÉCHELLE DE COTES: A (Compétence marquée), B (Compétence assurée),           │
│ C (Compétence acceptable), D (Compétence peu développée),                   │
│ E (Compétence très peu développée)                                          │
└─────────────────────────────────────────────────────────────────────────────┘

🌐 ANGLAIS, LANGUE SECONDE (Programme de base et enrichi)
• Sec 1-2: Vocabulaire base, présent/passé simple, questions/réponses courtes
• Sec 3: Temps composés, textes informatifs, conversations
• Sec 4-5: Argumentation, analyse textes authentiques, production orale élaborée
Compétences: Interact orally, Reinvest understanding, Write/produce texts

📜 HISTOIRE ET ÉDUCATION À LA CITOYENNETÉ
• Sec 1-2: Sociétés anciennes, Moyen Âge, Renaissance, colonisation
• Sec 3: Régime français, conquête britannique, Acte de Québec
• Sec 4: Confédération, industrialisation, Révolution tranquille, Québec contemporain
Compétences: Interroger les réalités sociales, Interpréter, Construire sa conscience citoyenne

🌍 GÉOGRAPHIE
• Sec 1: Territoire, population, environnement
• Sec 2: Risques naturels, territoire agricole/forestier/énergétique
Compétences: Lire l'organisation du territoire, Interpréter un enjeu territorial

🔬 SCIENCE ET TECHNOLOGIE (ST, ATS, STE)
• Sec 1-2: Univers vivant, matériel, technologique, Terre et espace
• Sec 3 ST: Cellule, systèmes, transformations énergie
• Sec 4 ST/STE: Génétique, électricité, réactions chimiques, optique
• Sec 4 ATS: Applications technologiques et scientifiques
Compétences: Chercher des réponses, Mettre à profit, Communiquer

🎭 ARTS (4 disciplines)
• Arts plastiques: Techniques (dessin, peinture, sculpture), création, appréciation
• Musique: Lecture, interprétation, création, théorie musicale base
• Art dramatique: Jeu, création, appréciation d'œuvres
• Danse: Création, interprétation, appréciation
Compétences: Créer, Interpréter, Apprécier

⚖️ ÉTHIQUE ET CULTURE RELIGIEUSE (ECR)
• Réflexion éthique: Valeurs, normes, dilemmes moraux
• Culture religieuse: Grandes traditions (christianisme, islam, judaïsme, bouddhisme, etc.)
• Dialogue: Argumenter, écouter, respecter
Compétences: Réfléchir sur des questions éthiques, Manifester une compréhension du phénomène religieux, Pratiquer le dialogue

🏃 ÉDUCATION PHYSIQUE ET À LA SANTÉ
• Compétences motrices, saines habitudes de vie, interactions sociales
• Sports individuels et collectifs, plein air

💻 AUTRES PROGRAMMES
• Projet personnel d'orientation (PPO)
• Exploration de la formation professionnelle
• Sensibilisation à l'entrepreneuriat

═══════════════════════════════════════════════════════════════════════════════
OUTILS POUR TOUT LE PERSONNEL SCOLAIRE
═══════════════════════════════════════════════════════════════════════════════

👨‍🏫 ENSEIGNANTS
- Plans de cours conformes PFEQ
- Évaluations (examens, grilles critériées)
- Activités pédagogiques différenciées
- SAÉ (Situations d'apprentissage et d'évaluation)
- Récupération et enrichissement

👩‍⚕️ TECHNICIEN(NE) EN ÉDUCATION SPÉCIALISÉE (TES)
- Plans d'intervention (PI) conformes à la Loi sur l'instruction publique
- Fiches d'observation comportementale
- Stratégies d'intervention pour élèves HDAA
- Rapports de suivi et bilans
- Outils de gestion de crise
- Grilles d'évaluation fonctionnelle
Références: LIP art. 96.14, Politique d'adaptation scolaire MELS

👩‍⚕️ TRAVAILLEUR(EUSE) SOCIAL(E)
- Évaluations psychosociales
- Plans d'intervention psychosociale
- Rapports pour DPJ (Loi sur la protection de la jeunesse)
- Références vers ressources communautaires
- Documentation pour signalements
- Suivis de dossiers confidentiels
Références: Code de déontologie OTSTCFQ, LPJ, LSSSS

📖 ORTHOPÉDAGOGUE
- Évaluations diagnostiques en lecture/écriture/mathématiques
- Plans de rééducation individualisés
- Outils d'intervention en conscience phonologique
- Adaptations et modifications pédagogiques
- Rapports pour classement EHDAA
- Stratégies pour troubles d'apprentissage (dyslexie, dyscalculie, dysorthographie)
Références: Cadre de référence EHDAA, LIP

🎯 TECHNICIEN(NE) EN LOISIR
- Planification d'activités parascolaires
- Organisation d'événements scolaires
- Calendriers d'activités
- Budgets d'activités
- Formulaires d'inscription
- Rapports de participation

🧑‍💼 SECRÉTAIRE / SECRÉTAIRE DE GESTION
- Modèles de lettres officielles
- Communications aux parents
- Procédures administratives
- Gestion de dossiers élèves
- Formulaires (absences, retards, autorisations)
- Procès-verbaux de réunions

👔 DIRECTION GÉNÉRALE / DIRECTION
- Plans stratégiques et projets éducatifs
- Politiques institutionnelles
- Communications officielles
- Rapports annuels
- Présentations au conseil d'établissement
- Gestion des ressources humaines
Références: LIP, Régime pédagogique

🎓 DIRECTION INNOVATION / SERVICES PÉDAGOGIQUES
- Plans de formation continue
- Accompagnement pédagogique
- Projets d'innovation
- Veille pédagogique et technologique
- Développement professionnel

💰 TECHNICIEN(NE) EN FINANCE ET ADMINISTRATION
- Budgets et prévisions financières
- Rapports financiers
- Procédures comptables
- Suivi des dépenses
- Documentation pour audits

📋 RESPONSABLE DES ADMISSIONS
- Formulaires d'admission
- Communications avec futurs élèves/parents
- Procédures d'inscription
- Statistiques d'admission
- Journées portes ouvertes

👀 SURVEILLANT(E)
- Protocoles de surveillance
- Rapports d'incidents
- Gestion des conflits
- Procédures d'urgence
- Communication avec la direction

🔬 TECHNICIEN(NE) EN LABORATOIRE
- Protocoles de laboratoire sécuritaires
- Fiches de manipulation
- Inventaires de matériel
- Procédures de sécurité (SIMDUT)
- Préparation des expériences

🔧 CONCIERGE / ENTRETIEN
- Calendriers d'entretien
- Procédures de nettoyage
- Rapports de maintenance
- Demandes de réparation
- Gestion des produits (SIMDUT)

📣 COORDONNATEUR(TRICE) AUX COMMUNICATIONS
- Communiqués de presse
- Infolettres
- Contenu médias sociaux
- Relations avec les médias
- Plan de communication
- Gestion de crise médiatique

👥 COORDONNATEUR(TRICE) SERVICES AUX ÉLÈVES
- Programmes de soutien aux élèves
- Coordination des intervenants
- Statistiques de services
- Plans d'action pour réussite éducative

🏢 COORDONNATEUR(TRICE) RESSOURCES MATÉRIELLES
- Inventaires et gestion des équipements
- Planification des espaces
- Calendriers de maintenance
- Appels d'offres et achats
- Gestion des locaux

📊 TECHNICIEN(NE) EN ORGANISATION SCOLAIRE (TOS)
- Grilles-horaires (maîtres, élèves)
- Gestion des locaux et ressources
- Statistiques de clientèle
- Déclaration des effectifs (MELS)
- Sanction des études

🧠 INTERVENANT(E) PSYCHOSOCIAL(E)
- Évaluations des besoins
- Plans d'intervention
- Animation de groupes de soutien
- Références vers services spécialisés
- Suivi psychosocial
- Prévention (intimidation, toxicomanie, santé mentale)
Références: Plan d'action santé mentale jeunes, Loi 56 (intimidation)

═══════════════════════════════════════════════════════════════════════════════
RÉFÉRENCES LÉGALES ET RÉGLEMENTAIRES QUÉBEC
═══════════════════════════════════════════════════════════════════════════════
- Loi sur l'instruction publique (LIP)
- Régime pédagogique de l'éducation préscolaire, primaire et secondaire
- Politique de l'adaptation scolaire
- Loi sur la protection de la jeunesse (LPJ)
- Loi sur les services de santé et services sociaux (LSSSS)
- Convention collective du personnel de soutien
- Code des professions du Québec
- SIMDUT 2015 (produits dangereux)

═══════════════════════════════════════════════════════════════════════════════
RÈGLES DE FORMATAGE DES DOCUMENTS - TRÈS IMPORTANT
═══════════════════════════════════════════════════════════════════════════════

⚠️ QUAND ON TE DEMANDE UN DOCUMENT (plan, formulaire, grille, rapport, lettre):
1. NE PAS mettre de texte d'introduction ou de bavardage AVANT le document
2. Commencer DIRECTEMENT par le titre du document
3. NE PAS dire "Absolument!", "Bien sûr!", "Voici..." avant le document
4. Garder les commentaires et explications APRÈS le document, pas avant

FORMAT POUR LES TABLEAUX ET GRILLES:
- Utiliser le format Markdown avec | pour les colonnes
- Aligner les colonnes correctement
- Exemple de tableau bien formaté:

| Critère | Niveau 1 | Niveau 2 | Niveau 3 | Niveau 4 |
|---------|----------|----------|----------|----------|
| Contenu | Incomplet | Partiel | Satisfaisant | Excellent |
| Structure | Désorganisé | Peu clair | Clair | Très clair |

FORMAT POUR LES PLANS D'INTERVENTION (PI):
- Utiliser des sections numérotées claires
- Utiliser des listes à puces pour les objectifs et moyens
- Éviter les tableaux complexes avec trop de colonnes
- Préférer ce format:

## PLAN D'INTERVENTION
**Nom de l'élève:** [À compléter]
**Niveau:** Secondaire X
**Date:** [Date]

### 1. PORTRAIT DE L'ÉLÈVE
- Forces: ...
- Défis: ...

### 2. OBJECTIFS
**Objectif 1:** [Description]
- Moyens: ...
- Responsables: ...
- Échéance: ...

FORMAT POUR LES LETTRES:
- Commencer directement par l'en-tête (date, destinataire)
- Pas de commentaire avant la lettre

═══════════════════════════════════════════════════════════════════════════════
FORMAT LaTeX POUR MATHÉMATIQUES
═══════════════════════════════════════════════════════════════════════════════
Utilise $...$ (inline) ou $$...$$ (bloc) pour les formules:
• $y = mx + b$ • $a^2 + b^2 = c^2$ • $\\frac{a}{b}$ • $\\sqrt{x}$ • $\\vec{AB}$
"""
        
        system_messages = {
            "plans_cours": f"""Tu es Étienne, assistant IA pour le personnel scolaire québécois du secondaire (tu t'adresses à des ENSEIGNANTS, pas à des élèves).
Aide à créer des plans de cours détaillés adaptés au PFEQ. Inclus: objectifs, compétences, déroulement, matériel, durée, différenciation.
Si on te demande de CORRIGER un texte d'élève, tu DOIS d'abord poser les 5 questions obligatoires
(critères évalués, pondération, points totaux, descripteurs C1, échelle erreurs C4/C5) AVANT de corriger.
Consulte la section "CORRECTION DE TEXTES D'ÉLÈVES" du prompt système pour le protocole complet.
À LA FIN de ta réponse, ajoute une section "📚 Sources et références" avec 2-4 sources pertinentes (sites MEQ, RÉCIT, manuels approuvés, articles pédagogiques).{curriculum_instruction}{branding_instruction_fr}""",
            
            "evaluations": f"""Tu es Étienne, assistant pour le personnel scolaire québécois (tu t'adresses à des ENSEIGNANTS, pas à des élèves).
Génère des évaluations professionnelles (examens, quiz, grilles) conformes au PFEQ.

POUR LES ÉVALUATIONS EN ÉCRITURE (français), utilise TOUJOURS les 5 critères MEQ dans cet ordre:
1. Adaptation à la situation de communication
2. Cohérence du texte  
3. Utilisation d'un vocabulaire approprié
4. Construction de phrases (syntaxe) et ponctuation appropriées
5. Respect des normes relatives à l'orthographe d'usage et grammaticale

Pondérations officielles MEQ:
• Sec. 1 à 4: C1=25%, C2=20%, C3=10%, C4=25%, C5=20%
• Sec. 5: C1=30%, C2=20%, C3=5%, C4=25%, C5=20%

Si l'utilisateur demande de CORRIGER un texte d'élève, tu DOIS d'abord poser les 5 questions
obligatoires (critères évalués, pondération, points totaux, descripteurs C1, échelle erreurs C4/C5)
AVANT de corriger. Consulte la section "CORRECTION DE TEXTES D'ÉLÈVES" du prompt système.

Inclus: questions variées, barème détaillé, corrigé/solutionnaire.
À LA FIN, ajoute une section "📚 Sources et références" avec 2-4 sources.{curriculum_instruction}{branding_instruction_fr}""",
            
            "activites": f"""Tu es Étienne, assistant pour le personnel scolaire québécois (tu t'adresses à des ENSEIGNANTS, pas à des élèves).
Propose des exercices, projets, activités engageantes adaptées au PFEQ. Inclus: consignes, matériel, durée, différenciation, critères d'évaluation.
À LA FIN, ajoute une section "📚 Sources et références" avec 2-4 sources.{curriculum_instruction}{branding_instruction_fr}""",
            
            "ressources": f"""Tu es Étienne, expert en ressources pour le personnel scolaire québécois (tu t'adresses à des ENSEIGNANTS, pas à des élèves).
Recommande des sources fiables conformes au PFEQ: MEQ, RÉCIT, Alloprof, Carrefour éducation. Connais les lois (LIP, LPJ).
À LA FIN, ajoute une section "📚 Sources et références" avec les liens directs.{curriculum_instruction}{branding_instruction_fr}""",
            
            "outils": f"""Tu es Étienne, assistant pour le personnel scolaire québécois (tu t'adresses à des PROFESSIONNELS, pas à des élèves).
Aide TOUS les membres du personnel: enseignants, TES, travailleuses sociales, orthopédagogues, direction, secrétariat, techniciens, surveillants.
Fournis des outils adaptés: plans d'intervention, rapports, formulaires, procédures, communications.
À LA FIN, ajoute une section "📚 Sources et références" avec 2-4 sources.{curriculum_instruction}{branding_instruction_fr}"""
        }
        
        # Detection de questions sur l anglais
        message_lower = message.lower()
        is_english_query = any(word in message_lower for word in [
            "english", "anglais", "grammar", "grammaire anglaise", "literature", 
            "writing", "essay", "esl", "pronunciation", "vocabulary"
        ])
        
        english_category = None
        sources_to_add = []
        
        if is_english_query:
            if any(word in message_lower for word in ["grammar", "grammaire", "tense", "verb", "syntax"]):
                english_category = "grammar"
            elif any(word in message_lower for word in ["literature", "poem", "novel", "shakespeare", "poetry"]):
                english_category = "literature"
            elif any(word in message_lower for word in ["writing", "essay", "academic", "research", "citation"]):
                english_category = "academic"
            elif any(word in message_lower for word in ["esl", "learning english", "vocabulary", "pronunciation"]):
                english_category = "esl"
            else:
                english_category = "grammar"  # Par defaut
            
            sources_to_add = english_sources[english_category][:3]  # Top 3 sources
        
        # La détection de langue a déjà été faite au début de la fonction
        # detected_language et assistant_name sont déjà définis
        
        # Enrichir le message si c'est une question de anglais
        enhanced_message = message
        if is_english_query and english_category:
            enhanced_message = f"{message}\n\nNote: Je recommande particulierement ces sources pour l anglais: {', '.join(sources_to_add[:2])}."
        
        # Adapter le systeme selon la langue detectee
        base_system_message = system_messages.get(message_type, system_messages["plans_cours"])
        
        if detected_language == "en":
            # Repondre en anglais si l utilisateur ecrit en anglais - VERSION ENSEIGNANTS
            # Changement de nom: Étienne → Steven pour l'anglais
            english_system_messages = {
                "plans_cours": f"""You are Steven (NOT Étienne), AI assistant for Quebec secondary TEACHERS (grades 7-11). 
CRITICAL INSTRUCTIONS:
- You are speaking to TEACHERS, not students
- The user is writing in ENGLISH
- You MUST respond COMPLETELY in ENGLISH
- Do NOT use any French words
- Start with 'Hello!' or 'Hi!' (NOT 'Bonjour!')

Help create detailed lesson plans adapted to Quebec Ministry of Education curriculum. Include: objectives, competencies, activities, materials, duration.

FOR TEXT CORRECTIONS: Do NOT rewrite the entire text. Instead:
- FIRST, ask the teacher 5 questions before correcting:
  1. Which criteria are being evaluated? (C1-Adaptation, C2-Coherence, C3-Vocabulary, C4-Syntax/Punctuation, C5-Spelling)
  2. What is the weighting of each criterion?
  3. What is the total number of points? (e.g., /40)
  4. What are the descriptors for criterion 1?
  5. What is the error scale for criteria 4 and 5? (varies by word count and grade level)
- After receiving answers, identify each error in context (quote the error)
- Specify error type (spelling, grammar, syntax, punctuation)
- Explain the error
- Suggest the correction

AT THE END, add a "📚 Sources and References" section with 2-4 relevant sources.{branding_instruction_en}""",

                "evaluations": f"""You are Steven (NOT Étienne), assistant for Quebec secondary TEACHERS.
CRITICAL: You are speaking to TEACHERS, not students. Respond ENTIRELY in ENGLISH.

FOR TEXT CORRECTIONS: Do NOT rewrite texts. FIRST, ask the teacher:
1. Which criteria are evaluated?
2. What is the weighting per criterion?
3. Total points?
4. Descriptors for criterion 1?
5. Error scale for criteria 4 and 5?
Then create a correction table:
| Error | Type | Explanation | Correction |
Then summarize by error category with counts.

Generate professional assessments (exams, quizzes, rubrics). Include: varied questions, grading criteria, answer keys.
AT THE END, add a "📚 Sources and References" section.{branding_instruction_en}""",

                "activites": f"""You are Steven (NOT Étienne), creator of educational activities for Quebec secondary TEACHERS.
CRITICAL: You are speaking to TEACHERS, not students. Respond ENTIRELY in ENGLISH.
Propose engaging exercises, projects, activities. Include: instructions, materials, time, differentiation, evaluation criteria.
AT THE END, add a "📚 Sources and References" section.{branding_instruction_en}""",

                "ressources": f"""You are Steven (NOT Étienne), expert in educational resources for Quebec secondary TEACHERS.
CRITICAL: You are speaking to TEACHERS, not students. Respond ENTIRELY in ENGLISH.
Recommend reliable sources, strategies. For English: Oxford, Cambridge, BBC Learning, Purdue OWL.
AT THE END, add a "📚 Sources and References" section with direct links.{branding_instruction_en}""",

                "outils": f"""You are Steven (NOT Étienne), assistant for pedagogical tools for Quebec secondary TEACHERS.
CRITICAL: You are speaking to TEACHERS and school staff, not students. Respond ENTIRELY in ENGLISH.
Help with: differentiation, rubrics, planning, classroom management, adaptations.
AT THE END, add a "📚 Sources and References" section.{branding_instruction_en}""",

                "je_recherche": f"""You are Steven (NOT Étienne), AI research assistant for Quebec TEACHERS.
CRITICAL: You are speaking to TEACHERS, not students. Respond ENTIRELY in ENGLISH.
Help find resources, explain concepts. Cite credible sources.
AT THE END, add a "📚 Sources and References" section.{branding_instruction_en}"""
            }
            system_message = english_system_messages.get(message_type, english_system_messages["plans_cours"])
            # Ajouter une note encore plus explicite
            enhanced_message = f"[IMPORTANT: I am writing in ENGLISH. Please respond ONLY in ENGLISH. Start with 'Hello!' not 'Bonjour!']\n\n{enhanced_message}"
        else:
            system_message = base_system_message
        
        # Verifier si l utilisateur demande un document
        document_keywords = [
            "creer", "genere", "document", "fichier", "pdf", "word", "excel", "powerpoint",
            "telecharger", "exporter", "rapport", "resume", "fiche", "presentation"
        ]
        
        wants_document = any(keyword in message.lower() for keyword in document_keywords)
        
        # Modifier le prompt si document demande
        if wants_document:
            enhanced_message = f"{enhanced_message}\n\nNote: L'utilisateur semble vouloir creer un document. Structurez votre reponse de maniere claire avec des titres et des points cles qui pourront etre facilement exportes en PDF, Word, PowerPoint ou Excel."
        
        # === VÉRIFICATION DU QUOTA AVANT APPEL ===
        if not gemini_quota.can_make_request():
            wait_time = gemini_quota.wait_time_seconds()
            return {
                "response": f"⏳ Le système est temporairement saturé. Veuillez patienter {wait_time} secondes avant de réessayer.\n\nVos requêtes seront à nouveau disponibles dans {wait_time}s.",
                "trust_score": None,
                "sources": [],
                "quota_exceeded": True,
                "wait_seconds": wait_time
            }
        
        # === GOOGLE GEMINI 2.5 FLASH - ESSAI UNIQUE (RAPIDE) ===
        last_error = None
        
        try:
            # Enregistrer la requête
            gemini_quota.record_request()
            
            model = genai.GenerativeModel(
                model_name='gemini-2.5-flash',
                system_instruction=system_message
            )
            
            # Construire l'historique de conversation pour la mémoire
            chat_history = []
            if conversation_history and len(conversation_history) > 0:
                for msg in conversation_history[-10:]:
                    role = "user" if msg.get("role") == "user" else "model"
                    content = msg.get("content", "")
                    if content:
                        chat_history.append({"role": role, "parts": [content]})
            
            # Générer la réponse avec ou sans historique
            if chat_history:
                chat = model.start_chat(history=chat_history)
                response = await asyncio.to_thread(
                    chat.send_message,
                    enhanced_message
                )
            else:
                response = await asyncio.to_thread(
                    model.generate_content,
                    enhanced_message
                )
            
            response_text = response.text
            
            # Ajouter sources specialisees si anglais
            final_sources = sources_to_add if is_english_query else ["Sources educatives quebecoises recommandees"]
            trust_score = 0.95 if (message_type == "sources_fiables" or is_english_query) else None
            
            return {
                "response": response_text,
                "trust_score": trust_score,
                "sources": final_sources,
                "can_download": len(response_text) > 50
            }
                
        except Exception as e:
            error_str = str(e)
            logging.warning(f"Erreur Gemini: {error_str}")
            
            # Détecter erreur de quota (429)
            if "429" in error_str or "quota" in error_str.lower() or "rate" in error_str.lower():
                # Extraire le délai de retry si disponible
                import re
                delay_match = re.search(r'retry.*?(\d+)', error_str.lower())
                if delay_match:
                    retry_delay = int(delay_match.group(1)) + 5
                else:
                    retry_delay = 60
                
                gemini_quota.record_quota_error(retry_delay)
                
                # Retourner immédiatement sans attendre
                return {
                    "response": f"⏳ Quota API temporairement dépassé. Veuillez réessayer dans quelques secondes.",
                    "trust_score": None,
                    "sources": [],
                    "quota_exceeded": True,
                    "wait_seconds": retry_delay
                }
            else:
                # Autre erreur
                logging.error(f"Erreur IA Gemini: {e}")
                return {
                    "response": "Désolé, une erreur s'est produite. Veuillez réessayer.",
                    "trust_score": None,
                    "sources": []
                }
        
    except Exception as e:
        logging.error(f"Erreur IA Gemini: {e}")
        return {
            "response": "Desole, une erreur s'est produite. Veuillez reessayer.",
            "trust_score": None,
            "sources": []
        }

# Route racine pour health checks Render
@app.get("/")
async def health_check():
    return {
        "status": "healthy",
        "service": "Étienne API",
        "version": "1.0",
        "endpoints": "/api/*"
    }

# Routes API
@api_router.get("/")
async def root():
    return {"message": "API Etienne - Assistant IA pour les etudiants quebecois fourni par le College Champagneur"}

@api_router.get("/quota-status")
async def get_quota_status():
    """Obtenir le statut du quota de requêtes Gemini"""
    return gemini_quota.get_status()

@api_router.post("/chat", response_model=ChatMessage)
async def chat_with_ai(request: ChatRequest, http_request: Request):
    """Endpoint principal pour le chat avec l IA"""
    try:
        # Extraction de l IP et verification du rate limiting
        client_ip = get_client_ip(http_request)
        user_agent = http_request.headers.get("User-Agent", "Unknown")
        
        # Verifier le rate limiting
        if not check_rate_limit(client_ip):
            raise HTTPException(
                status_code=429,
                detail="Trop de requetes. Veuillez patienter avant de reessayer."
            )
        
        # Generation de un session_id si non fourni
        session_id = request.session_id or str(uuid.uuid4())
        
        # 🚨 DETECTION DE CONTENU ILLEGAL
        illegal_check = await detect_illegal_content(
            message=request.message,
            session_id=session_id,
            ip_address=client_ip,
            user_agent=user_agent
        )
        
        if illegal_check["detected"]:
            # Bloquer et retourner un message de avertissement
            warning_message = (
                "⚠️ ATTENTION: Votre demande a ete detectee comme potentiellement dangereuse ou illegale.\n\n"
                "Je ne peux pas vous aider avec ce type de contenu. "
                "Si vous avez des questions legitimes sur la securite ou l education, "
                "veuillez reformuler votre demande de maniere appropriee.\n\n"
                "Cette tentative a ete enregistree conformement a nos politiques de securite."
            )
            
            chat_message = ChatMessage(
                session_id=session_id,
                message=request.message,
                response=warning_message,
                message_type=request.message_type,
                trust_score=0.0,
                sources=["Systeme de securite Etienne"]
            )
            
            # Ne pas sauvegarder le message utilisateur, seulement l avertissement
            return chat_message
        
        # Obtention de la reponse IA (si contenu legal) avec historique de conversation
        conversation_history = request.conversation_history if request.conversation_history else []
        ai_result = await get_ai_response(request.message, request.message_type, conversation_history)
        
        # Creation de l objet ChatMessage
        chat_message = ChatMessage(
            session_id=session_id,
            message=request.message,
            response=ai_result["response"],
            message_type=request.message_type,
            trust_score=ai_result["trust_score"],
            sources=ai_result["sources"],
            image_base64=ai_result.get("image_base64"),
            images=ai_result.get("images")  # Pour plusieurs images
        )
        
        # Sauvegarde en base de donnees
        await db.chat_messages.insert_one(chat_message.dict())
        
        # Enregistrer la requête dans les statistiques
        # Extraire l'email de l'utilisateur du token si disponible
        user_email = None
        user_license = None
        auth_header = http_request.headers.get("Authorization")
        if auth_header and auth_header.startswith("Bearer "):
            try:
                token = auth_header.split(" ")[1]
                payload = jwt.decode(token, JWT_SECRET, algorithms=["HS256"])
                user_email = payload.get("email")
                # Récupérer la licence de l'utilisateur
                user_doc = await db.users.find_one({"email": user_email})
                if user_doc:
                    user_license = user_doc.get("license_key")
            except:
                pass
        
        # Enregistrer dans request_logs
        await db.request_logs.insert_one({
            "user_email": user_email,
            "license_key": user_license,
            "message_type": request.message_type,
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "year": datetime.now(timezone.utc).year,
            "month": datetime.now(timezone.utc).month,
            "ip_address": client_ip
        })
        
        # Incrémenter le compteur de messages de l'utilisateur
        if user_email:
            await db.users.update_one(
                {"email": user_email},
                {"$inc": {"message_count": 1}}
            )
        
        return chat_message
        
    except Exception as e:
        logging.error(f"Erreur chat: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors du traitement de la demande")

@api_router.get("/chat/history/{session_id}", response_model=List[ChatMessage])
async def get_chat_history(session_id: str):
    """Recupere l historique de une session de chat"""
    try:
        messages = await db.chat_messages.find(
            {"session_id": session_id}
        ).sort("timestamp", 1).to_list(100)
        
        return [ChatMessage(**message) for message in messages]
        
    except Exception as e:
        logging.error(f"Erreur historique: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de la recuperation de l historique")

@api_router.post("/sources/analyze")
async def analyze_sources(sources: List[str]):
    """Analyse la fiabilite de une liste de sources"""
    try:
        analyzed_sources = []
        
        for source_url in sources:
            trust_score = calculate_trust_score(source_url)
            analyzed_sources.append({
                "url": source_url,
                "trust_score": trust_score,
                "trust_level": "Tres fiable" if trust_score >= 0.8 else 
                              "Fiable" if trust_score >= 0.6 else 
                              "Moderement fiable" if trust_score >= 0.4 else "Peu fiable",
                "recommendation": "Source recommandee" if trust_score >= 0.7 else 
                                "Verifier avec de autres sources" if trust_score >= 0.5 else 
                                "Source non recommandee"
            })
        
        return {"analyzed_sources": analyzed_sources}
        
    except Exception as e:
        logging.error(f"Erreur analyse sources: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de l analyse des sources")

def generate_pdf_document(title: str, content: str) -> BytesIO:
    """Genere un document PDF structure et elegant avec logo et parsing Markdown"""
    from reportlab.platypus import Image as RLImage, ListFlowable, ListItem
    from reportlab.lib.colors import HexColor, black, white
    from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY, TA_LEFT
    from reportlab.platypus import PageBreak, HRFlowable
    import re
    
    def clean_latex_for_export(text):
        """Nettoie le code LaTeX - ASCII-safe uniquement"""
        if not text:
            return text
        cleaned = text
        # Fractions
        for _ in range(5):
            cleaned = re.sub(r'\\frac\{([^{}]*)\}\{([^{}]*)\}', r'(\1/\2)', cleaned)
        # Racines
        cleaned = re.sub(r'\\sqrt\[(\d+)\]\{([^{}]*)\}', r'racine_\1(\2)', cleaned)
        cleaned = re.sub(r'\\sqrt\{([^{}]*)\}', r'racine(\1)', cleaned)
        # Délimiteurs $
        cleaned = re.sub(r'\$\$([^$]+)\$\$', r' \1 ', cleaned)
        cleaned = re.sub(r'\$([^$]+)\$', r'\1', cleaned)
        # Puissances et indices
        cleaned = re.sub(r'\^{([^}]*)}', r'^\1', cleaned)
        cleaned = re.sub(r'_{([^}]*)}', r'_\1', cleaned)
        # Opérateurs
        cleaned = re.sub(r'\\times', ' x ', cleaned)
        cleaned = re.sub(r'\\cdot', '.', cleaned)
        cleaned = re.sub(r'\\div', ' / ', cleaned)
        cleaned = re.sub(r'\\pm', '+/-', cleaned)
        cleaned = re.sub(r'\\leq', '<=', cleaned)
        cleaned = re.sub(r'\\geq', '>=', cleaned)
        cleaned = re.sub(r'\\neq', '!=', cleaned)
        cleaned = re.sub(r'\\approx', '~=', cleaned)
        cleaned = re.sub(r'\\infty', 'infini', cleaned)
        # Lettres grecques
        cleaned = re.sub(r'\\pi', 'pi', cleaned)
        cleaned = re.sub(r'\\alpha', 'alpha', cleaned)
        cleaned = re.sub(r'\\beta', 'beta', cleaned)
        cleaned = re.sub(r'\\theta', 'theta', cleaned)
        cleaned = re.sub(r'\\Delta', 'Delta', cleaned)
        # Flèches
        cleaned = re.sub(r'\\rightarrow', '->', cleaned)
        cleaned = re.sub(r'\\Rightarrow', '=>', cleaned)
        # Autres
        cleaned = re.sub(r'\\text\{([^}]*)\}', r'\1', cleaned)
        cleaned = re.sub(r'\\textbf\{([^}]*)\}', r'\1', cleaned)
        cleaned = re.sub(r'\\left', '', cleaned)
        cleaned = re.sub(r'\\right', '', cleaned)
        cleaned = re.sub(r'\\quad', '  ', cleaned)
        cleaned = re.sub(r'\\,', ' ', cleaned)
        cleaned = re.sub(r'\\;', ' ', cleaned)
        cleaned = re.sub(r'\\\\', ' ', cleaned)
        cleaned = re.sub(r'\\([a-zA-Z]+)', r'\1', cleaned)
        cleaned = re.sub(r'\{\}', '', cleaned)
        for _ in range(3):
            cleaned = re.sub(r'\{([^{}]*)\}', r'\1', cleaned)
        cleaned = re.sub(r'  +', ' ', cleaned)
        return cleaned
    
    # Nettoyer le LaTeX du contenu
    content = clean_latex_for_export(content)
    
    buffer = BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter, topMargin=50, bottomMargin=72)
    styles = getSampleStyleSheet()
    
    # Styles personnalisés
    title_style = styles['Title'].clone('CustomTitle')
    title_style.fontSize = 22
    title_style.spaceAfter = 20
    title_style.textColor = HexColor('#8B0000')
    title_style.alignment = TA_CENTER
    
    heading1_style = styles['Heading1'].clone('CustomH1')
    heading1_style.fontSize = 16
    heading1_style.textColor = HexColor('#8B0000')
    heading1_style.spaceAfter = 10
    heading1_style.spaceBefore = 15
    
    heading2_style = styles['Heading2'].clone('CustomH2')
    heading2_style.fontSize = 14
    heading2_style.textColor = HexColor('#2563eb')
    heading2_style.spaceAfter = 8
    heading2_style.spaceBefore = 12
    
    heading3_style = styles['Heading3'].clone('CustomH3')
    heading3_style.fontSize = 12
    heading3_style.textColor = HexColor('#6b7280')
    heading3_style.spaceAfter = 6
    heading3_style.spaceBefore = 10
    
    body_style = styles['Normal'].clone('CustomBody')
    body_style.fontSize = 11
    body_style.spaceAfter = 6
    body_style.alignment = TA_JUSTIFY
    
    bullet_style = styles['Normal'].clone('CustomBullet')
    bullet_style.fontSize = 11
    bullet_style.leftIndent = 20
    bullet_style.spaceAfter = 4
    
    story = []
    
    # Logo
    logo_path = Path(__file__).parent / 'logo_champagneur.jpg'
    if logo_path.exists():
        try:
            logo = RLImage(str(logo_path), width=80, height=80)
            logo.hAlign = 'CENTER'
            story.append(logo)
            story.append(Spacer(1, 10))
        except Exception as e:
            logging.warning(f"Impossible de charger le logo: {e}")
    
    # Titre principal
    header_table = Table([[title]], colWidths=[500])
    header_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTSIZE', (0, 0), (0, 0), 20),
        ('TEXTCOLOR', (0, 0), (0, 0), HexColor('#8B0000')),
        ('FONTNAME', (0, 0), (0, 0), 'Helvetica-Bold'),
    ]))
    story.append(header_table)
    story.append(Spacer(1, 10))
    
    # Sous-titre
    subtitle_table = Table([['Collège Champagneur - Étienne Assistant IA']], colWidths=[500])
    subtitle_table.setStyle(TableStyle([
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTSIZE', (0, 0), (0, 0), 10),
        ('TEXTCOLOR', (0, 0), (0, 0), HexColor('#6b7280')),
        ('FONTNAME', (0, 0), (0, 0), 'Helvetica-Oblique'),
    ]))
    story.append(subtitle_table)
    story.append(Spacer(1, 15))
    story.append(HRFlowable(width="100%", thickness=2, color=HexColor('#f97316')))
    story.append(Spacer(1, 20))
    
    def format_text_with_bold(text):
        """Convertit **texte** en <b>texte</b> pour ReportLab"""
        text = re.sub(r'\*\*([^*]+)\*\*', r'<b>\1</b>', text)
        return text
    
    # Traitement ligne par ligne
    lines = content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i].rstrip()
        
        if not line:
            i += 1
            continue
        
        # Titre niveau 1
        if line.startswith('# ') and not line.startswith('## '):
            clean_title = line[2:].strip()
            story.append(Paragraph(clean_title, heading1_style))
            i += 1
            continue
        
        # Titre niveau 2
        if line.startswith('## '):
            clean_title = line[3:].strip()
            story.append(Paragraph(clean_title, heading2_style))
            i += 1
            continue
        
        # Titre niveau 3
        if line.startswith('### '):
            clean_title = line[4:].strip()
            story.append(Paragraph(clean_title, heading3_style))
            i += 1
            continue
        
        # Ligne de séparation
        if line.strip() == '---' or line.strip().startswith('───'):
            story.append(Spacer(1, 5))
            story.append(HRFlowable(width="80%", thickness=1, color=HexColor('#e5e7eb')))
            story.append(Spacer(1, 5))
            i += 1
            continue
        
        # Liste à puces
        if re.match(r'^[\-\*•]\s+', line):
            item_text = re.sub(r'^[\-\*•]\s+', '', line)
            formatted_text = format_text_with_bold(item_text)
            story.append(Paragraph(f"• {formatted_text}", bullet_style))
            i += 1
            continue
        
        # Sous-liste (indentée)
        if re.match(r'^\s+[\-\*•]\s+', line):
            item_text = re.sub(r'^\s+[\-\*•]\s+', '', line)
            formatted_text = format_text_with_bold(item_text)
            indent_style = bullet_style.clone('IndentBullet')
            indent_style.leftIndent = 40
            story.append(Paragraph(f"  ○ {formatted_text}", indent_style))
            i += 1
            continue
        
        # Liste numérotée
        if re.match(r'^\d+[\.\)]\s+', line):
            item_text = re.sub(r'^\d+[\.\)]\s+', '', line)
            formatted_text = format_text_with_bold(item_text)
            story.append(Paragraph(formatted_text, bullet_style))
            i += 1
            continue
        
        # Tableau Markdown
        if line.startswith('|') and '|' in line[1:]:
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                table_lines.append(lines[i].strip())
                i += 1
            
            if len(table_lines) >= 2:
                header_cells = [cell.strip().replace('**', '') for cell in table_lines[0].split('|') if cell.strip()]
                data_start = 2 if len(table_lines) > 1 and re.match(r'^[\|\-:\s]+$', table_lines[1]) else 1
                
                table_data = [header_cells]
                for table_line in table_lines[data_start:]:
                    cells = [cell.strip().replace('**', '') for cell in table_line.split('|') if cell.strip()]
                    if cells:
                        table_data.append(cells)
                
                if table_data:
                    num_cols = max(len(row) for row in table_data)
                    col_width = 480 / num_cols
                    
                    table = Table(table_data, colWidths=[col_width] * num_cols)
                    table.setStyle(TableStyle([
                        ('BACKGROUND', (0, 0), (-1, 0), HexColor('#f3f4f6')),
                        ('TEXTCOLOR', (0, 0), (-1, 0), HexColor('#1f2937')),
                        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                        ('FONTSIZE', (0, 0), (-1, -1), 9),
                        ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                        ('GRID', (0, 0), (-1, -1), 0.5, HexColor('#d1d5db')),
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('TOPPADDING', (0, 0), (-1, -1), 6),
                        ('BOTTOMPADDING', (0, 0), (-1, -1), 6),
                    ]))
                    story.append(table)
                    story.append(Spacer(1, 10))
            continue
        
        # Paragraphe normal
        formatted_text = format_text_with_bold(line)
        story.append(Paragraph(formatted_text, body_style))
        i += 1
    
    # Footer
    story.append(Spacer(1, 30))
    story.append(HRFlowable(width="100%", thickness=1, color=HexColor('#e5e7eb')))
    story.append(Spacer(1, 10))
    
    footer_table = Table([
        ['Généré par Étienne', f"{datetime.now().strftime('%d/%m/%Y à %H:%M')}", 'Collège Champagneur'],
    ], colWidths=[150, 200, 150])
    footer_table.setStyle(TableStyle([
        ('FONTSIZE', (0, 0), (-1, -1), 9),
        ('TEXTCOLOR', (0, 0), (-1, -1), HexColor('#6b7280')),
        ('ALIGN', (0, 0), (0, 0), 'LEFT'),
        ('ALIGN', (1, 0), (1, 0), 'CENTER'),
        ('ALIGN', (2, 0), (2, 0), 'RIGHT'),
    ]))
    story.append(footer_table)
    
    doc.build(story)
    buffer.seek(0)
    return buffer

def generate_docx_document(title: str, content: str) -> BytesIO:
    """Genere un document Word avec belle presentation, logo et parsing Markdown"""
    from docx.shared import Cm, Inches
    from docx.enum.style import WD_STYLE_TYPE
    import re
    
    def clean_latex_for_export(text):
        """Nettoie le code LaTeX - ASCII-safe uniquement"""
        if not text:
            return text
        cleaned = text
        for _ in range(5):
            cleaned = re.sub(r'\\frac\{([^{}]*)\}\{([^{}]*)\}', r'(\1/\2)', cleaned)
        cleaned = re.sub(r'\\sqrt\[(\d+)\]\{([^{}]*)\}', r'racine_\1(\2)', cleaned)
        cleaned = re.sub(r'\\sqrt\{([^{}]*)\}', r'racine(\1)', cleaned)
        cleaned = re.sub(r'\$\$([^$]+)\$\$', r' \1 ', cleaned)
        cleaned = re.sub(r'\$([^$]+)\$', r'\1', cleaned)
        cleaned = re.sub(r'\^{([^}]*)}', r'^\1', cleaned)
        cleaned = re.sub(r'_{([^}]*)}', r'_\1', cleaned)
        cleaned = re.sub(r'\\times', ' x ', cleaned)
        cleaned = re.sub(r'\\cdot', '.', cleaned)
        cleaned = re.sub(r'\\div', ' / ', cleaned)
        cleaned = re.sub(r'\\pm', '+/-', cleaned)
        cleaned = re.sub(r'\\leq', '<=', cleaned)
        cleaned = re.sub(r'\\geq', '>=', cleaned)
        cleaned = re.sub(r'\\neq', '!=', cleaned)
        cleaned = re.sub(r'\\approx', '~=', cleaned)
        cleaned = re.sub(r'\\infty', 'infini', cleaned)
        cleaned = re.sub(r'\\pi', 'pi', cleaned)
        cleaned = re.sub(r'\\alpha', 'alpha', cleaned)
        cleaned = re.sub(r'\\beta', 'beta', cleaned)
        cleaned = re.sub(r'\\theta', 'theta', cleaned)
        cleaned = re.sub(r'\\Delta', 'Delta', cleaned)
        cleaned = re.sub(r'\\rightarrow', '->', cleaned)
        cleaned = re.sub(r'\\Rightarrow', '=>', cleaned)
        cleaned = re.sub(r'\\text\{([^}]*)\}', r'\1', cleaned)
        cleaned = re.sub(r'\\textbf\{([^}]*)\}', r'\1', cleaned)
        cleaned = re.sub(r'\\left', '', cleaned)
        cleaned = re.sub(r'\\right', '', cleaned)
        cleaned = re.sub(r'\\quad', '  ', cleaned)
        cleaned = re.sub(r'\\,', ' ', cleaned)
        cleaned = re.sub(r'\\;', ' ', cleaned)
        cleaned = re.sub(r'\\\\', ' ', cleaned)
        cleaned = re.sub(r'\\([a-zA-Z]+)', r'\1', cleaned)
        cleaned = re.sub(r'\{\}', '', cleaned)
        for _ in range(3):
            cleaned = re.sub(r'\{([^{}]*)\}', r'\1', cleaned)
        cleaned = re.sub(r'  +', ' ', cleaned)
        return cleaned
    
    # Nettoyer le LaTeX du contenu
    content = clean_latex_for_export(content)
    
    doc = Document()
    
    # Logo du Collège Champagneur
    logo_path = Path(__file__).parent / 'logo_champagneur.jpg'
    if logo_path.exists():
        try:
            logo_para = doc.add_paragraph()
            logo_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            logo_run = logo_para.add_run()
            logo_run.add_picture(str(logo_path), width=Cm(3))
        except Exception as e:
            logging.warning(f"Impossible de charger le logo Word: {e}")
    
    # En-tete avec titre
    header_para = doc.add_paragraph()
    header_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title_run = header_para.add_run(title)
    title_run.bold = True
    title_run.font.size = DocxPt(20)
    title_run.font.color.rgb = RGBColor(139, 0, 0)
    
    # Sous-titre
    subtitle_para = doc.add_paragraph()
    subtitle_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle_run = subtitle_para.add_run("Collège Champagneur - Étienne Assistant IA")
    subtitle_run.font.size = DocxPt(10)
    subtitle_run.font.color.rgb = RGBColor(107, 114, 128)
    subtitle_run.italic = True
    
    # Ligne de separation
    separator_para = doc.add_paragraph()
    separator_run = separator_para.add_run("─" * 60)
    separator_run.font.color.rgb = RGBColor(139, 0, 0)
    
    doc.add_paragraph()
    
    def add_formatted_text(paragraph, text):
        """Ajoute du texte avec formatage gras (**texte**)"""
        # Pattern pour détecter **texte**
        pattern = r'\*\*([^*]+)\*\*'
        last_end = 0
        
        for match in re.finditer(pattern, text):
            # Ajouter le texte avant le match
            if match.start() > last_end:
                run = paragraph.add_run(text[last_end:match.start()])
                run.font.size = DocxPt(11)
            
            # Ajouter le texte en gras
            bold_run = paragraph.add_run(match.group(1))
            bold_run.bold = True
            bold_run.font.size = DocxPt(11)
            bold_run.font.color.rgb = RGBColor(37, 99, 235)  # Bleu
            
            last_end = match.end()
        
        # Ajouter le reste du texte
        if last_end < len(text):
            run = paragraph.add_run(text[last_end:])
            run.font.size = DocxPt(11)
    
    # Traitement du contenu ligne par ligne
    lines = content.split('\n')
    i = 0
    
    while i < len(lines):
        line = lines[i].rstrip()
        
        # Ignorer les lignes vides au début
        if not line:
            i += 1
            continue
        
        # Titre niveau 1: # Titre
        if line.startswith('# ') and not line.startswith('## '):
            clean_title = line[2:].strip()
            heading = doc.add_heading(clean_title, level=1)
            for run in heading.runs:
                run.font.color.rgb = RGBColor(139, 0, 0)
                run.font.size = DocxPt(18)
            i += 1
            continue
        
        # Titre niveau 2: ## Titre
        if line.startswith('## '):
            clean_title = line[3:].strip()
            heading = doc.add_heading(clean_title, level=2)
            for run in heading.runs:
                run.font.color.rgb = RGBColor(37, 99, 235)
                run.font.size = DocxPt(16)
            i += 1
            continue
        
        # Titre niveau 3: ### Titre
        if line.startswith('### '):
            clean_title = line[4:].strip()
            heading = doc.add_heading(clean_title, level=3)
            for run in heading.runs:
                run.font.color.rgb = RGBColor(107, 114, 128)
                run.font.size = DocxPt(14)
            i += 1
            continue
        
        # Ligne de séparation ---
        if line.strip() == '---' or line.strip().startswith('───'):
            sep_para = doc.add_paragraph()
            sep_run = sep_para.add_run("─" * 50)
            sep_run.font.color.rgb = RGBColor(200, 200, 200)
            i += 1
            continue
        
        # Liste à puces: - item ou * item ou • item
        if re.match(r'^[\-\*•]\s+', line):
            # Calculer l'indentation
            indent_level = 0
            original_line = lines[i]
            stripped = original_line.lstrip()
            indent_spaces = len(original_line) - len(stripped)
            indent_level = indent_spaces // 4
            
            item_text = re.sub(r'^[\-\*•]\s+', '', stripped)
            para = doc.add_paragraph(style='List Bullet')
            para.paragraph_format.left_indent = Inches(0.25 * (indent_level + 1))
            add_formatted_text(para, item_text)
            i += 1
            continue
        
        # Liste numérotée: 1. item ou 1) item
        if re.match(r'^\d+[\.\)]\s+', line):
            item_text = re.sub(r'^\d+[\.\)]\s+', '', line)
            para = doc.add_paragraph(style='List Number')
            add_formatted_text(para, item_text)
            i += 1
            continue
        
        # Tableau Markdown: | col1 | col2 |
        if line.startswith('|') and '|' in line[1:]:
            # Collecter toutes les lignes du tableau
            table_lines = []
            while i < len(lines) and lines[i].strip().startswith('|'):
                table_lines.append(lines[i].strip())
                i += 1
            
            if len(table_lines) >= 2:
                # Parser les colonnes
                header_cells = [cell.strip() for cell in table_lines[0].split('|') if cell.strip()]
                
                # Ignorer la ligne de séparation (|---|---|)
                data_start = 1
                if len(table_lines) > 1 and re.match(r'^[\|\-:\s]+$', table_lines[1]):
                    data_start = 2
                
                # Créer le tableau
                num_cols = len(header_cells)
                num_rows = len(table_lines) - data_start + 1
                
                table = doc.add_table(rows=num_rows, cols=num_cols)
                table.style = 'Table Grid'
                
                # En-tête
                header_row = table.rows[0]
                for j, cell_text in enumerate(header_cells):
                    cell = header_row.cells[j]
                    cell.text = cell_text.replace('**', '')
                    for para in cell.paragraphs:
                        for run in para.runs:
                            run.bold = True
                            run.font.size = DocxPt(10)
                
                # Données
                for row_idx, table_line in enumerate(table_lines[data_start:], start=1):
                    if row_idx < num_rows:
                        cells = [cell.strip() for cell in table_line.split('|') if cell.strip()]
                        for j, cell_text in enumerate(cells):
                            if j < num_cols:
                                cell = table.rows[row_idx].cells[j]
                                cell.text = cell_text.replace('**', '')
                                for para in cell.paragraphs:
                                    for run in para.runs:
                                        run.font.size = DocxPt(10)
                
                doc.add_paragraph()
            continue
        
        # Paragraphe normal
        para = doc.add_paragraph()
        para.alignment = WD_ALIGN_PARAGRAPH.JUSTIFY
        add_formatted_text(para, line)
        i += 1
    
    # Footer
    doc.add_paragraph()
    footer_separator = doc.add_paragraph()
    footer_separator_run = footer_separator.add_run("─" * 80)
    footer_separator_run.font.color.rgb = RGBColor(229, 231, 235)
    
    footer_para = doc.add_paragraph()
    footer_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_left = footer_para.add_run("Généré par Étienne")
    footer_left.font.size = DocxPt(9)
    footer_left.font.color.rgb = RGBColor(107, 114, 128)
    footer_center = footer_para.add_run(f" • {datetime.now().strftime('%d/%m/%Y à %H:%M')} • ")
    footer_center.font.size = DocxPt(9)
    footer_center.font.color.rgb = RGBColor(107, 114, 128)
    footer_right = footer_para.add_run("Collège Champagneur")
    footer_right.font.size = DocxPt(9)
    footer_right.font.color.rgb = RGBColor(107, 114, 128)
    footer_right.italic = True
    
    buffer = BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def generate_pptx_document(title: str, content: str) -> BytesIO:
    """Genere une belle presentation PowerPoint"""
    prs = Presentation()
    
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    
    # === SLIDE DE TITRE ===
    slide_layout = prs.slide_layouts[0]  # Title slide
    title_slide = prs.slides.add_slide(slide_layout)
    
    # Titre principal
    title_shape = title_slide.shapes.title
    title_shape.text = f"🎓 {title} 📚"
    title_paragraph = title_shape.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.color.rgb = RGBColor(249, 115, 22)  # Orange
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    subtitle_shape = title_slide.placeholders[1]
    subtitle_shape.text = f"Presente par Etienne\nAssistant IA Educatif\nCollege Champagneur\n\n{datetime.now().strftime('%d %B %Y')}"
    subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
    subtitle_paragraph.font.size = Pt(18)
    subtitle_paragraph.font.color.rgb = RGBColor(37, 99, 235)  # Bleu
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    # === TRAITEMENT DU CONTENU ===
    sections = content.split('\n\n')
    current_slide_content = []
    slide_counter = 1
    
    for section in sections:
        if section.strip():
            current_slide_content.append(section.strip())
            
            # Creer une nouvelle slide tous les 4 points ou si section trop longue
            if len(current_slide_content) >= 4 or len(section) > 300:
                slide_counter += 1
                slide = prs.slides.add_slide(prs.slide_layouts[1])  # Content layout
                
                # Titre de la slide
                slide.shapes.title.text = f"📖 Section {slide_counter - 1}"
                title_para = slide.shapes.title.text_frame.paragraphs[0]
                title_para.font.size = Pt(32)
                title_para.font.color.rgb = RGBColor(37, 99, 235)
                
                # Contenu
                content_placeholder = slide.placeholders[1]
                text_frame = content_placeholder.text_frame
                text_frame.clear()  # Nettoyer le contenu par defaut
                
                for i, item in enumerate(current_slide_content):
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    
                    # Detecter si c'est un titre
                    if (item.startswith('#') or 
                        (len(item) < 80 and item.isupper()) or
                        (i == 0 and len(item) < 100)):
                        # Style titre
                        p.text = f"• {item.replace('#', '').strip()}"
                        p.font.size = Pt(20)
                        p.font.color.rgb = RGBColor(249, 115, 22)  # Orange
                        p.font.bold = True
                    else:
                        # Style normal
                        p.text = f"  ◦ {item[:200]}{'...' if len(item) > 200 else ''}"
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor(55, 65, 81)  # Gris fonce
                    
                    p.space_after = Pt(12)
                
                current_slide_content = []
    
    # Derniere slide s'il reste du contenu
    if current_slide_content:
        slide_counter += 1
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"📖 Section {slide_counter - 1}"
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        title_para.font.size = Pt(32)
        title_para.font.color.rgb = RGBColor(37, 99, 235)
        
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        
        for i, item in enumerate(current_slide_content):
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()
            
            p.text = f"• {item[:150]}{'...' if len(item) > 150 else ''}"
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(55, 65, 81)
            p.space_after = Pt(12)
    
    # === SLIDE DE CONCLUSION ===
    conclusion_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
    
    # Ajouter une zone de texte pour la conclusion
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(4)
    
    textbox = conclusion_slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    
    # Titre de conclusion
    conclusion_title = text_frame.add_paragraph()
    conclusion_title.text = "🎉 Merci pour votre attention !"
    conclusion_title.font.size = Pt(36)
    conclusion_title.font.color.rgb = RGBColor(249, 115, 22)
    conclusion_title.font.bold = True
    conclusion_title.alignment = PP_ALIGN.CENTER
    
    # Texte de conclusion
    conclusion_text = text_frame.add_paragraph()
    conclusion_text.text = "\n\nEtienne - Assistant IA Educatif\nCollege Champagneur\n\n📚 Continuez a apprendre et a explorer ! 🎓"
    conclusion_text.font.size = Pt(18)
    conclusion_text.font.color.rgb = RGBColor(37, 99, 235)
    conclusion_text.alignment = PP_ALIGN.CENTER
    
    buffer = BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

def generate_xlsx_document(title: str, content: str) -> BytesIO:
    """Genere un fichier Excel"""
    buffer = BytesIO()
    wb = Workbook()
    ws = wb.active
    ws.title = "Document Etienne"
    
    # En-tetes
    ws['A1'] = title
    ws['A1'].font = Font(bold=True, size=16)
    ws['A1'].fill = PatternFill(start_color='366092', end_color='366092', fill_type='solid')
    ws['A1'].font = Font(bold=True, size=16, color='FFFFFF')
    
    # Contenu
    paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
    
    ws['A3'] = "Points cles :"
    ws['A3'].font = Font(bold=True)
    
    row = 4
    for i, para in enumerate(paragraphs):
        ws[f'A{row}'] = f"{i + 1}."
        ws[f'B{row}'] = para
        ws[f'A{row}'].font = Font(bold=True)
        row += 1
    
    # Ajuster la largeur des colonnes
    ws.column_dimensions['A'].width = 5
    ws.column_dimensions['B'].width = 80
    
    # Pied de page
    ws[f'A{row + 2}'] = f"Genere par Etienne le {datetime.now().strftime('%d/%m/%Y a %H:%M')}"
    ws[f'A{row + 2}'].font = Font(italic=True)
    
    wb.save(buffer)
    buffer.seek(0)
    return buffer

@api_router.post("/generate-document")
async def generate_document(request: DocumentRequest):
    """Genere un document dans le format demande"""
    try:
        # Validation du format
        allowed_formats = ['pdf', 'docx', 'pptx', 'xlsx']
        if request.format not in allowed_formats:
            raise HTTPException(status_code=400, detail=f"Format non supporte. Formats autorises: {', '.join(allowed_formats)}")
        
        # Generation du nom de fichier
        if not request.filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{request.title.replace(' ', '_')}_{timestamp}.{request.format}"
        else:
            filename = request.filename if request.filename.endswith(f".{request.format}") else f"{request.filename}.{request.format}"
        
        # Generation du document selon le format
        if request.format == 'pdf':
            buffer = generate_pdf_document(request.title, request.content)
            media_type = "application/pdf"
        elif request.format == 'docx':
            buffer = generate_docx_document(request.title, request.content)
            media_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        elif request.format == 'pptx':
            buffer = generate_pptx_document(request.title, request.content)
            media_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        elif request.format == 'xlsx':
            buffer = generate_xlsx_document(request.title, request.content)
            media_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
        
        return StreamingResponse(
            BytesIO(buffer.read()),
            media_type=media_type,
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
        
    except Exception as e:
        logging.error(f"Erreur generation document: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de la generation du document: {str(e)}")

@api_router.post("/upload-file")
async def upload_and_extract_file(file: UploadFile = File(...)):
    """Upload un fichier et extrait son contenu texte"""
    try:
        # Verifier la taille du fichier (max 10MB)
        max_size = 10 * 1024 * 1024  # 10MB
        file_size = 0
        
        # Lire le fichier pour verifier la taille
        content = await file.read()
        file_size = len(content)
        
        if file_size > max_size:
            raise HTTPException(
                status_code=413,
                detail="Fichier trop volumineux. Taille maximale: 10MB"
            )
        
        # Remettre le pointeur au debut pour l extraction
        file.file = BytesIO(content)
        
        # Extraire le texte
        extracted_text = await extract_text_from_file(file)
        
        # Limiter la longueur du texte extrait (pour eviter les tokens excessifs)
        max_text_length = 10000  # ~2500 mots
        if len(extracted_text) > max_text_length:
            extracted_text = extracted_text[:max_text_length] + "\n\n[...Texte tronque pour optimiser l analyse...]"
        
        return {
            "filename": file.filename,
            "file_size": file_size,
            "extracted_text": extracted_text,
            "text_length": len(extracted_text),
            "message": "Fichier traite avec succes. Vous pouvez maintenant poser votre question."
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur upload fichier: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Erreur lors du traitement du fichier: {str(e)}"
        )

@api_router.post("/analyze-file", response_model=ChatMessage)
async def analyze_file_with_question(request: FileAnalysisRequest):
    """Analyse un fichier avec une question specifique - 100% GRATUIT avec Gemini - BILINGUE"""
    try:
        # DETECTION DE LA LANGUE DE LA QUESTION
        detected_language = detect_language(request.question)
        assistant_name = "Steven" if detected_language == "en" else "Étienne"
        
        # Preparer le prompt avec le contenu du fichier (adapté à la langue)
        if detected_language == "en":
            enhanced_message = f"""
CONTEXT: The user has uploaded a document ({request.filename}) and asks the following question:

QUESTION: {request.question}

DOCUMENT CONTENT:
{request.extracted_text}

INSTRUCTIONS: 
- Analyze the document content in relation to the question asked
- Provide a precise answer based on the document content
- If the answer is not in the document, mention it clearly
- Structure your response in a clear and pedagogical manner
"""
            system_message = """You are Steven, an AI assistant specialized in document analysis for Quebec teachers. 
IMPORTANT: The user is writing in ENGLISH, so you MUST respond ENTIRELY in ENGLISH.
You must analyze the provided content and answer the user's question precisely and pedagogically.
Adapt your language to Quebec secondary education level (grades 7-11).

BRANDING: If asked who created you or how you work, say that you are an AI assistant developed and programmed by Collège Champagneur to help Quebec teachers. Never mention Google, Gemini, OpenAI or any other underlying AI technology."""
        else:
            enhanced_message = f"""
CONTEXTE: L'utilisateur a uploadé un document ({request.filename}) et pose la question suivante:

QUESTION: {request.question}

CONTENU DU DOCUMENT:
{request.extracted_text}

INSTRUCTIONS: 
- Analysez le contenu du document en relation avec la question posée
- Fournissez une réponse précise basée sur le contenu du document
- Si la réponse n'est pas dans le document, mentionnez-le clairement
- Structurez votre réponse de manière claire et pédagogique
"""
            system_message = """Tu es Étienne, un assistant IA spécialisé dans l'analyse de documents pour les enseignants québécois. 
Tu dois analyser le contenu fourni et répondre à la question de l'utilisateur de manière précise et pédagogique.
Adapte ton langage au niveau du secondaire québécois et utilise un français accessible.

IMAGE DE MARQUE: Si on te demande qui t'a créé ou comment tu fonctionnes, réponds que tu es un assistant IA développé et programmé par le Collège Champagneur pour aider les enseignants québécois. Ne mentionne jamais Google, Gemini, OpenAI ou toute autre technologie IA sous-jacente."""

        # === GOOGLE GEMINI 2.5 FLASH (100% GRATUIT - Derniere version stable) ===
        model = genai.GenerativeModel(
            model_name='gemini-2.5-flash',
            system_instruction=system_message
        )
        
        # Generer la reponse
        response = await asyncio.to_thread(
            model.generate_content,
            enhanced_message
        )
        
        response_text = response.text
        
        # Creer l objet ChatMessage
        chat_message = ChatMessage(
            session_id=f"file-{request.filename}-{uuid.uuid4()}",
            message=f"📎 Analyse du fichier '{request.filename}': {request.question}",
            response=response_text,
            message_type=request.message_type,
            trust_score=0.90,  # Score eleve pour l analyse de documents
            sources=[request.filename]
        )
        
        # Sauvegarder en base de donnees
        await db.chat_messages.insert_one(chat_message.dict())
        
        return chat_message
        
    except Exception as e:
        logging.error(f"Erreur analyse fichier: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"Erreur lors de l analyse du fichier: {str(e)}"
        )

@api_router.get("/subjects")
async def get_school_subjects():
    """Retourne la liste des matieres du secondaire quebecois (Sec 1 a 5) pour ENSEIGNANTS"""
    subjects = {
        "langues": {
            "name": "Langues",
            "subjects": ["Français", "Anglais", "Anglais enrichi", "Espagnol"]
        },
        "mathematiques": {
            "name": "Mathématiques",
            "subjects": ["Mathématiques CST", "Mathématiques TS", "Mathématiques SN (Sciences naturelles)"]
        },
        "sciences_humaines": {
            "name": "Sciences Humaines",
            "subjects": ["Histoire", "Géographie", "Monde contemporain", "Culture et société québécoise"]
        },
        "sciences": {
            "name": "Sciences",
            "subjects": ["Sciences et technologies", "Applications technologiques et scientifiques"]
        },
        "developpement": {
            "name": "Développement Personnel",
            "subjects": ["Éducation financière", "Méthodologie", "Éducation physique et à la santé"]
        },
        "arts": {
            "name": "Arts",
            "subjects": ["Art dramatique", "Arts plastiques", "Danse", "Musique"]
        }
    }
    return subjects

# Modeles pour les nouvelles fonctionnalites
class AIDetectionRequest(BaseModel):
    text: str

class AIFeedbackRequest(BaseModel):
    text: str
    predicted_probability: float
    actual_is_ai: bool  # True si c'etait vraiment IA, False si humain
    analysis_id: Optional[str] = None

# ==================== NOUVEAUX MODELES POUR SYSTEME DE LICENCES ====================

class LicenseCreate(BaseModel):
    organization_name: str
    license_key: str  # Ex: ETIENNE-ECOLE-MONTMORENCY-2024
    max_users: int
    expiry_date: str  # Format: YYYY-MM-DD
    notes: Optional[str] = None

class LicenseUpdate(BaseModel):
    max_users: Optional[int] = None
    expiry_date: Optional[str] = None
    is_active: Optional[bool] = None
    notes: Optional[str] = None
    license_key: Optional[str] = None  # Pour modifier la clé de licence

class UserSignup(BaseModel):
    full_name: str
    email: str
    password: str
    license_key: str

class UserLogin(BaseModel):
    email: str
    password: str

class BlockedWordCreate(BaseModel):
    word: str
    category: str  # "violence", "drugs", "illegal", "custom", etc.
    severity: str  # "low", "medium", "high", "critical"
    is_exception: bool = False  # Si True, le mot est autorisé malgré détection

class BlockedWordUpdate(BaseModel):
    category: Optional[str] = None
    severity: Optional[str] = None
    is_exception: Optional[bool] = None
    is_active: Optional[bool] = None

# ==================== NOUVEAUX MODELES POUR ADMIN MULTI-NIVEAUX ====================

class LicenseAdminUpdate(BaseModel):
    """Pour désigner ou retirer un admin de licence"""
    license_admin_email: Optional[str] = None  # Email de l'admin de licence

class LicenseKeyUpdate(BaseModel):
    """Pour modifier la clé de licence (admin licence ou super admin)"""
    new_license_key: str

class UserCreate(BaseModel):
    """Pour créer un utilisateur manuellement (admin de licence)"""
    full_name: str
    email: str
    password: str

class UserRemove(BaseModel):
    """Pour retirer un utilisateur (admin de licence)"""
    user_id: str

class ChangeEmailRequest(BaseModel):
    """Pour changer l'email d'un utilisateur"""
    old_email: str
    new_email: str

class UserSelfChangeEmail(BaseModel):
    """Pour qu'un utilisateur change son propre email"""
    new_email: str
    password: str  # Confirmation du mot de passe

class UserSelfChangePassword(BaseModel):
    """Pour qu'un utilisateur change son propre mot de passe"""
    current_password: str
    new_password: str

class AdminChangeUserPassword(BaseModel):
    """Pour qu'un admin change le mot de passe d'un utilisateur"""
    user_email: str
    new_password: str


@api_router.post("/detect-ai")
async def detect_ai_endpoint(request: AIDetectionRequest):
    """Detecte si un texte a ete genere par IA - VERSION AMELIOREE AVEC CLAUDE"""
    try:
        # Utiliser la nouvelle detection LLM
        result = await detect_ai_content_with_llm(request.text)
        
        return {
            "success": True,
            "detection_result": result,
            "message": f"Analyse terminee. Probabilite IA: {result['ai_probability']*100}%"
        }
    except Exception as e:
        logging.error(f"Erreur detection IA: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de la detection IA: {str(e)}")

@api_router.post("/analyze-text")
async def analyze_text_complete(request: AIDetectionRequest):
    """Analyse complete d'un texte (IA + Langue) - VERSION SIMPLIFIEE"""
    try:
        # Detection de langue
        detected_language = detect_language(request.text)
        
        # Detection IA avec le nouveau systeme Claude
        ai_result = await detect_ai_content_with_llm(request.text)
        
        # Analyse combinee
        overall_risk = "Low"
        if ai_result["is_likely_ai"]:
            overall_risk = "High"
        elif ai_result["ai_probability"] > 0.3:
            overall_risk = "Medium"
        
        recommendations = []
        if ai_result["is_likely_ai"]:
            recommendations.append(f"Ce texte semble genere par IA ({ai_result['ai_probability']*100:.0f}% de probabilite). Verifiez l originalite.")
        if not recommendations:
            recommendations.append("Le texte semble original et authentique.")
        
        return {
            "success": True,
            "language": detected_language,
            "ai_detection": ai_result,
            "overall_assessment": {
                "risk_level": overall_risk,
                "recommendations": recommendations,
                "text_length": len(request.text),
                "word_count": len(request.text.split())
            }
        }
        
    except Exception as e:
        logging.error(f"Erreur analyse complete: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de l analyse complete: {str(e)}")

@api_router.post("/ai-feedback")
async def submit_ai_feedback(request: AIFeedbackRequest):
    """
    Permet aux utilisateurs de corriger la detection IA pour ameliorer le systeme.
    Le systeme apprend de ces corrections pour devenir plus precis.
    """
    try:
        # Enregistrer le feedback dans la base de donnees
        feedback_doc = {
            "id": str(uuid.uuid4()),
            "text": request.text,
            "predicted_probability": request.predicted_probability,
            "actual_is_ai": request.actual_is_ai,
            "analysis_id": request.analysis_id,
            "timestamp": datetime.now(timezone.utc)
        }
        
        await db.ai_feedback.insert_one(feedback_doc)
        
        # Faire apprendre le detecteur avance
        advanced_detector.learn_from_feedback(
            text=request.text,
            actual_label=request.actual_is_ai,
            predicted_prob=request.predicted_probability
        )
        
        # Sauvegarder les nouveaux poids dans la base de donnees
        await save_detector_weights()
        
        # Calculer les statistiques de apprentissage
        total_feedback = await db.ai_feedback.count_documents({})
        
        logging.info(f"Feedback enregistre. Total feedbacks: {total_feedback}")
        
        return {
            "success": True,
            "message": "Merci! Votre feedback a ete enregistre et le systeme s'est ameliore.",
            "total_feedback_count": total_feedback,
            "learning_status": "Poids du detecteur mis a jour"
        }
        
    except Exception as e:
        logging.error(f"Erreur enregistrement feedback: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de l enregistrement du feedback: {str(e)}")

@api_router.get("/ai-stats")
async def get_ai_detection_stats():
    """Retourne les statistiques du systeme de auto-apprentissage"""
    try:
        # Compter le nombre de feedbacks
        total_feedback = await db.ai_feedback.count_documents({})
        
        # Recuperer les poids actuels
        weights_doc = await db.detector_weights.find_one({"version": "latest"})
        
        # Calculer l accuracy si possible
        if total_feedback > 0:
            feedbacks = await db.ai_feedback.find().to_list(length=1000)
            
            correct_predictions = sum(
                1 for fb in feedbacks 
                if (fb["predicted_probability"] > 0.5) == fb["actual_is_ai"]
            )
            
            accuracy = (correct_predictions / len(feedbacks)) * 100 if feedbacks else 0
        else:
            accuracy = 0
        
        return {
            "success": True,
            "total_feedback_received": total_feedback,
            "estimated_accuracy": round(accuracy, 1) if total_feedback > 0 else "Pas encore de donnees",
            "last_weights_update": weights_doc.get("updated_at") if weights_doc else None,
            "learning_status": "Actif" if total_feedback > 0 else "En attente de feedback",
            "message": f"Le systeme a appris de {total_feedback} exemples"
        }
        
    except Exception as e:
        logging.error(f"Erreur stats: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de la recuperation des stats: {str(e)}")



# ============================================
# ENDPOINTS ADMIN - LOGS DE SECURITE
# ============================================

@api_router.get("/admin/illegal-content-logs")
async def get_illegal_content_logs(
    limit: int = 100,
    severity: Optional[str] = None,
    category: Optional[str] = None,
    current_admin: dict = Depends(get_current_admin)
):
    """
    Recupere les logs de tentatives de contenu illegal
    AUTHENTIFICATION REQUISE
    """


# ============================================
# ENDPOINTS D'AUTHENTIFICATION ADMIN
# ============================================

class AdminLoginRequest(BaseModel):
    username: str
    password: str

class AdminLoginResponse(BaseModel):
    success: bool
    token: str
    token_type: str = "bearer"
    expires_in: int = 86400  # 24 heures en secondes

@api_router.post("/admin/login", response_model=AdminLoginResponse)
async def admin_login(credentials: AdminLoginRequest):
    """Login admin - genere un JWT token"""
    from security_advanced import ADMIN_USERNAME, ADMIN_PASSWORD_HASH
    
    if credentials.username != ADMIN_USERNAME:
        raise HTTPException(status_code=401, detail="Identifiants invalides")
    
    if not verify_password(credentials.password, ADMIN_PASSWORD_HASH):
        raise HTTPException(status_code=401, detail="Identifiants invalides")
    
    # Creer le token JWT
    token = create_access_token({"username": credentials.username})
    
    return AdminLoginResponse(
        success=True,
        token=token
    )

@api_router.post("/admin/send-daily-report")
async def trigger_daily_report(current_admin: dict = Depends(get_current_admin)):
    """
    Genere et envoie manuellement le rapport quotidien
    AUTHENTIFICATION REQUISE
    """
    try:
        # Recuperer tous les logs
        logs_cursor = await db.security_illegal_logs.find({}).sort("timestamp", -1).to_list(1000)
        
        logs = []
        for log in logs_cursor:
            log['_id'] = str(log.get('_id', ''))
            log['timestamp'] = log.get('timestamp', datetime.now(timezone.utc)).isoformat() if isinstance(log.get('timestamp'), datetime) else str(log.get('timestamp'))
            logs.append(log)
        
        # Statistiques
        stats = {
            "total_logs": await db.security_illegal_logs.count_documents({}),
            "critical_count": await db.security_illegal_logs.count_documents({"severity": "critical"}),
            "high_count": await db.security_illegal_logs.count_documents({"severity": "high"}),
            "category_stats": {}
        }
        
        for category_name in ILLEGAL_CONTENT_PATTERNS.keys():
            count = await db.security_illegal_logs.count_documents({"detected_category": category_name})
            if count > 0:
                stats["category_stats"][category_name] = count
        
        # Envoyer le rapport
        success = await asyncio.to_thread(send_daily_report, logs, stats)
        
        # Sauvegarder localement
        if logs:
            import os
            report_dir = Path(__file__).parent / "security_reports"
            report_dir.mkdir(exist_ok=True)
            
            # CSV
            csv_content = generate_csv_report(logs)
            csv_path = report_dir / f"report_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.csv"
            with open(csv_path, 'wb') as f:
                f.write(csv_content)
            
            # PDF
            pdf_content = generate_pdf_report(logs, stats)
            pdf_path = report_dir / f"report_{datetime.now(timezone.utc).strftime('%Y%m%d_%H%M%S')}.pdf"
            with open(pdf_path, 'wb') as f:
                f.write(pdf_content)
        
        return {
            "success": success,
            "message": "Rapport envoye par email et sauvegarde localement" if success else "Erreur lors de l envoi",
            "logs_count": len(logs),
            "report_saved": True if logs else False
        }
        
    except Exception as e:
        logging.error(f"Erreur generation rapport: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/admin/export-csv")
async def export_logs_csv(
    days: int = 30,
    current_admin: dict = Depends(get_current_admin)
):
    """
    Exporte les logs en CSV
    AUTHENTIFICATION REQUISE
    """
    try:
        cutoff_date = datetime.now(timezone.utc) - timedelta(days=days)
        
        logs_cursor = await db.security_illegal_logs.find(
            {"timestamp": {"$gte": cutoff_date}}
        ).sort("timestamp", -1).to_list(10000)
        
        logs = []
        for log in logs_cursor:
            log['_id'] = str(log.get('_id', ''))
            log['timestamp'] = log.get('timestamp', datetime.now(timezone.utc)).isoformat() if isinstance(log.get('timestamp'), datetime) else str(log.get('timestamp'))
            logs.append(log)
        
        csv_content = generate_csv_report(logs)
        
        return StreamingResponse(
            iter([csv_content]),
            media_type="text/csv",
            headers={
                "Content-Disposition": f"attachment; filename=security_logs_{datetime.now(timezone.utc).strftime('%Y%m%de')}.csv"
            }
        )
        
    except Exception as e:
        logging.error(f"Erreur export CSV: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/admin/export-pdf")
async def export_logs_pdf(
    days: int = 30,
    current_admin: dict = Depends(get_current_admin)
):
    """
    Exporte les logs en PDF
    AUTHENTIFICATION REQUISE
    """
    try:
        cutoff_date = datetime.now(timezone.utc) - timedelta(days=days)
        
        logs_cursor = await db.security_illegal_logs.find(
            {"timestamp": {"$gte": cutoff_date}}
        ).sort("timestamp", -1).to_list(10000)
        
        logs = []
        for log in logs_cursor:
            log['_id'] = str(log.get('_id', ''))
            log['timestamp'] = log.get('timestamp', datetime.now(timezone.utc)).isoformat() if isinstance(log.get('timestamp'), datetime) else str(log.get('timestamp'))
            logs.append(log)
        
        # Statistiques
        stats = {
            "total_logs": len(logs),
            "critical_count": sum(1 for log in logs if log.get('severity') == 'critical'),
            "high_count": sum(1 for log in logs if log.get('severity') == 'high'),
            "category_stats": {}
        }
        
        for log in logs:
            cat = log.get('detected_category', 'unknown')
            stats["category_stats"][cat] = stats["category_stats"].get(cat, 0) + 1
        
        pdf_content = generate_pdf_report(logs, stats)
        
        return StreamingResponse(
            iter([pdf_content]),
            media_type="application/pdf",
            headers={
                "Content-Disposition": f"attachment; filename=security_report_{datetime.now(timezone.utc).strftime('%Y%m%de')}.pdf"
            }
        )
        
    except Exception as e:
        logging.error(f"Erreur export PDF: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/admin/rate-limit-stats")
async def get_rate_limit_stats(current_admin: dict = Depends(get_current_admin)):
    """
    Obtient les statistiques de rate limiting
    AUTHENTIFICATION REQUISE
    """
    from security_advanced import rate_limit_storage
    
    stats = {
        "total_ips_tracked": len(rate_limit_storage),
        "ips_near_limit": 0,
        "ips_blocked": 0,
        "top_ips": []
    }
    
    from security_advanced import RATE_LIMIT_ATTEMPTS
    
    for ip, attempts in rate_limit_storage.items():
        attempt_count = len(attempts)
        if attempt_count >= RATE_LIMIT_ATTEMPTS:
            stats["ips_blocked"] += 1
        elif attempt_count >= RATE_LIMIT_ATTEMPTS * 0.7:
            stats["ips_near_limit"] += 1
        
        stats["top_ips"].append({
            "ip": ip,
            "attempts": attempt_count,
            "status": "blocked" if attempt_count >= RATE_LIMIT_ATTEMPTS else "active"
        })
    
    # Trier par nombre de tentatives
    stats["top_ips"] = sorted(stats["top_ips"], key=lambda x: x["attempts"], reverse=True)[:20]
    
    return stats

@api_router.delete("/admin/illegal-content-logs/clear")
async def clear_illegal_content_logs(older_than_days: int = 30, current_admin: dict = Depends(get_current_admin)):
    """
    Supprime les anciens logs (par defaut > 30 jours)
    ADMIN ONLY - A proteger avec authentification en production
    """
    try:
        cutoff_date = datetime.now(timezone.utc) - timedelta(days=older_than_days)
        
        result = await db.security_illegal_logs.delete_many({
            "timestamp": {"$lt": cutoff_date}
        })
        
        return {
            "success": True,
            "deleted_count": result.deleted_count,
            "older_than_days": older_than_days
        }
        
    except Exception as e:
        logging.error(f"Erreur suppression logs: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de la suppression des logs: {str(e)}")



# ============================================
# GENERATION D'IMAGES IA (Hugging Face - GRATUIT)
# ============================================

import requests
import base64

HUGGINGFACE_API_KEY = os.environ.get('HUGGINGFACE_API_KEY', '')
# Nouveau modèle HF qui fonctionne (FLUX.1 - rapide et gratuit)
HUGGINGFACE_API_URL = "https://api-inference.huggingface.co/models/black-forest-labs/FLUX.1-schnell"

# Système de limite pour rester gratuit (1000 images/mois max)
IMAGE_GENERATION_LIMIT = int(os.getenv("IMAGE_GENERATION_LIMIT", "800"))  # Marge de sécurité
image_counter = defaultdict(int)
counter_lock = Lock()

def check_image_limit() -> bool:
    """Vérifie si la limite mensuelle d'images n'est pas atteinte"""
    with counter_lock:
        current_month = datetime.now(timezone.utc).strftime("%Y-%m")
        count = image_counter[current_month]
        
        if count >= IMAGE_GENERATION_LIMIT:
            logging.warning(f"⚠️ Limite mensuelle d'images atteinte: {count}/{IMAGE_GENERATION_LIMIT}")
            return False
        
        image_counter[current_month] += 1
        logging.info(f"📊 Images générées ce mois: {count + 1}/{IMAGE_GENERATION_LIMIT}")
        return True

def get_remaining_images() -> dict:
    """Retourne le nombre d'images restantes ce mois"""
    with counter_lock:
        current_month = datetime.now(timezone.utc).strftime("%Y-%m")
        used = image_counter[current_month]
        remaining = IMAGE_GENERATION_LIMIT - used
        return {
            "used": used,
            "limit": IMAGE_GENERATION_LIMIT,
            "remaining": max(0, remaining),
            "month": current_month
        }

class ImageGenerationRequest(BaseModel):
    prompt: str
    negative_prompt: Optional[str] = "blurry, bad quality, distorted"

class ImageGenerationResponse(BaseModel):
    success: bool
    image_url: Optional[str] = None
    image_base64: Optional[str] = None
    error: Optional[str] = None

async def generate_image_huggingface(prompt: str, negative_prompt: str = "") -> bytes:
    """Génère une image avec Hugging Face FLUX.1-schnell (GRATUIT avec limite)"""
    # Nouveau modèle: FLUX.1-schnell (rapide et gratuit)
    # Limite: 800 images/mois pour rester dans le free tier HF
    
    if not HUGGINGFACE_API_KEY:
        raise HTTPException(
            status_code=503,
            detail="❌ Génération d'images désactivée. Clé API Hugging Face manquante. Ajoutez HUGGINGFACE_API_KEY dans les variables d'environnement."
        )
    
    # Vérifier la limite mensuelle
    if not check_image_limit():
        stats = get_remaining_images()
        raise HTTPException(
            status_code=429,
            detail=f"⚠️ Limite mensuelle atteinte: {stats['used']}/{stats['limit']} images utilisées. Le compteur se réinitialise le 1er du mois prochain. Utilisez les graphiques mathématiques à la place!"
        )
    
    headers = {"Authorization": f"Bearer {HUGGINGFACE_API_KEY}"}
    
    payload = {
        "inputs": prompt,
        "parameters": {
            "num_inference_steps": 4,  # FLUX.1-schnell est optimisé pour 4 steps
            "guidance_scale": 0.0,      # FLUX.1-schnell n'utilise pas de guidance
        }
    }
    
    try:
        response = requests.post(HUGGINGFACE_API_URL, headers=headers, json=payload, timeout=60)
        response.raise_for_status()
        return response.content
    except requests.exceptions.HTTPError as e:
        if e.response.status_code == 503:
            # Le modèle est en train de charger
            logging.warning("Modèle en chargement, réessai dans 20s...")
            await asyncio.sleep(20)
            response = requests.post(HUGGINGFACE_API_URL, headers=headers, json=payload, timeout=60)
            response.raise_for_status()
            return response.content
        else:
            raise
    except Exception as e:
        logging.error(f"Erreur génération image HF: {e}")
        raise HTTPException(
            status_code=500,
            detail=f"❌ Erreur lors de la génération de l'image: {str(e)}"
        )

@api_router.get("/image-quota")
async def get_image_quota():
    """Retourne le quota d'images restant pour ce mois (FREE TIER)"""
    stats = get_remaining_images()
    return {
        "quota": {
            "used": stats['used'],
            "limit": stats['limit'],
            "remaining": stats['remaining'],
            "month": stats['month']
        },
        "message": f"📊 {stats['remaining']} images restantes ce mois (gratuit)"
    }

@api_router.post("/generate-image", response_model=ImageGenerationResponse)
async def generate_image(request: ImageGenerationRequest):
    """
    Genere une image IA a partir de prompt
    Utilise Hugging Face Stable Diffusion (100% GRATUIT)
    """
    try:
        logging.info(f"Generation image: {request.prompt}")
        
        # Generer image
        image_bytes = await generate_image_huggingface(request.prompt, request.negative_prompt)
        
        # Convertir en base64 pour affichage
        image_base64 = base64.b64encode(image_bytes).decode('utf-8')
        
        return ImageGenerationResponse(
            success=True,
            image_base64=image_base64,
            error=None
        )
        
    except Exception as e:
        logging.error(f"Erreur generation image: {e}")
        return ImageGenerationResponse(
            success=False,
            image_base64=None,
            error=str(e)
        )


# ==================== ROUTES AUTHENTIFICATION ET LICENCES ====================

@api_router.post("/auth/signup")
async def signup(user_data: UserSignup):
    """Inscription d'un nouvel utilisateur avec clé de licence"""
    try:
        # Vérifier si l'email existe déjà
        existing_user = await db.users.find_one({"email": user_data.email})
        if existing_user:
            raise HTTPException(status_code=400, detail="Cet email est déjà utilisé")
        
        # Vérifier la validité de la licence
        license_doc = await db.licenses.find_one({"license_key": user_data.license_key})
        if not license_doc:
            raise HTTPException(status_code=400, detail="Clé de licence invalide")
        
        if not license_doc.get("is_active", True):
            raise HTTPException(status_code=400, detail="Cette licence est désactivée")
        
        # Vérifier la date d'expiration
        expiry_date_str = license_doc["expiry_date"]
        # Parser la date et ajouter timezone UTC si elle n'en a pas
        if 'T' in expiry_date_str:
            expiry_date = datetime.fromisoformat(expiry_date_str)
        else:
            # Format YYYY-MM-DD seulement, ajouter l'heure de fin de journée
            expiry_date = datetime.strptime(expiry_date_str, "%Y-%m-%d").replace(hour=23, minute=59, second=59, tzinfo=timezone.utc)
        
        # S'assurer que expiry_date a une timezone
        if expiry_date.tzinfo is None:
            expiry_date = expiry_date.replace(tzinfo=timezone.utc)
        
        if expiry_date < datetime.now(timezone.utc):
            raise HTTPException(status_code=400, detail="Cette licence a expiré")
        
        # Vérifier le nombre maximum d'utilisateurs
        current_users = await db.users.count_documents({"license_key": user_data.license_key})
        if current_users >= license_doc["max_users"]:
            raise HTTPException(status_code=400, detail=f"Limite d'utilisateurs atteinte pour cette licence ({license_doc['max_users']} max)")
        
        # Hasher le mot de passe
        hashed_password = hash_password(user_data.password)
        
        # Créer l'utilisateur
        user_doc = {
            "id": str(uuid.uuid4()),
            "full_name": user_data.full_name,
            "email": user_data.email,
            "password": hashed_password,
            "license_key": user_data.license_key,
            "organization": license_doc["organization_name"],
            "is_active": True,
            "created_at": datetime.now(timezone.utc).isoformat(),
            "last_login": None,
            "message_count": 0
        }
        
        await db.users.insert_one(user_doc)
        
        # Créer un token JWT
        token = create_access_token({"user_id": user_doc["id"], "email": user_doc["email"]})
        
        return {
            "success": True,
            "message": "Compte créé avec succès",
            "token": token,
            "user": {
                "id": user_doc["id"],
                "full_name": user_doc["full_name"],
                "email": user_doc["email"],
                "organization": user_doc["organization"]
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur inscription: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de l'inscription: {str(e)}")

@api_router.post("/auth/login")
async def login(credentials: UserLogin):
    """Connexion d'un utilisateur"""
    try:
        # Trouver l'utilisateur
        user = await db.users.find_one({"email": credentials.email})
        if not user:
            raise HTTPException(status_code=401, detail="Email ou mot de passe incorrect")
        
        # Vérifier le mot de passe
        if not verify_password(credentials.password, user["password"]):
            raise HTTPException(status_code=401, detail="Email ou mot de passe incorrect")
        
        # Vérifier si l'utilisateur est actif
        if not user.get("is_active", True):
            raise HTTPException(status_code=403, detail="Votre compte a été désactivé")
        
        # Vérifier la licence
        license_doc = await db.licenses.find_one({"license_key": user["license_key"]})
        if not license_doc or not license_doc.get("is_active", True):
            raise HTTPException(status_code=403, detail="Votre licence est inactive")
        
        # Vérifier l'expiration
        expiry_date_str = license_doc["expiry_date"]
        # Parser la date et ajouter timezone UTC si elle n'en a pas
        if 'T' in expiry_date_str:
            expiry_date = datetime.fromisoformat(expiry_date_str)
        else:
            expiry_date = datetime.strptime(expiry_date_str, "%Y-%m-%d").replace(hour=23, minute=59, second=59, tzinfo=timezone.utc)
        
        if expiry_date.tzinfo is None:
            expiry_date = expiry_date.replace(tzinfo=timezone.utc)
        
        if expiry_date < datetime.now(timezone.utc):
            raise HTTPException(status_code=403, detail="Votre licence a expiré")
        
        # Mettre à jour la dernière connexion et incrémenter le compteur
        await db.users.update_one(
            {"id": user["id"]},
            {
                "$set": {"last_login": datetime.now(timezone.utc).isoformat()},
                "$inc": {"login_count": 1}
            }
        )
        
        # Enregistrer la connexion dans les logs
        await db.login_logs.insert_one({
            "user_id": user["id"],
            "email": user["email"],
            "license_key": user["license_key"],
            "timestamp": datetime.now(timezone.utc).isoformat(),
            "year": datetime.now(timezone.utc).year,
            "month": datetime.now(timezone.utc).month
        })
        
        # Créer un token JWT
        token = create_access_token({"user_id": user["id"], "email": user["email"]})
        
        return {
            "success": True,
            "token": token,
            "user": {
                "id": user["id"],
                "full_name": user["full_name"],
                "email": user["email"],
                "organization": user.get("organization", "")
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur connexion: {e}")
        raise HTTPException(status_code=500, detail=f"Erreur lors de la connexion: {str(e)}")

# ==================== ROUTES ADMIN - GESTION DES LICENCES ====================

@api_router.post("/admin/licenses")
async def create_license(license_data: LicenseCreate, admin = Depends(get_current_admin)):
    """Créer une nouvelle licence (Admin seulement)"""
    try:
        # Vérifier si la clé existe déjà
        existing = await db.licenses.find_one({"license_key": license_data.license_key})
        if existing:
            raise HTTPException(status_code=400, detail="Cette clé de licence existe déjà")
        
        license_doc = {
            "id": str(uuid.uuid4()),
            "organization_name": license_data.organization_name,
            "license_key": license_data.license_key,
            "max_users": license_data.max_users,
            "expiry_date": license_data.expiry_date,
            "is_active": True,
            "notes": license_data.notes or "",
            "created_at": datetime.now(timezone.utc).isoformat(),
            "created_by": admin.get("email", admin.get("username", "admin"))
        }
        
        await db.licenses.insert_one(license_doc)
        
        return {
            "success": True,
            "message": "Licence créée avec succès",
            "license": license_doc
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur création licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/admin/licenses")
async def get_all_licenses(admin = Depends(get_current_admin)):
    """Obtenir toutes les licences avec statistiques (Admin seulement)"""
    try:
        licenses = await db.licenses.find({}, {"_id": 0}).to_list(1000)
        
        # Ajouter les stats pour chaque licence
        for license in licenses:
            user_count = await db.users.count_documents({"license_key": license["license_key"]})
            license["current_users"] = user_count
            
            # Parser la date d'expiration avec gestion des timezones
            expiry_date_str = license["expiry_date"]
            if 'T' in expiry_date_str:
                expiry_date = datetime.fromisoformat(expiry_date_str)
            else:
                # Format YYYY-MM-DD seulement, ajouter l'heure de fin de journée
                expiry_date = datetime.strptime(expiry_date_str, "%Y-%m-%d").replace(hour=23, minute=59, second=59, tzinfo=timezone.utc)
            
            # S'assurer que expiry_date a une timezone
            if expiry_date.tzinfo is None:
                expiry_date = expiry_date.replace(tzinfo=timezone.utc)
            
            license["is_expired"] = expiry_date < datetime.now(timezone.utc)
        
        return {
            "success": True,
            "licenses": licenses,
            "total_licenses": len(licenses)
        }
        
    except Exception as e:
        logging.error(f"Erreur récupération licences: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.put("/admin/licenses/{license_key}")
async def update_license(license_key: str, update_data: LicenseUpdate, admin = Depends(get_current_admin)):
    """Mettre à jour une licence (Admin seulement)"""
    try:
        license_doc = await db.licenses.find_one({"license_key": license_key})
        if not license_doc:
            raise HTTPException(status_code=404, detail="Licence non trouvée")
        
        update_fields = {}
        new_key = None
        
        if update_data.max_users is not None:
            update_fields["max_users"] = update_data.max_users
        if update_data.expiry_date is not None:
            update_fields["expiry_date"] = update_data.expiry_date
        if update_data.is_active is not None:
            update_fields["is_active"] = update_data.is_active
        if update_data.notes is not None:
            update_fields["notes"] = update_data.notes
        if update_data.license_key is not None and update_data.license_key != license_key:
            # Vérifier que la nouvelle clé n'existe pas déjà
            existing = await db.licenses.find_one({"license_key": update_data.license_key})
            if existing:
                raise HTTPException(status_code=400, detail="Cette clé de licence existe déjà")
            new_key = update_data.license_key
            update_fields["license_key"] = new_key
        
        update_fields["updated_at"] = datetime.now(timezone.utc).isoformat()
        
        await db.licenses.update_one(
            {"license_key": license_key},
            {"$set": update_fields}
        )
        
        # Si la clé a changé, mettre à jour tous les utilisateurs
        if new_key:
            await db.users.update_many(
                {"license_key": license_key},
                {"$set": {"license_key": new_key}}
            )
        
        return {
            "success": True,
            "message": "Licence mise à jour avec succès" + (f" (nouvelle clé: {new_key})" if new_key else "")
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur mise à jour licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/admin/licenses/{license_key}/users")
async def get_license_users(license_key: str, admin = Depends(get_current_admin)):
    """Obtenir tous les utilisateurs d'une licence (Admin seulement)"""
    try:
        users = await db.users.find(
            {"license_key": license_key},
            {"_id": 0, "password": 0}
        ).to_list(1000)
        
        return {
            "success": True,
            "users": users,
            "total_users": len(users)
        }
        
    except Exception as e:
        logging.error(f"Erreur récupération utilisateurs: {e}")
        raise HTTPException(status_code=500, detail=str(e))

class TransferUsersRequest(BaseModel):
    user_emails: list[str]
    target_license_key: str

@api_router.post("/admin/transfer-users")
async def transfer_users(data: TransferUsersRequest, admin = Depends(get_current_admin)):
    """Transférer un ou plusieurs utilisateurs vers une autre licence (Super Admin seulement)"""
    try:
        # Vérifier que la licence cible existe et est active
        target_license = await db.licenses.find_one({"license_key": data.target_license_key})
        if not target_license:
            raise HTTPException(status_code=404, detail="Licence cible introuvable")
        
        if not target_license.get("is_active", False):
            raise HTTPException(status_code=400, detail="La licence cible est inactive")
        
        # Vérifier le nombre d'utilisateurs actuels sur la licence cible
        current_users_count = await db.users.count_documents({"license_key": data.target_license_key})
        max_users = target_license.get("max_users", 10)
        
        # Calculer combien d'utilisateurs on peut transférer
        available_slots = max_users - current_users_count
        if available_slots < len(data.user_emails):
            raise HTTPException(
                status_code=400, 
                detail=f"Espace insuffisant sur la licence cible. Places disponibles: {available_slots}, demandé: {len(data.user_emails)}"
            )
        
        # Effectuer le transfert
        transferred = []
        errors = []
        
        for email in data.user_emails:
            user = await db.users.find_one({"email": email.lower()})
            if not user:
                errors.append(f"Utilisateur {email} introuvable")
                continue
            
            old_license = user.get("license_key", "N/A")
            
            # Mettre à jour la licence de l'utilisateur
            result = await db.users.update_one(
                {"email": email.lower()},
                {
                    "$set": {
                        "license_key": data.target_license_key,
                        "transferred_at": datetime.now(timezone.utc).isoformat(),
                        "transferred_from": old_license
                    }
                }
            )
            
            if result.modified_count > 0:
                transferred.append({
                    "email": email,
                    "from_license": old_license,
                    "to_license": data.target_license_key
                })
                logging.info(f"Utilisateur {email} transféré de {old_license} vers {data.target_license_key}")
            else:
                errors.append(f"Échec du transfert pour {email}")
        
        return {
            "success": True,
            "message": f"{len(transferred)} utilisateur(s) transféré(s) avec succès",
            "transferred": transferred,
            "errors": errors if errors else None
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur transfert utilisateurs: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== ROUTES ADMIN - GESTION DES MOTS BLOQUÉS ====================

@api_router.post("/admin/blocked-words")
async def add_blocked_word(word_data: BlockedWordCreate, admin = Depends(get_current_admin)):
    """Ajouter un mot à la liste des mots bloqués (Admin seulement)"""
    try:
        # Vérifier si le mot existe déjà
        existing = await db.blocked_words.find_one({"word": word_data.word.lower()})
        if existing:
            raise HTTPException(status_code=400, detail="Ce mot est déjà dans la liste")
        
        word_doc = {
            "id": str(uuid.uuid4()),
            "word": word_data.word.lower(),
            "category": word_data.category,
            "severity": word_data.severity,
            "is_exception": word_data.is_exception,
            "is_active": True,
            "created_at": datetime.now(timezone.utc).isoformat(),
            "created_by": admin.get("email", admin.get("username", "admin"))
        }
        
        await db.blocked_words.insert_one(word_doc)
        
        return {
            "success": True,
            "message": "Mot ajouté avec succès",
            "word": word_doc
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur ajout mot bloqué: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/admin/blocked-words")
async def get_blocked_words(admin = Depends(get_current_admin)):
    """Obtenir tous les mots bloqués (Admin seulement)"""
    try:
        words = await db.blocked_words.find({}, {"_id": 0}).to_list(5000)
        
        # Grouper par catégorie
        by_category = {}
        for word in words:
            category = word["category"]
            if category not in by_category:
                by_category[category] = []
            by_category[category].append(word)
        
        return {
            "success": True,
            "words": words,
            "by_category": by_category,
            "total_words": len(words)
        }
        
    except Exception as e:
        logging.error(f"Erreur récupération mots bloqués: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.put("/admin/blocked-words/{word_id}")
async def update_blocked_word(word_id: str, update_data: BlockedWordUpdate, admin = Depends(get_current_admin)):
    """Mettre à jour un mot bloqué (Admin seulement)"""
    try:
        word_doc = await db.blocked_words.find_one({"id": word_id})
        if not word_doc:
            raise HTTPException(status_code=404, detail="Mot non trouvé")
        
        update_fields = {}
        if update_data.category is not None:
            update_fields["category"] = update_data.category
        if update_data.severity is not None:
            update_fields["severity"] = update_data.severity
        if update_data.is_exception is not None:
            update_fields["is_exception"] = update_data.is_exception
        if update_data.is_active is not None:
            update_fields["is_active"] = update_data.is_active
        
        update_fields["updated_at"] = datetime.now(timezone.utc).isoformat()
        
        await db.blocked_words.update_one(
            {"id": word_id},
            {"$set": update_fields}
        )
        
        return {
            "success": True,
            "message": "Mot mis à jour avec succès"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur mise à jour mot bloqué: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.delete("/admin/blocked-words/{word_id}")
async def delete_blocked_word(word_id: str, admin = Depends(get_current_admin)):
    """Supprimer un mot de la liste (Admin seulement)"""
    try:
        result = await db.blocked_words.delete_one({"id": word_id})
        if result.deleted_count == 0:
            raise HTTPException(status_code=404, detail="Mot non trouvé")
        
        return {
            "success": True,
            "message": "Mot supprimé avec succès"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur suppression mot bloqué: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== SYSTÈME ADMIN MULTI-NIVEAUX ====================

async def get_current_license_admin(request: Request):
    """Vérifie si l'utilisateur est admin d'une licence ou super admin"""
    token = request.headers.get("Authorization", "").replace("Bearer ", "")
    if not token:
        raise HTTPException(status_code=401, detail="Token manquant")
    
    try:
        # D'abord vérifier si c'est un super admin
        payload = jwt.decode(token, os.environ.get("JWT_SECRET_KEY", ""), algorithms=["HS256"])
        
        # Si c'est un super admin (token admin)
        if "username" in payload:
            return {"is_super_admin": True, "username": payload["username"]}
        
        # Sinon, c'est un token utilisateur - vérifier s'il est license_admin
        user_id = payload.get("user_id")
        if not user_id:
            raise HTTPException(status_code=401, detail="Token invalide")
        
        user = await db.users.find_one({"id": user_id})
        if not user:
            raise HTTPException(status_code=401, detail="Utilisateur non trouvé")
        
        # Vérifier si l'utilisateur est admin de sa licence
        license_doc = await db.licenses.find_one({"license_key": user["license_key"]})
        if not license_doc:
            raise HTTPException(status_code=403, detail="Licence non trouvée")
        
        if license_doc.get("license_admin_id") == user_id:
            return {
                "is_super_admin": False,
                "is_license_admin": True,
                "user_id": user_id,
                "license_key": user["license_key"],
                "email": user["email"]
            }
        
        raise HTTPException(status_code=403, detail="Vous n'êtes pas administrateur de cette licence")
        
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expiré")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Token invalide")

@api_router.post("/admin/licenses/{license_key}/set-admin")
async def set_license_admin(license_key: str, data: LicenseAdminUpdate, admin = Depends(get_current_admin)):
    """Désigner un admin de licence (Super Admin seulement)"""
    try:
        license_doc = await db.licenses.find_one({"license_key": license_key})
        if not license_doc:
            raise HTTPException(status_code=404, detail="Licence non trouvée")
        
        if data.license_admin_email:
            # Trouver l'utilisateur par email
            user = await db.users.find_one({"email": data.license_admin_email, "license_key": license_key})
            if not user:
                raise HTTPException(status_code=404, detail="Utilisateur non trouvé dans cette licence")
            
            await db.licenses.update_one(
                {"license_key": license_key},
                {"$set": {
                    "license_admin_id": user["id"],
                    "license_admin_email": data.license_admin_email,
                    "updated_at": datetime.now(timezone.utc).isoformat()
                }}
            )
            
            return {
                "success": True,
                "message": f"Admin de licence désigné: {data.license_admin_email}"
            }
        else:
            # Retirer l'admin de licence
            await db.licenses.update_one(
                {"license_key": license_key},
                {"$unset": {"license_admin_id": "", "license_admin_email": ""},
                 "$set": {"updated_at": datetime.now(timezone.utc).isoformat()}}
            )
            
            return {
                "success": True,
                "message": "Admin de licence retiré"
            }
            
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur désignation admin licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.put("/license-admin/update-key")
async def update_license_key(data: LicenseKeyUpdate, admin = Depends(get_current_license_admin)):
    """Modifier la clé de licence (Admin de licence ou Super Admin)"""
    try:
        # Déterminer la licence à modifier
        if admin.get("is_super_admin"):
            raise HTTPException(status_code=400, detail="Super admin: utilisez /admin/licenses/{license_key} pour modifier")
        
        license_key = admin.get("license_key")
        
        # Vérifier que la nouvelle clé n'existe pas déjà
        existing = await db.licenses.find_one({"license_key": data.new_license_key})
        if existing:
            raise HTTPException(status_code=400, detail="Cette clé de licence existe déjà")
        
        # Mettre à jour la clé dans la licence
        await db.licenses.update_one(
            {"license_key": license_key},
            {"$set": {
                "license_key": data.new_license_key,
                "updated_at": datetime.now(timezone.utc).isoformat()
            }}
        )
        
        # Mettre à jour la clé pour tous les utilisateurs de cette licence
        await db.users.update_many(
            {"license_key": license_key},
            {"$set": {"license_key": data.new_license_key}}
        )
        
        return {
            "success": True,
            "message": f"Clé de licence modifiée: {license_key} → {data.new_license_key}"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur modification clé licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/license-admin/my-license")
async def get_my_license(admin = Depends(get_current_license_admin)):
    """Obtenir les infos de sa propre licence (Admin de licence)"""
    try:
        if admin.get("is_super_admin"):
            raise HTTPException(status_code=400, detail="Super admin: utilisez /admin/licenses")
        
        license_key = admin.get("license_key")
        license_doc = await db.licenses.find_one({"license_key": license_key}, {"_id": 0})
        
        if not license_doc:
            raise HTTPException(status_code=404, detail="Licence non trouvée")
        
        # Compter les utilisateurs
        user_count = await db.users.count_documents({"license_key": license_key})
        license_doc["current_users"] = user_count
        
        # Vérifier expiration
        expiry_date_str = license_doc["expiry_date"]
        if 'T' in expiry_date_str:
            expiry_date = datetime.fromisoformat(expiry_date_str)
        else:
            expiry_date = datetime.strptime(expiry_date_str, "%Y-%m-%d").replace(hour=23, minute=59, second=59, tzinfo=timezone.utc)
        
        if expiry_date.tzinfo is None:
            expiry_date = expiry_date.replace(tzinfo=timezone.utc)
        
        license_doc["is_expired"] = expiry_date < datetime.now(timezone.utc)
        
        return {
            "success": True,
            "license": license_doc
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur récupération licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/license-admin/users")
async def get_my_license_users(admin = Depends(get_current_license_admin)):
    """Obtenir les utilisateurs de sa licence (Admin de licence)"""
    try:
        if admin.get("is_super_admin"):
            raise HTTPException(status_code=400, detail="Super admin: utilisez /admin/licenses/{license_key}/users")
        
        license_key = admin.get("license_key")
        users = await db.users.find(
            {"license_key": license_key},
            {"_id": 0, "password": 0}
        ).to_list(1000)
        
        return {
            "success": True,
            "users": users,
            "total_users": len(users)
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur récupération utilisateurs: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/license-admin/users")
async def add_user_to_license(user_data: UserCreate, admin = Depends(get_current_license_admin)):
    """Ajouter un utilisateur à sa licence (Admin de licence)"""
    try:
        if admin.get("is_super_admin"):
            raise HTTPException(status_code=400, detail="Super admin: les utilisateurs s'inscrivent via /signup")
        
        license_key = admin.get("license_key")
        
        # Vérifier si l'email existe déjà
        existing = await db.users.find_one({"email": user_data.email})
        if existing:
            raise HTTPException(status_code=400, detail="Cet email est déjà utilisé")
        
        # Vérifier la limite d'utilisateurs
        license_doc = await db.licenses.find_one({"license_key": license_key})
        current_users = await db.users.count_documents({"license_key": license_key})
        
        if current_users >= license_doc["max_users"]:
            raise HTTPException(
                status_code=400, 
                detail=f"Limite d'utilisateurs atteinte ({license_doc['max_users']} max). Contactez le super administrateur pour augmenter la limite."
            )
        
        # Créer l'utilisateur
        hashed_password = hash_password(user_data.password)
        
        user_doc = {
            "id": str(uuid.uuid4()),
            "full_name": user_data.full_name,
            "email": user_data.email,
            "password": hashed_password,
            "license_key": license_key,
            "organization": license_doc["organization_name"],
            "is_active": True,
            "created_at": datetime.now(timezone.utc).isoformat(),
            "created_by": admin.get("email", "license_admin")
        }
        
        await db.users.insert_one(user_doc)
        
        return {
            "success": True,
            "message": f"Utilisateur {user_data.email} ajouté avec succès",
            "user": {
                "id": user_doc["id"],
                "full_name": user_doc["full_name"],
                "email": user_doc["email"]
            }
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur ajout utilisateur: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.delete("/license-admin/users/{user_id}")
async def remove_user_from_license(user_id: str, admin = Depends(get_current_license_admin)):
    """Retirer un utilisateur de sa licence (Admin de licence)"""
    try:
        if admin.get("is_super_admin"):
            # Super admin peut supprimer n'importe quel utilisateur
            result = await db.users.delete_one({"id": user_id})
            if result.deleted_count == 0:
                raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
            return {"success": True, "message": "Utilisateur supprimé"}
        
        license_key = admin.get("license_key")
        
        # Vérifier que l'utilisateur appartient à cette licence
        user = await db.users.find_one({"id": user_id, "license_key": license_key})
        if not user:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé dans votre licence")
        
        # Ne pas permettre de se supprimer soi-même
        if user_id == admin.get("user_id"):
            raise HTTPException(status_code=400, detail="Vous ne pouvez pas vous supprimer vous-même")
        
        await db.users.delete_one({"id": user_id})
        
        return {
            "success": True,
            "message": f"Utilisateur {user['email']} retiré de la licence"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur suppression utilisateur: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== CONVERSATIONS CLOUD (SYNCHRONISÉES) ====================

async def get_current_user_from_token(request: Request):
    """Extrait l'utilisateur du token JWT"""
    token = request.headers.get("Authorization", "").replace("Bearer ", "")
    if not token:
        raise HTTPException(status_code=401, detail="Token manquant")
    
    try:
        payload = jwt.decode(token, os.environ.get("JWT_SECRET_KEY", ""), algorithms=["HS256"])
        user_id = payload.get("user_id")
        if not user_id:
            raise HTTPException(status_code=401, detail="Token invalide")
        return {"user_id": user_id}
    except jwt.ExpiredSignatureError:
        raise HTTPException(status_code=401, detail="Token expiré")
    except jwt.InvalidTokenError:
        raise HTTPException(status_code=401, detail="Token invalide")

@api_router.get("/conversations")
async def get_user_conversations(user = Depends(get_current_user_from_token)):
    """Récupérer toutes les conversations de l'utilisateur"""
    try:
        conversations = await db.conversations.find(
            {"user_id": user["user_id"]},
            {"_id": 0}
        ).sort("updated_at", -1).to_list(100)
        
        return {
            "success": True,
            "conversations": conversations
        }
    except Exception as e:
        logging.error(f"Erreur récupération conversations: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/conversations")
async def create_conversation(data: ConversationCreate, user = Depends(get_current_user_from_token)):
    """Créer une nouvelle conversation"""
    try:
        conv_id = f"conv_{uuid.uuid4().hex[:12]}"
        
        conversation = {
            "id": conv_id,
            "user_id": user["user_id"],
            "title": data.title or "Nouvelle conversation",
            "messages": [],
            "is_pinned": False,
            "is_favorite": False,
            "tags": [],
            "created_at": datetime.now(timezone.utc).isoformat(),
            "updated_at": datetime.now(timezone.utc).isoformat()
        }
        
        await db.conversations.insert_one(conversation)
        
        return {
            "success": True,
            "conversation_id": conv_id,
            "conversation": {k: v for k, v in conversation.items() if k != "_id"}
        }
    except Exception as e:
        logging.error(f"Erreur création conversation: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/conversations/{conversation_id}")
async def get_conversation(conversation_id: str, user = Depends(get_current_user_from_token)):
    """Récupérer une conversation spécifique avec ses messages"""
    try:
        conversation = await db.conversations.find_one(
            {"id": conversation_id, "user_id": user["user_id"]},
            {"_id": 0}
        )
        
        if not conversation:
            raise HTTPException(status_code=404, detail="Conversation non trouvée")
        
        return {
            "success": True,
            "conversation": conversation
        }
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur récupération conversation: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/conversations/{conversation_id}/messages")
async def add_message_to_conversation(
    conversation_id: str, 
    message: ConversationMessage, 
    user = Depends(get_current_user_from_token)
):
    """Ajouter un message à une conversation"""
    try:
        # Vérifier que la conversation existe et appartient à l'utilisateur
        conversation = await db.conversations.find_one(
            {"id": conversation_id, "user_id": user["user_id"]}
        )
        
        if not conversation:
            raise HTTPException(status_code=404, detail="Conversation non trouvée")
        
        # Créer le message
        new_message = {
            "role": message.role,
            "content": message.content,
            "timestamp": message.timestamp or datetime.now(timezone.utc).isoformat()
        }
        
        if message.image_base64:
            new_message["image_base64"] = message.image_base64
        
        # Ajouter le message et mettre à jour la date
        await db.conversations.update_one(
            {"id": conversation_id},
            {
                "$push": {"messages": new_message},
                "$set": {"updated_at": datetime.now(timezone.utc).isoformat()}
            }
        )
        
        return {
            "success": True,
            "message": "Message ajouté"
        }
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur ajout message: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.put("/conversations/{conversation_id}")
async def update_conversation(
    conversation_id: str, 
    data: ConversationUpdate, 
    user = Depends(get_current_user_from_token)
):
    """Mettre à jour une conversation (titre, épinglé, favoris, tags)"""
    try:
        update_fields = {"updated_at": datetime.now(timezone.utc).isoformat()}
        
        if data.title is not None:
            update_fields["title"] = data.title
        if data.is_pinned is not None:
            update_fields["is_pinned"] = data.is_pinned
        if data.is_favorite is not None:
            update_fields["is_favorite"] = data.is_favorite
        if data.tags is not None:
            update_fields["tags"] = data.tags
        
        result = await db.conversations.update_one(
            {"id": conversation_id, "user_id": user["user_id"]},
            {"$set": update_fields}
        )
        
        if result.matched_count == 0:
            raise HTTPException(status_code=404, detail="Conversation non trouvée")
        
        return {
            "success": True,
            "message": "Conversation mise à jour"
        }
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur mise à jour conversation: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.delete("/conversations/{conversation_id}")
async def delete_conversation(conversation_id: str, user = Depends(get_current_user_from_token)):
    """Supprimer une conversation"""
    try:
        result = await db.conversations.delete_one(
            {"id": conversation_id, "user_id": user["user_id"]}
        )
        
        if result.deleted_count == 0:
            raise HTTPException(status_code=404, detail="Conversation non trouvée")
        
        return {
            "success": True,
            "message": "Conversation supprimée"
        }
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur suppression conversation: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== RÉINITIALISATION MOT DE PASSE ====================

def send_reset_email_sync(email: str, reset_token: str, user_name: str, smtp_config: dict) -> bool:
    """Fonction synchrone pour envoyer l'email via API HTTP smtp2go (plus fiable que SMTP sur cloud)"""
    import requests
    
    try:
        api_key = smtp_config.get("api_key", "")
        smtp_from = smtp_config.get("from_email", "noreply@champagneur.qc.ca")
        frontend_url = smtp_config.get("frontend_url", "")
        
        # Log de debug détaillé
        logging.info(f"📧 smtp2go API - From: {smtp_from}, API Key: {'✓' if api_key else '✗'}")
        logging.info(f"📧 Frontend URL: {frontend_url}")
        
        if not api_key:
            logging.error("❌ SMTP2GO_API_KEY manquant - impossible d'envoyer l'email")
            return False
        
        if not frontend_url:
            logging.warning("⚠️ FRONTEND_URL manquant - le lien de réinitialisation pourrait être invalide")
        
        reset_url = f"{frontend_url}?reset_token={reset_token}"
        logging.info(f"📧 Reset URL généré: {reset_url[:60]}...")
        
        # Contenu de l'email
        text_content = f"""Bonjour {user_name},

Vous avez demandé une réinitialisation de votre mot de passe pour votre compte Étienne.

Cliquez sur le lien ci-dessous pour créer un nouveau mot de passe :
{reset_url}

Ce lien expirera dans 1 heure.

Si vous n'avez pas demandé cette réinitialisation, ignorez simplement cet email.

Cordialement,
L'équipe Étienne
Assistant IA pour enseignants
"""
        
        html_content = f"""
<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: Arial, sans-serif; line-height: 1.6; color: #333; }}
        .container {{ max-width: 600px; margin: 0 auto; padding: 20px; }}
        .header {{ background: linear-gradient(135deg, #f97316, #ea580c); color: white; padding: 20px; text-align: center; border-radius: 10px 10px 0 0; }}
        .content {{ background: #f9fafb; padding: 30px; border: 1px solid #e5e7eb; }}
        .button {{ display: inline-block; background: #f97316; color: white; padding: 12px 30px; text-decoration: none; border-radius: 5px; margin: 20px 0; }}
        .footer {{ background: #1f2937; color: #9ca3af; padding: 15px; text-align: center; font-size: 12px; border-radius: 0 0 10px 10px; }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>🎓 Étienne</h1>
            <p>Assistant IA pour enseignants</p>
        </div>
        <div class="content">
            <h2>Réinitialisation de mot de passe</h2>
            <p>Bonjour <strong>{user_name}</strong>,</p>
            <p>Vous avez demandé une réinitialisation de votre mot de passe.</p>
            <p>Cliquez sur le bouton ci-dessous pour créer un nouveau mot de passe :</p>
            <p style="text-align: center;">
                <a href="{reset_url}" class="button">Réinitialiser mon mot de passe</a>
            </p>
            <p><small>Ce lien expirera dans <strong>1 heure</strong>.</small></p>
            <p><small>Si vous n'avez pas demandé cette réinitialisation, ignorez simplement cet email.</small></p>
        </div>
        <div class="footer">
            <p>© 2025 Étienne - Collège Champagneur</p>
        </div>
    </div>
</body>
</html>
"""
        
        # Appel API smtp2go
        logging.info(f"📧 Envoi via API smtp2go à {email}...")
        
        response = requests.post(
            "https://api.smtp2go.com/v3/email/send",
            json={
                "api_key": api_key,
                "to": [email],
                "sender": smtp_from,
                "subject": "Réinitialisation de votre mot de passe - Étienne",
                "text_body": text_content,
                "html_body": html_content
            },
            timeout=30
        )
        
        result = response.json()
        
        if response.status_code == 200 and result.get("data", {}).get("succeeded", 0) > 0:
            logging.info(f"✅ Email envoyé avec succès à {email} via API smtp2go")
            return True
        else:
            logging.error(f"❌ Échec API smtp2go: {result}")
            return False
        
    except requests.exceptions.Timeout:
        logging.error(f"❌ Timeout API smtp2go")
        return False
    except requests.exceptions.RequestException as e:
        logging.error(f"❌ Erreur réseau API smtp2go: {e}")
        return False
    except Exception as e:
        logging.error(f"❌ Erreur envoi email inattendue: {type(e).__name__}: {e}")
        return False

async def send_reset_email(email: str, reset_token: str, user_name: str) -> bool:
    """Envoyer l'email de réinitialisation via API HTTP smtp2go (asynchrone)"""
    try:
        smtp_config = {
            "api_key": os.environ.get("SMTP2GO_API_KEY", ""),
            "from_email": os.environ.get("SMTP_FROM_EMAIL", "noreply@champagneur.qc.ca"),
            "frontend_url": os.environ.get("FRONTEND_URL", "https://etienne-saas.preview.emergentagent.com")
        }
        
        if not smtp_config["api_key"]:
            logging.error("Configuration SMTP2GO manquante - variable SMTP2GO_API_KEY requise")
            return False
        
        # Exécuter l'envoi dans un thread séparé pour ne pas bloquer
        result = await asyncio.to_thread(
            send_reset_email_sync,
            email,
            reset_token,
            user_name,
            smtp_config
        )
        return result
        
    except Exception as e:
        logging.error(f"❌ Erreur async send_reset_email: {e}")
        return False

@api_router.post("/auth/forgot-password")
async def forgot_password(data: PasswordResetRequest):
    """Demander une réinitialisation de mot de passe"""
    try:
        # Vérifier si l'utilisateur existe
        user = await db.users.find_one({"email": data.email.lower()})
        
        # Toujours retourner succès pour ne pas révéler si l'email existe
        if not user:
            return {
                "success": True,
                "message": "Si cette adresse email est associée à un compte, vous recevrez un lien de réinitialisation."
            }
        
        # Générer un token unique
        reset_token = secrets.token_urlsafe(32)
        expiry = datetime.now(timezone.utc) + timedelta(hours=1)
        
        # Sauvegarder le token
        await db.password_resets.delete_many({"email": data.email.lower()})  # Supprimer anciens tokens
        await db.password_resets.insert_one({
            "email": data.email.lower(),
            "token": reset_token,
            "expires_at": expiry.isoformat(),
            "created_at": datetime.now(timezone.utc).isoformat()
        })
        
        # Envoyer l'email EN ARRIÈRE-PLAN (ne pas bloquer la réponse)
        # Utiliser asyncio.create_task pour fire-and-forget
        async def send_email_background():
            try:
                email_sent = await send_reset_email(data.email.lower(), reset_token, user.get("full_name", "Utilisateur"))
                if not email_sent:
                    logging.warning(f"Échec envoi email pour {data.email}")
                else:
                    logging.info(f"✅ Email de réinitialisation envoyé à {data.email}")
            except Exception as e:
                logging.error(f"Erreur envoi email background: {e}")
        
        # Lancer l'envoi en arrière-plan sans attendre
        asyncio.create_task(send_email_background())
        
        # Retourner immédiatement la réponse
        return {
            "success": True,
            "message": "Si cette adresse email est associée à un compte, vous recevrez un lien de réinitialisation."
        }
        
    except Exception as e:
        logging.error(f"Erreur forgot password: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de la demande de réinitialisation")

@api_router.post("/auth/reset-password")
async def reset_password(data: PasswordResetConfirm):
    """Réinitialiser le mot de passe avec le token"""
    try:
        # Vérifier le token
        reset_doc = await db.password_resets.find_one({"token": data.token})
        
        if not reset_doc:
            raise HTTPException(status_code=400, detail="Lien de réinitialisation invalide ou expiré")
        
        # Vérifier l'expiration
        expiry = datetime.fromisoformat(reset_doc["expires_at"].replace('Z', '+00:00'))
        if expiry.tzinfo is None:
            expiry = expiry.replace(tzinfo=timezone.utc)
        
        if datetime.now(timezone.utc) > expiry:
            await db.password_resets.delete_one({"token": data.token})
            raise HTTPException(status_code=400, detail="Lien de réinitialisation expiré")
        
        # Valider le nouveau mot de passe
        if len(data.new_password) < 6:
            raise HTTPException(status_code=400, detail="Le mot de passe doit contenir au moins 6 caractères")
        
        # Hasher le nouveau mot de passe
        hashed_password = hash_password(data.new_password)
        
        # Mettre à jour le mot de passe
        result = await db.users.update_one(
            {"email": reset_doc["email"]},
            {"$set": {
                "password": hashed_password,
                "updated_at": datetime.now(timezone.utc).isoformat()
            }}
        )
        
        if result.matched_count == 0:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
        
        # Supprimer le token utilisé
        await db.password_resets.delete_one({"token": data.token})
        
        return {
            "success": True,
            "message": "Mot de passe réinitialisé avec succès. Vous pouvez maintenant vous connecter."
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur reset password: {e}")
        raise HTTPException(status_code=500, detail="Erreur lors de la réinitialisation")

@api_router.get("/auth/verify-reset-token/{token}")
async def verify_reset_token(token: str):
    """Vérifier si un token de réinitialisation est valide"""
    try:
        reset_doc = await db.password_resets.find_one({"token": token})
        
        if not reset_doc:
            return {"valid": False, "message": "Token invalide"}
        
        expiry = datetime.fromisoformat(reset_doc["expires_at"].replace('Z', '+00:00'))
        if expiry.tzinfo is None:
            expiry = expiry.replace(tzinfo=timezone.utc)
        
        if datetime.now(timezone.utc) > expiry:
            return {"valid": False, "message": "Token expiré"}
        
        return {"valid": True, "email": reset_doc["email"]}
        
    except Exception as e:
        logging.error(f"Erreur vérification token: {e}")
        return {"valid": False, "message": "Erreur de vérification"}

# ==================== CHANGEMENT D'EMAIL ====================

@api_router.post("/auth/change-email")
async def user_change_own_email(data: UserSelfChangeEmail, request: Request):
    """Permet à un utilisateur de changer son propre email"""
    try:
        # Récupérer l'utilisateur actuel via le token
        auth_header = request.headers.get("Authorization")
        if not auth_header or not auth_header.startswith("Bearer "):
            raise HTTPException(status_code=401, detail="Non authentifié")
        
        token = auth_header.split(" ")[1]
        payload = jwt.decode(token, JWT_SECRET, algorithms=["HS256"])
        current_email = payload.get("email")
        
        if not current_email:
            raise HTTPException(status_code=401, detail="Token invalide")
        
        # Vérifier que l'utilisateur existe
        user = await db.users.find_one({"email": current_email.lower()})
        if not user:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
        
        # Vérifier le mot de passe
        if not verify_password(data.password, user["password"]):
            raise HTTPException(status_code=401, detail="Mot de passe incorrect")
        
        # Vérifier que le nouvel email n'est pas déjà utilisé
        existing = await db.users.find_one({"email": data.new_email.lower()})
        if existing:
            raise HTTPException(status_code=400, detail="Cet email est déjà utilisé par un autre compte")
        
        # Valider le format de l'email
        if not re.match(r'^[\w\.-]+@[\w\.-]+\.\w+$', data.new_email):
            raise HTTPException(status_code=400, detail="Format d'email invalide")
        
        # Mettre à jour l'email
        old_email = current_email.lower()
        result = await db.users.update_one(
            {"email": old_email},
            {
                "$set": {
                    "email": data.new_email.lower(),
                    "previous_email": old_email,
                    "email_changed_at": datetime.now(timezone.utc).isoformat()
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=500, detail="Erreur lors de la mise à jour")
        
        # Générer un nouveau token avec le nouvel email
        new_token = jwt.encode({
            "email": data.new_email.lower(),
            "user_id": user.get("id"),
            "exp": datetime.now(timezone.utc) + timedelta(hours=24)
        }, JWT_SECRET, algorithm="HS256")
        
        logging.info(f"Email changé: {old_email} -> {data.new_email.lower()}")
        
        return {
            "success": True,
            "message": "Email modifié avec succès",
            "new_email": data.new_email.lower(),
            "new_token": new_token
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur changement email utilisateur: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/admin/change-user-email")
async def admin_change_user_email(data: ChangeEmailRequest, admin = Depends(get_current_admin)):
    """Permet au super admin de changer l'email d'un utilisateur"""
    try:
        # Vérifier que l'utilisateur existe
        user = await db.users.find_one({"email": data.old_email.lower()})
        if not user:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
        
        # Vérifier que le nouvel email n'est pas déjà utilisé
        existing = await db.users.find_one({"email": data.new_email.lower()})
        if existing:
            raise HTTPException(status_code=400, detail="Cet email est déjà utilisé par un autre compte")
        
        # Valider le format
        if not re.match(r'^[\w\.-]+@[\w\.-]+\.\w+$', data.new_email):
            raise HTTPException(status_code=400, detail="Format d'email invalide")
        
        # Mettre à jour
        old_email = data.old_email.lower()
        result = await db.users.update_one(
            {"email": old_email},
            {
                "$set": {
                    "email": data.new_email.lower(),
                    "previous_email": old_email,
                    "email_changed_at": datetime.now(timezone.utc).isoformat(),
                    "email_changed_by": "super_admin"
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=500, detail="Erreur lors de la mise à jour")
        
        logging.info(f"Email changé par admin: {old_email} -> {data.new_email.lower()}")
        
        return {
            "success": True,
            "message": f"Email modifié: {old_email} → {data.new_email.lower()}",
            "old_email": old_email,
            "new_email": data.new_email.lower()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur changement email par admin: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/license-admin/change-user-email")
async def license_admin_change_user_email(data: ChangeEmailRequest, admin = Depends(get_current_license_admin)):
    """Permet à l'admin de licence de changer l'email d'un utilisateur de sa licence"""
    try:
        # Vérifier que l'utilisateur existe et appartient à la même licence
        user = await db.users.find_one({"email": data.old_email.lower()})
        if not user:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
        
        # Vérifier que l'utilisateur appartient à la licence de l'admin
        if user.get("license_key") != admin.get("license_key"):
            raise HTTPException(status_code=403, detail="Vous ne pouvez modifier que les utilisateurs de votre licence")
        
        # Vérifier que le nouvel email n'est pas déjà utilisé
        existing = await db.users.find_one({"email": data.new_email.lower()})
        if existing:
            raise HTTPException(status_code=400, detail="Cet email est déjà utilisé par un autre compte")
        
        # Valider le format
        if not re.match(r'^[\w\.-]+@[\w\.-]+\.\w+$', data.new_email):
            raise HTTPException(status_code=400, detail="Format d'email invalide")
        
        # Mettre à jour
        old_email = data.old_email.lower()
        result = await db.users.update_one(
            {"email": old_email},
            {
                "$set": {
                    "email": data.new_email.lower(),
                    "previous_email": old_email,
                    "email_changed_at": datetime.now(timezone.utc).isoformat(),
                    "email_changed_by": "license_admin"
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=500, detail="Erreur lors de la mise à jour")
        
        logging.info(f"Email changé par admin licence: {old_email} -> {data.new_email.lower()}")
        
        return {
            "success": True,
            "message": f"Email modifié: {old_email} → {data.new_email.lower()}",
            "old_email": old_email,
            "new_email": data.new_email.lower()
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur changement email par admin licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== ROUTES - CHANGEMENT MOT DE PASSE ====================

@api_router.post("/auth/change-password")
async def user_change_own_password(data: UserSelfChangePassword, request: Request):
    """Permet à un utilisateur de changer son propre mot de passe"""
    try:
        # Récupérer l'utilisateur actuel via le token
        auth_header = request.headers.get("Authorization")
        if not auth_header or not auth_header.startswith("Bearer "):
            raise HTTPException(status_code=401, detail="Non authentifié")
        
        token = auth_header.split(" ")[1]
        payload = jwt.decode(token, JWT_SECRET, algorithms=["HS256"])
        current_email = payload.get("email")
        
        if not current_email:
            raise HTTPException(status_code=401, detail="Token invalide")
        
        # Vérifier que l'utilisateur existe
        user = await db.users.find_one({"email": current_email.lower()})
        if not user:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
        
        # Vérifier le mot de passe actuel
        if not verify_password(data.current_password, user["password"]):
            raise HTTPException(status_code=401, detail="Mot de passe actuel incorrect")
        
        # Valider le nouveau mot de passe
        if len(data.new_password) < 6:
            raise HTTPException(status_code=400, detail="Le nouveau mot de passe doit contenir au moins 6 caractères")
        
        # Hasher le nouveau mot de passe
        hashed_password = hash_password(data.new_password)
        
        # Mettre à jour le mot de passe
        result = await db.users.update_one(
            {"email": current_email.lower()},
            {
                "$set": {
                    "password": hashed_password,
                    "password_changed_at": datetime.now(timezone.utc).isoformat()
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=500, detail="Erreur lors de la mise à jour")
        
        logging.info(f"Mot de passe changé pour: {current_email}")
        
        return {
            "success": True,
            "message": "Mot de passe modifié avec succès"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur changement mot de passe utilisateur: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/admin/change-user-password")
async def admin_change_user_password(data: AdminChangeUserPassword, admin = Depends(get_current_admin)):
    """Permet au super admin de changer le mot de passe d'un utilisateur"""
    try:
        # Vérifier que l'utilisateur existe
        user = await db.users.find_one({"email": data.user_email.lower()})
        if not user:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
        
        # Valider le nouveau mot de passe
        if len(data.new_password) < 6:
            raise HTTPException(status_code=400, detail="Le mot de passe doit contenir au moins 6 caractères")
        
        # Hasher le nouveau mot de passe
        hashed_password = hash_password(data.new_password)
        
        # Mettre à jour
        result = await db.users.update_one(
            {"email": data.user_email.lower()},
            {
                "$set": {
                    "password": hashed_password,
                    "password_changed_at": datetime.now(timezone.utc).isoformat(),
                    "password_changed_by": "admin"
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=500, detail="Erreur lors de la mise à jour")
        
        logging.info(f"Admin a changé le mot de passe de: {data.user_email}")
        
        return {
            "success": True,
            "message": f"Mot de passe modifié pour {data.user_email}"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur changement mot de passe par admin: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.post("/license-admin/change-user-password")
async def license_admin_change_user_password(data: AdminChangeUserPassword, admin = Depends(get_current_license_admin)):
    """Permet à l'admin de licence de changer le mot de passe d'un utilisateur de sa licence"""
    try:
        # Récupérer la licence de l'admin
        admin_user = await db.users.find_one({"email": admin["email"]})
        if not admin_user or not admin_user.get("is_license_admin"):
            raise HTTPException(status_code=403, detail="Non autorisé")
        
        admin_license_id = admin_user.get("license_id")
        
        # Vérifier que l'utilisateur existe et appartient à la même licence
        user = await db.users.find_one({"email": data.user_email.lower()})
        if not user:
            raise HTTPException(status_code=404, detail="Utilisateur non trouvé")
        
        if user.get("license_id") != admin_license_id:
            raise HTTPException(status_code=403, detail="Cet utilisateur n'appartient pas à votre licence")
        
        # Valider le nouveau mot de passe
        if len(data.new_password) < 6:
            raise HTTPException(status_code=400, detail="Le mot de passe doit contenir au moins 6 caractères")
        
        # Hasher le nouveau mot de passe
        hashed_password = hash_password(data.new_password)
        
        # Mettre à jour
        result = await db.users.update_one(
            {"email": data.user_email.lower()},
            {
                "$set": {
                    "password": hashed_password,
                    "password_changed_at": datetime.now(timezone.utc).isoformat(),
                    "password_changed_by": f"license_admin:{admin['email']}"
                }
            }
        )
        
        if result.modified_count == 0:
            raise HTTPException(status_code=500, detail="Erreur lors de la mise à jour")
        
        logging.info(f"Admin licence {admin['email']} a changé le mot de passe de: {data.user_email}")
        
        return {
            "success": True,
            "message": f"Mot de passe modifié pour {data.user_email}"
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logging.error(f"Erreur changement mot de passe par admin licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))

# ==================== ROUTES ADMIN - STATISTIQUES ====================

@api_router.get("/admin/stats")
async def get_admin_stats(admin = Depends(get_current_admin)):
    """Obtenir les statistiques globales (Admin seulement)"""
    try:
        now = datetime.now(timezone.utc)
        current_year = now.year
        current_month = now.month
        
        total_licenses = await db.licenses.count_documents({})
        active_licenses = await db.licenses.count_documents({"is_active": True})
        total_users = await db.users.count_documents({})
        active_users = await db.users.count_documents({"is_active": True})
        total_blocked_words = await db.blocked_words.count_documents({"is_active": True})
        
        # Licences expirées
        expired_licenses = await db.licenses.count_documents({
            "expiry_date": {"$lt": now.isoformat()}
        })
        
        # Total connexions (login_count de tous les utilisateurs)
        login_pipeline = [
            {"$group": {"_id": None, "total_logins": {"$sum": "$login_count"}}}
        ]
        login_result = await db.users.aggregate(login_pipeline).to_list(1)
        total_logins = login_result[0]["total_logins"] if login_result else 0
        
        # Connexions ce mois-ci
        logins_this_month = await db.login_logs.count_documents({
            "year": current_year,
            "month": current_month
        })
        
        # Connexions cette année
        logins_this_year = await db.login_logs.count_documents({
            "year": current_year
        })
        
        # Requêtes ce mois-ci
        requests_this_month = await db.request_logs.count_documents({
            "year": current_year,
            "month": current_month
        })
        
        # Requêtes cette année
        requests_this_year = await db.request_logs.count_documents({
            "year": current_year
        })
        
        # Total requêtes
        total_requests = await db.request_logs.count_documents({})
        
        # Requêtes par mois (12 derniers mois)
        requests_by_month = []
        for i in range(12):
            month = current_month - i
            year = current_year
            if month <= 0:
                month += 12
                year -= 1
            count = await db.request_logs.count_documents({"year": year, "month": month})
            requests_by_month.append({
                "year": year,
                "month": month,
                "count": count
            })
        
        # Top 10 utilisateurs par requêtes ce mois
        top_users_pipeline = [
            {"$match": {"year": current_year, "month": current_month, "user_email": {"$ne": None}}},
            {"$group": {"_id": "$user_email", "count": {"$sum": 1}}},
            {"$sort": {"count": -1}},
            {"$limit": 10}
        ]
        top_users_month = await db.request_logs.aggregate(top_users_pipeline).to_list(10)
        
        # Top 10 utilisateurs par requêtes cette année
        top_users_year_pipeline = [
            {"$match": {"year": current_year, "user_email": {"$ne": None}}},
            {"$group": {"_id": "$user_email", "count": {"$sum": 1}}},
            {"$sort": {"count": -1}},
            {"$limit": 10}
        ]
        top_users_year = await db.request_logs.aggregate(top_users_year_pipeline).to_list(10)
        
        # Statistiques par utilisateur (tous)
        users_stats_pipeline = [
            {"$match": {"user_email": {"$ne": None}}},
            {"$group": {
                "_id": "$user_email",
                "total_requests": {"$sum": 1},
                "requests_this_year": {
                    "$sum": {"$cond": [{"$eq": ["$year", current_year]}, 1, 0]}
                },
                "requests_this_month": {
                    "$sum": {"$cond": [
                        {"$and": [{"$eq": ["$year", current_year]}, {"$eq": ["$month", current_month]}]},
                        1, 0
                    ]}
                }
            }},
            {"$sort": {"total_requests": -1}}
        ]
        users_stats = await db.request_logs.aggregate(users_stats_pipeline).to_list(100)
        
        return {
            "success": True,
            "stats": {
                "licenses": {
                    "total": total_licenses,
                    "active": active_licenses,
                    "expired": expired_licenses
                },
                "users": {
                    "total": total_users,
                    "active": active_users
                },
                "blocked_words": total_blocked_words,
                "logins": {
                    "total": total_logins,
                    "this_month": logins_this_month,
                    "this_year": logins_this_year
                },
                "requests": {
                    "total": total_requests,
                    "this_month": requests_this_month,
                    "this_year": requests_this_year,
                    "by_month": requests_by_month
                },
                "top_users": {
                    "this_month": [{"email": u["_id"], "count": u["count"]} for u in top_users_month],
                    "this_year": [{"email": u["_id"], "count": u["count"]} for u in top_users_year]
                },
                "users_detailed": [
                    {
                        "email": u["_id"],
                        "total": u["total_requests"],
                        "this_year": u["requests_this_year"],
                        "this_month": u["requests_this_month"]
                    } for u in users_stats
                ]
            },
            "period": {
                "current_year": current_year,
                "current_month": current_month
            }
        }
        
    except Exception as e:
        logging.error(f"Erreur récupération stats: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@api_router.get("/license-admin/stats")
async def get_license_admin_stats(admin = Depends(get_current_license_admin)):
    """Obtenir les statistiques de la licence (Admin de licence seulement)"""
    try:
        license_key = admin.get("license_key")
        now = datetime.now(timezone.utc)
        current_year = now.year
        current_month = now.month
        
        # Nombre d'utilisateurs de la licence
        total_users = await db.users.count_documents({"license_key": license_key})
        active_users = await db.users.count_documents({"license_key": license_key, "is_active": True})
        
        # Total connexions des utilisateurs de la licence
        login_pipeline = [
            {"$match": {"license_key": license_key}},
            {"$group": {"_id": None, "total_logins": {"$sum": "$login_count"}}}
        ]
        login_result = await db.users.aggregate(login_pipeline).to_list(1)
        total_logins = login_result[0]["total_logins"] if login_result else 0
        
        # Requêtes ce mois-ci pour cette licence
        requests_this_month = await db.request_logs.count_documents({
            "license_key": license_key,
            "year": current_year,
            "month": current_month
        })
        
        # Requêtes cette année pour cette licence
        requests_this_year = await db.request_logs.count_documents({
            "license_key": license_key,
            "year": current_year
        })
        
        # Total requêtes pour cette licence
        total_requests = await db.request_logs.count_documents({
            "license_key": license_key
        })
        
        # Requêtes par mois (12 derniers mois)
        requests_by_month = []
        for i in range(12):
            month = current_month - i
            year = current_year
            if month <= 0:
                month += 12
                year -= 1
            count = await db.request_logs.count_documents({
                "license_key": license_key,
                "year": year,
                "month": month
            })
            requests_by_month.append({
                "year": year,
                "month": month,
                "count": count
            })
        
        # Top utilisateurs de la licence ce mois
        top_users_pipeline = [
            {"$match": {"license_key": license_key, "year": current_year, "month": current_month, "user_email": {"$ne": None}}},
            {"$group": {"_id": "$user_email", "count": {"$sum": 1}}},
            {"$sort": {"count": -1}},
            {"$limit": 10}
        ]
        top_users_month = await db.request_logs.aggregate(top_users_pipeline).to_list(10)
        
        # Statistiques détaillées par utilisateur de la licence
        users_stats_pipeline = [
            {"$match": {"license_key": license_key, "user_email": {"$ne": None}}},
            {"$group": {
                "_id": "$user_email",
                "total_requests": {"$sum": 1},
                "requests_this_year": {
                    "$sum": {"$cond": [{"$eq": ["$year", current_year]}, 1, 0]}
                },
                "requests_this_month": {
                    "$sum": {"$cond": [
                        {"$and": [{"$eq": ["$year", current_year]}, {"$eq": ["$month", current_month]}]},
                        1, 0
                    ]}
                }
            }},
            {"$sort": {"total_requests": -1}}
        ]
        users_stats = await db.request_logs.aggregate(users_stats_pipeline).to_list(100)
        
        return {
            "success": True,
            "license_key": license_key,
            "stats": {
                "users": {
                    "total": total_users,
                    "active": active_users
                },
                "logins": {
                    "total": total_logins
                },
                "requests": {
                    "total": total_requests,
                    "this_month": requests_this_month,
                    "this_year": requests_this_year,
                    "by_month": requests_by_month
                },
                "top_users_month": [{"email": u["_id"], "count": u["count"]} for u in top_users_month],
                "users_detailed": [
                    {
                        "email": u["_id"],
                        "total": u["total_requests"],
                        "this_year": u["requests_this_year"],
                        "this_month": u["requests_this_month"]
                    } for u in users_stats
                ]
            },
            "period": {
                "current_year": current_year,
                "current_month": current_month
            }
        }
        
    except Exception as e:
        logging.error(f"Erreur récupération stats licence: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# ==================== SECTION 25: NOTES DE VERSION (NOUVEAUTÉS) ====================

class ReleaseNoteCreate(BaseModel):
    title: str
    version: str = ""
    date: str = ""
    description: str = ""
    how_to_use: str = ""

@api_router.get("/release-notes")
async def get_release_notes():
    """Récupère toutes les notes de version, triées par date décroissante"""
    notes = await db.release_notes.find({}, {"_id": 0}).sort("date", -1).to_list(100)
    return notes

@api_router.post("/release-notes")
async def create_release_note(note: ReleaseNoteCreate, admin=Depends(get_current_admin)):
    """Crée une nouvelle note de version (admin seulement)"""
    note_data = note.dict()
    note_data["id"] = str(uuid.uuid4())
    note_data["created_at"] = datetime.now(timezone.utc).isoformat()
    if not note_data["date"]:
        note_data["date"] = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    await db.release_notes.insert_one(note_data)
    return {"success": True, "id": note_data["id"]}

@api_router.put("/release-notes/{note_id}")
async def update_release_note(note_id: str, note: ReleaseNoteCreate, admin=Depends(get_current_admin)):
    """Met à jour une note de version (admin seulement)"""
    update_data = note.dict()
    update_data["updated_at"] = datetime.now(timezone.utc).isoformat()
    result = await db.release_notes.update_one({"id": note_id}, {"$set": update_data})
    if result.modified_count == 0:
        raise HTTPException(status_code=404, detail="Note non trouvée")
    return {"success": True}

@api_router.delete("/release-notes/{note_id}")
async def delete_release_note(note_id: str, admin=Depends(get_current_admin)):
    """Supprime une note de version (admin seulement)"""
    result = await db.release_notes.delete_one({"id": note_id})
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="Note non trouvée")
    return {"success": True}

# Include the router in the main app
app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
)
logger = logging.getLogger(__name__)


@app.on_event("startup")
async def startup_event():
    """Charge les poids du detecteur au demarrage et initialise les données par défaut"""
    await load_detector_weights()
    await init_default_blocked_words()
    await init_release_notes()
    logging.info("Application demarree et poids du detecteur charges")

async def init_default_blocked_words():
    """Initialise des mots bloqués par défaut si la collection est vide"""
    try:
        count = await db.blocked_words.count_documents({})
        if count == 0:
            default_words = [
                # Violence
                {"word": "tuer quelqu'un", "category": "violence", "severity": "critical", "is_exception": False, "is_active": True},
                {"word": "fabriquer une arme", "category": "violence", "severity": "critical", "is_exception": False, "is_active": True},
                {"word": "blesser gravement", "category": "violence", "severity": "high", "is_exception": False, "is_active": True},
                # Drogues
                {"word": "fabriquer de la drogue", "category": "drugs", "severity": "critical", "is_exception": False, "is_active": True},
                {"word": "acheter des stupéfiants", "category": "drugs", "severity": "high", "is_exception": False, "is_active": True},
                {"word": "cultiver du cannabis", "category": "drugs", "severity": "high", "is_exception": False, "is_active": True},
                # Piratage
                {"word": "pirater un compte", "category": "hacking", "severity": "high", "is_exception": False, "is_active": True},
                {"word": "voler des mots de passe", "category": "hacking", "severity": "high", "is_exception": False, "is_active": True},
                {"word": "cracker un logiciel", "category": "hacking", "severity": "medium", "is_exception": False, "is_active": True},
                # Contenu inapproprié
                {"word": "contenu pour adulte", "category": "inappropriate", "severity": "medium", "is_exception": False, "is_active": True},
                {"word": "insulter un enseignant", "category": "inappropriate", "severity": "medium", "is_exception": False, "is_active": True},
                # Triche scolaire
                {"word": "tricher à l'examen", "category": "custom", "severity": "medium", "is_exception": False, "is_active": True},
                {"word": "acheter un travail", "category": "custom", "severity": "medium", "is_exception": False, "is_active": True},
                {"word": "plagier un texte", "category": "custom", "severity": "low", "is_exception": False, "is_active": True},
            ]
            
            for word_data in default_words:
                word_data["id"] = str(uuid.uuid4())
                word_data["created_at"] = datetime.now(timezone.utc).isoformat()
            
            await db.blocked_words.insert_many(default_words)
            logging.info(f"✅ {len(default_words)} mots bloqués par défaut initialisés")
    except Exception as e:
        logging.error(f"Erreur initialisation mots bloqués: {e}")

# Initialiser les notes par défaut au démarrage
async def init_release_notes():
    """Insère les notes initiales si la collection est vide"""
    count = await db.release_notes.count_documents({})
    if count == 0:
        default_notes = [
            {
                "id": str(uuid.uuid4()),
                "version": "1.0",
                "date": "2025-01-22",
                "title": "Lancement d'Étienne",
                "description": "- Assistant IA bilingue (Étienne en français, Steven en anglais) pour le personnel scolaire québécois\n- Système d'authentification avec clés de licence par organisation\n- Panneau d'administration multi-niveaux (Super Admin et Admin de licence)\n- 5 catégories de chat : Plans de cours, Évaluations, Activités, Ressources, Outils\n- Réponses adaptées au Programme de formation de l'école québécoise (PFEQ)",
                "how_to_use": "1. Connectez-vous avec votre courriel scolaire et la clé de licence fournie par votre organisation\n2. Choisissez une catégorie de conversation (Plans de cours, Évaluations, etc.)\n3. Posez votre question à Étienne dans la zone de texte\n4. Étienne répond en tenant compte du curriculum québécois",
                "created_at": "2025-01-22T00:00:00Z"
            },
            {
                "id": str(uuid.uuid4()),
                "version": "1.5",
                "date": "2025-01-30",
                "title": "Conversations cloud et LaTeX",
                "description": "- Conversations synchronisées dans le cloud (accès multi-appareils)\n- Historique de conversations avec recherche, épinglage et favoris\n- Export des conversations en fichier texte (.txt)\n- Support des formules mathématiques LaTeX dans les réponses\n- Réinitialisation du mot de passe par courriel\n- Export des réponses en PDF et Word avec logo Champagneur",
                "how_to_use": "1. Vos conversations sont automatiquement sauvegardées dans le cloud\n2. Utilisez la barre latérale gauche pour retrouver vos anciennes conversations\n3. Cliquez sur l'étoile pour mettre une conversation en favori\n4. Les formules mathématiques s'affichent automatiquement (ex: fractions, exposants)\n5. Utilisez les boutons PDF/Word sous chaque réponse pour télécharger",
                "created_at": "2025-01-30T00:00:00Z"
            },
            {
                "id": str(uuid.uuid4()),
                "version": "2.0",
                "date": "2025-02-16",
                "title": "Modération et améliorations IA",
                "description": "- Système de modération de contenu (mots bloqués) pour les administrateurs\n- Étienne ajoute des sources et références à la fin de ses réponses\n- Amélioration de la détection de langue (ne bascule plus en anglais par erreur)\n- Étienne se souvient qu'il s'adresse à des enseignants, pas à des élèves\n- Corrections de textes présentées en tableau (sans réécrire le texte entier)\n- Corrections de bugs : panneau Admin, impression LaTeX, page blanche licence",
                "how_to_use": "1. Les administrateurs peuvent gérer les mots bloqués via le panneau Admin > Mots bloqués\n2. Chaque réponse d'Étienne inclut maintenant une section Sources à la fin\n3. Étienne adapte son langage pour s'adresser à des professionnels de l'éducation",
                "created_at": "2025-02-16T00:00:00Z"
            },
            {
                "id": str(uuid.uuid4()),
                "version": "3.0",
                "date": "2025-12-01",
                "title": "Correction de textes MEQ",
                "description": "- Nouveau bouton « Corriger un texte » dédié avec formulaire complet\n- Protocole officiel du MEQ : Étienne pose 5 questions obligatoires avant de corriger\n- Barèmes officiels C4 et C5 intégrés (CS Laval / MEQ) pour tous les niveaux (sec 1 à 5)\n- Tableaux complets de 17 tranches de mots (101 à 501+), non modifiables\n- Pondérations officielles pré-remplies selon le niveau scolaire\n- Règles ministérielles de comptage d'erreurs (syntaxe, ponctuation, orthographe)\n- Attribution automatique des cotes (A à E) et calcul de la note finale\n- Les fautes répétées du même type sont soulignées mais non comptées\n- Upload de fichiers : PDF, Word, texte ou image (OCR automatique)\n- Téléchargement de la correction en PDF ou Word (correction seule ou avec texte original)\n- Carte dédiée « Corriger un texte » dans les fonctionnalités principales",
                "how_to_use": "1. Cliquez sur le bouton rouge « Corriger un texte » en haut à droite du chat\n2. Sélectionnez le niveau scolaire (Secondaire 1 à 5) — les pondérations et barèmes se remplissent automatiquement\n3. Ajustez les critères et pondérations si nécessaire\n4. Entrez le nombre total de points (ex: /40)\n5. Optionnel : précisez les descripteurs du critère 1\n6. Entrez le texte de l'élève : copier-coller OU joindre un fichier (PDF, Word, image)\n7. Cliquez « Corriger le texte »\n8. Étienne analyse le texte, compte les erreurs, attribue les cotes et calcule la note\n9. Téléchargez la correction en PDF ou Word avec le bouton sous la réponse",
                "created_at": "2025-12-01T00:00:00Z"
            }
        ]
        await db.release_notes.insert_many(default_notes)
        logging.info(f"✅ {len(default_notes)} notes de version initiales créées")

@app.on_event("shutdown")
async def shutdown_db_client():
    client.close()
